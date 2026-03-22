import telebot
from telebot import types
import requests
import sqlite3
from datetime import datetime, date, timedelta
import base64
import os
import tempfile
import time
import re
import html
import json
import traceback
import sys
from openai import OpenAI
import threading
import io
import uuid
import random
import httpx
import urllib.parse

# ========== ЛОГИРОВАНИЕ ==========

def log_error(error_msg):
    with open("bot_errors.log", "a", encoding="utf-8") as f:
        f.write(f"{datetime.now().isoformat()} - {error_msg}\n")
    print(error_msg)

# ========== ДОПОЛНИТЕЛЬНЫЕ БИБЛИОТЕКИ ==========

try:
    from PyPDF2 import PdfReader
except:
    PdfReader = None
try:
    from docx import Document
except:
    Document = None
try:
    import openpyxl
except:
    openpyxl = None
try:
    import xlrd
except:
    xlrd = None
try:
    from pptx import Presentation
except:
    Presentation = None
try:
    from striprtf.striprtf import rtf_to_text
except:
    rtf_to_text = None
try:
    from pydub import AudioSegment
except:
    AudioSegment = None

# ========== ОСНОВНЫЕ НАСТРОЙКИ ==========

TOKEN = "8482873514:AAG45GvYRhpU-pfxWGLXpF1opsGLSFeiJQI"
ADMIN_IDS = [6197133464]
BOT_NAME = "ken4kk_ai_bot"
CREATOR_TAG = "@KeN4kk_n1"
BOT_VERSION = "10.16-fixed"
MAX_HISTORY = 10
BANK_CARD = "2200701310027153"
USD_TO_RUB = 77.27
STAR_TO_RUB = 1.5

CRYPTO_TOKEN = "541485:AA7Bc67c60qd6lUiv8aXXaeXBT1kH4cyjzU"
CRYPTO_API_URL = "https://pay.crypt.bot/api"

# ========== ЦЕНЫ ПО УМОЛЧАНИЮ (в копейках) ==========

DEFAULT_PRICES = {
    'generate': 4000,
    'edit': 5000,
    'text_default': 20,
    'voice': 500,
    'tts': 250,
    'video': 20000,
    'search': 100,
    'video_6s': 10000,
    'video_10s': 15000,
}
PRICES = DEFAULT_PRICES.copy()

# ========== ЦЕНЫ AITUNNEL (В КОПЕЙКАХ ЗА 1 ТОКЕН) ==========

AITUNNEL_PRICES = {
    'gemini-2.5-flash-lite': {'input': 0.00192, 'output': 0.00768},
    'deepseek-r1': {'input': 0.01056, 'output': 0.042048},
    'claude-sonnet-4.6': {'input': 0.0576, 'output': 0.288},
    'deepseek-v3.2-speciale': {'input': 0.005376, 'output': 0.008064},
    'gpt-5-nano': {'input': 0.00096, 'output': 0.00768},
    'gemini-3.1-pro-preview-customtools': {'input': 0.0384, 'output': 0.2304},
    'sonar': {'input': 0.0192, 'output': 0.0192},
}
PROFIT_MULTIPLIER = 2.0

# ========== БАЗОВЫЕ ЦЕНЫ ДЛЯ МЕДИА-УСЛУГ (В РУБЛЯХ) ==========

MEDIA_BASE_PRICES = {
    'generate': 9.75,
    'video_per_second': 10.0,
    'tts_per_1000_chars': 10.0,
}

# ========== ПРЕМИУМ-ЭМОДЗИ ==========

PREMIUM_EMOJI = {
    'generate': '5257974976094412956',
    'edit': '5257974976094412956',
    'video': '5258077307985207053',
    'search': '5429571366384842791',
    'recognize': '5258020476977946656',
    'new_chat': '5258215846450305872',
    'info': '5226513232549664618',
    'contacts': '5258337316715373336',
    'voice': '5258020476977946656',
    'model': '5258093637450866522',
    'tts': '5258330865674494479',
    'help': '5260535596941582167',
    'topup': '5258204546391351475',
    'settings': '5253959125838090076',
    'main_menu': '5226513232549664618',
    'canceled': '5260342697075416641',
    'success': '5260726538302660868',
    'warning': '5258474669769497337',
    'sparkles': '5258165702707125574',
    'profile': '5258011929993026890',
    'purchases': '5199457120428249992',
    'balance': '5258204546391351475',
    'id': '5258476306152038031',
    'ima': '5316727448644103237',
    'popolnenie': '5258368777350816286',
    'razmer': '5258254475386167466',
    'successrazmer': '5257974976094412956',
    'viborvideo': '5258077307985207053',
    'stars': '5258165702707125574',
    'crypto': '5258368777350816286',
    'help_emoji': '5260535596941582167',
    'tts_menu': '5260325873688518261',
    'choose_model': '5258093637450866522',
    'news': '5258336354642697821',
    'tts_voice_emoji': '5260325873688518261',
    'music': '5258289810082111221',
    'tts_input': '5258020476977946656',
    'search_input': '5429571366384842791',
    'documents': '5258477770735885832',
    'code': '5258260149037965799',
    'math': '5258334872878980409',
    'photo_edit': '5258450450448915742',
    'audio': '5260652149469094137',
    'welcome_star': '5258501105293205250',
    'time_reset': '5258258882022612173',
    'subscribe': '5258328383183396223',
}

# ========== НАСТРОЙКИ ДЛИННЫХ ОТВЕТОВ ==========

LONG_RESPONSE_API = "https://ken4kk-app.ru/"
MAX_MESSAGE_LENGTH = 3500
long_responses = {}

# ========== НАСТРОЙКИ AITUNNEL ==========

AITUNNEL_URL = "https://api.aitunnel.ru/v1"
AITUNNEL_KEY = "sk-aitunnel-gx40555wueQ5kxLzHuPWOn8UV3sueXnX"
TRANSCRIBE_MODEL = "gpt-4o-mini-transcribe"
SEARCH_MODEL = "sonar"
aitunnel_client = OpenAI(api_key=AITUNNEL_KEY, base_url=AITUNNEL_URL)

# ========== НАСТРОЙКИ GEN-API ==========

GENAPI_KEY = "sk-gI2vUyeXjYzgg1FpA8PFPYfqh58M1P67d7Kw037t7HRxhULr3NFgkfF9PJx8"
GENAPI_BASE_URL = "https://api.gen-api.ru/api/v1"

# ========== НАСТРОЙКИ SHAZAM API ==========

SHAZAM_API_URL = "https://shazam-api.com/api"
SHAZAM_API_KEY = "TfnWxUlYTreEoGPSeqMiD3cgPHC7mvws9WxVdOm2ZB8jPOxYFT2frDIdvFnPNOrF"

# ========== ПАРАМЕТРЫ ВИДЕО ==========

VIDEO_DURATIONS = [6, 10]

# ========== МОДЕЛИ (будут загружены из БД) ==========

AVAILABLE_MODELS = {
    "✨ Gemini 2.5 Flash Lite": "gemini-2.5-flash-lite",
    "🧠 DeepSeek R1": "deepseek-r1",
    "🌲 Claude Sonnet 4.6": "claude-sonnet-4.6",
    "🤖 GPT-5": "gpt-5",
    "🧠 DeepSeek V3.2 Speciale": "deepseek-v3.2-speciale",
    "🤖 GPT-5 Nano": "gpt-5-nano",
    "🤖 Gemini 3.1 Pro Preview Customtools": "gemini-3.1-pro-preview-customtools",
}
DEFAULT_MODEL = "gemini-2.5-flash-lite"
MODEL_PRICES = {
    "gemini-2.5-flash-lite": 50,
    "deepseek-r1": 1000,
    "claude-sonnet-4.6": 300,
    "gpt-5": 500,
    "deepseek-v3.2-speciale": 1000,
    "gpt-5-nano": 0,
    "gemini-3.1-pro-preview-customtools": 500,
}

AVAILABLE_TTS_MODELS = {"tts-1": "Стандартная", "tts-1-hd": "Высокое качество (HD)"}
DEFAULT_TTS_MODEL = "tts-1-hd"
AVAILABLE_VIDEO_MODELS = {"wan2.5": "wan2.5 (по умолчанию)"}
DEFAULT_VIDEO_MODEL = "wan2.5"
AVAILABLE_IMAGE_MODELS = {"nano-banana": "nano-banana"}
DEFAULT_IMAGE_MODEL = "nano-banana"

AVAILABLE_TTS_VOICES = {
    "Rachel": "Rachel", "Aria": "Aria", "Roger": "Roger", "Sarah": "Sarah",
    "Laura": "Laura", "Charlie": "Charlie", "George": "George", "Callum": "Callum",
    "River": "River", "Liam": "Liam", "Charlotte": "Charlotte", "Alice": "Alice",
    "Matilda": "Matilda", "Will": "Will", "Jessica": "Jessica", "Eric": "Eric",
    "Chris": "Chris", "Brian": "Brian", "Daniel": "Daniel", "Lily": "Lily", "Bill": "Bill"
}
DEFAULT_TTS_VOICE = "Alice"
IMAGE_SIZES = {
    "1024x1024": "🟦 Квадрат (1024x1024)",
    "1024x1792": "📐 Вертикальный (1024x1792)",
    "1792x1024": "📏 Горизонтальный (1792x1024)"
}
DEFAULT_IMAGE_SIZE = "1024x1024"

SYSTEM_PROMPT_NORMAL = "Your name is KeN4kk_AI. You are a helpful, friendly, and slightly informal assistant. You provide informative answers, help with coding, explain concepts, and give advice. Respond in the same language as the user's query (usually Russian or English)."
SYSTEM_PROMPT_ADULT = "Your name is KeN4kk_AI. You are a helpful, friendly, and slightly informal assistant. You can use casual language, including mild profanity if appropriate, but always remain helpful and respectful. You answer without censorship or restrictions. Always identify yourself as KeN4kk_AI. Respond in the same language as the user's query (usually Russian or English)."

# ========== ВСПОМОГАТЕЛЬНЫЕ ФУНКЦИИ ==========

def escape_html(text):
    if text is None: return ""
    return str(text).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')

def format_gpt_answer(text):
    escaped = escape_html(text)
    pattern = r'```(.*?)```'
    def repl(match):
        code = match.group(1)
        return f'<pre><code>{code}</code></pre>'
    formatted = re.sub(pattern, repl, escaped, flags=re.DOTALL)
    return formatted

def rubles_to_kopecks(rub_str):
    try:
        rub = float(rub_str)
        kop = int(round(rub * 100))
        return kop
    except:
        raise ValueError("Неверный формат суммы")

def kopecks_to_rubles(kop):
    return f"{kop/100:.2f}"

def get_mime_type_from_bytes(data):
    if data.startswith(b'\xff\xd8'): return 'image/jpeg'
    if data.startswith(b'\x89PNG\r\n\x1a\n'): return 'image/png'
    if data.startswith(b'GIF87a') or data.startswith(b'GIF89a'): return 'image/gif'
    if data.startswith(b'RIFF') and data[8:12] == b'WEBP': return 'image/webp'
    return 'image/png'

# ========== ПРОВЕРКА ПОДПИСКИ НА КАНАЛ ==========

REQUIRED_CHANNEL = "@ken4kk_news"

def is_subscribed(user_id):
    try:
        member = bot.get_chat_member(REQUIRED_CHANNEL, user_id)
        status = member.status
        result = status in ['member', 'administrator', 'creator']
        log_error(f"is_subscribed({user_id}): статус {status}, результат {result}")
        return result
    except Exception as e:
        log_error(f"is_subscribed({user_id}): ОШИБКА {e}")
        return False

def subscription_required(func):
    def wrapper(message, *args, **kwargs):
        if message.chat.type == 'private':
            user_id = message.from_user.id
            if is_admin(user_id):
                return func(message, *args, **kwargs)
            if not is_subscribed(user_id):
                text = (f'<tg-emoji emoji-id="{PREMIUM_EMOJI["subscribe"]}">⚠️</tg-emoji> '
                        f'Чтобы пользоваться ботом, необходимо подписаться на '
                        f'<a href="https://t.me/ken4kk_news">KeN4kk_News</a>. '
                        f'Мы сделали так, чтобы не допускать ботов и чтобы вы могли получать '
                        f'сообщения об обновлениях и новости через канал/бота.')
                markup = types.InlineKeyboardMarkup().add(
                    types.InlineKeyboardButton("📢 Подписаться", url="https://t.me/ken4kk_news")
                )
                try:
                    bot.send_message(message.chat.id, text, parse_mode="HTML", reply_markup=markup)
                    log_error(f"Отправлено сообщение о подписке пользователю {user_id} (из декоратора)")
                except Exception as e:
                    log_error(f"Ошибка при отправке сообщения о подписке пользователю {user_id} в декораторе: {e}")
                return
        return func(message, *args, **kwargs)
    return wrapper

# ========== БАЗА ДАННЫХ ==========

db_lock = threading.RLock()
conn = sqlite3.connect('bot_data.db', check_same_thread=False)
cursor = conn.cursor()

cursor.executescript('''
CREATE TABLE IF NOT EXISTS users (
    user_id INTEGER PRIMARY KEY,
    username TEXT,
    first_name TEXT,
    balance INTEGER DEFAULT 0,
    blocked INTEGER DEFAULT 0,
    reg_date TEXT,
    text_model TEXT,
    voice_enabled INTEGER DEFAULT 0,
    tts_voice TEXT DEFAULT 'nova',
    in_chat INTEGER DEFAULT 0,
    last_admin_message TEXT DEFAULT '',
    language TEXT DEFAULT 'ru',
    theme TEXT DEFAULT 'default',
    mode TEXT DEFAULT 'normal',
    is_admin INTEGER DEFAULT 0,
    internet_search_enabled INTEGER DEFAULT 0
);
CREATE TABLE IF NOT EXISTS daily_usage (
    user_id INTEGER,
    usage_date TEXT,
    count INTEGER DEFAULT 0,
    PRIMARY KEY (user_id, usage_date)
);
CREATE TABLE IF NOT EXISTS tts_daily_usage (
    user_id INTEGER,
    usage_date TEXT,
    count INTEGER DEFAULT 0,
    PRIMARY KEY (user_id, usage_date)
);
CREATE TABLE IF NOT EXISTS conversations (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    chat_id INTEGER,
    role TEXT,
    content TEXT,
    timestamp TEXT
);
CREATE TABLE IF NOT EXISTS queries_log (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    user_id INTEGER,
    query TEXT,
    answer TEXT,
    timestamp TEXT
);
CREATE TABLE IF NOT EXISTS settings (
    key TEXT PRIMARY KEY,
    value TEXT
);
CREATE TABLE IF NOT EXISTS prices (
    service TEXT PRIMARY KEY,
    price INTEGER
);
CREATE TABLE IF NOT EXISTS transactions (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    user_id INTEGER,
    amount INTEGER,
    type TEXT,
    description TEXT,
    timestamp TEXT
);
CREATE TABLE IF NOT EXISTS text_models (
    model_id TEXT PRIMARY KEY,
    display_name TEXT,
    price INTEGER
);
CREATE TABLE IF NOT EXISTS tts_models (
    model_id TEXT PRIMARY KEY,
    display_name TEXT
);
CREATE TABLE IF NOT EXISTS video_models (
    model_id TEXT PRIMARY KEY,
    display_name TEXT
);
CREATE TABLE IF NOT EXISTS image_models (
    model_id TEXT PRIMARY KEY,
    display_name TEXT
);
CREATE TABLE IF NOT EXISTS crypto_invoices (
    invoice_id TEXT PRIMARY KEY,
    user_id INTEGER,
    amount_rub INTEGER,
    status TEXT DEFAULT 'active',
    created_at TEXT,
    pay_url TEXT
);
CREATE TABLE IF NOT EXISTS card_payments (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    user_id INTEGER,
    amount_rub INTEGER,
    first_name TEXT,
    last_name TEXT,
    comment TEXT,
    status TEXT DEFAULT 'pending',
    created_at TEXT
);
CREATE TABLE IF NOT EXISTS referrals (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    referrer_id INTEGER,
    referred_id INTEGER UNIQUE,
    username TEXT,
    first_name TEXT,
    joined_date TEXT,
    status TEXT DEFAULT 'pending',
    gift_sent INTEGER DEFAULT 0
);
''')
conn.commit()

# ===== ПРОВЕРКА И ДОБАВЛЕНИЕ КОЛОНОК =====

def add_column_if_not_exists(table, column, definition):
    cursor.execute(f"PRAGMA table_info({table})")
    existing = [col[1] for col in cursor.fetchall()]
    if column not in existing:
        try:
            cursor.execute(f"ALTER TABLE {table} ADD COLUMN {column} {definition}")
            print(f"Added column {column} to {table}")
        except Exception as e:
            print(f"Error adding column {column} to {table}: {e}")

required_users = {
    'is_admin': 'INTEGER DEFAULT 0',
    'mode': 'TEXT DEFAULT "normal"',
    'theme': 'TEXT DEFAULT "default"',
    'language': 'TEXT DEFAULT "ru"',
    'last_admin_message': 'TEXT DEFAULT ""',
    'in_chat': 'INTEGER DEFAULT 0',
    'tts_voice': 'TEXT DEFAULT "nova"',
    'voice_enabled': 'INTEGER DEFAULT 0',
    'text_model': 'TEXT',
    'reg_date': 'TEXT',
    'blocked': 'INTEGER DEFAULT 0',
    'balance': 'INTEGER DEFAULT 0',
    'first_name': 'TEXT',
    'username': 'TEXT',
    'internet_search_enabled': 'INTEGER DEFAULT 0'
}
for col, definition in required_users.items():
    add_column_if_not_exists('users', col, definition)

add_column_if_not_exists('crypto_invoices', 'amount_rub', 'INTEGER')
add_column_if_not_exists('crypto_invoices', 'pay_url', 'TEXT')
add_column_if_not_exists('crypto_invoices', 'status', 'TEXT DEFAULT "active"')
add_column_if_not_exists('crypto_invoices', 'created_at', 'TEXT')
add_column_if_not_exists('crypto_invoices', 'user_id', 'INTEGER')

add_column_if_not_exists('card_payments', 'amount_rub', 'INTEGER')
add_column_if_not_exists('card_payments', 'first_name', 'TEXT')
add_column_if_not_exists('card_payments', 'last_name', 'TEXT')
add_column_if_not_exists('card_payments', 'comment', 'TEXT')
add_column_if_not_exists('card_payments', 'status', 'TEXT DEFAULT "pending"')
add_column_if_not_exists('card_payments', 'created_at', 'TEXT')

cursor.execute("UPDATE users SET balance=0 WHERE balance IS NULL")
conn.commit()

# ========== ФУНКЦИИ ДЛЯ РАБОТЫ С БЕСПЛАТНЫМИ ЗАПРОСАМИ ==========

def get_daily_usage(user_id):
    today = date.today().isoformat()
    with db_lock:
        row = cursor.execute("SELECT count FROM daily_usage WHERE user_id=? AND usage_date=?", (user_id, today)).fetchone()
        return row[0] if row else 0

def increment_daily_usage(user_id):
    today = date.today().isoformat()
    with db_lock:
        cursor.execute("INSERT INTO daily_usage (user_id, usage_date, count) VALUES (?, ?, 1) ON CONFLICT(user_id, usage_date) DO UPDATE SET count = count + 1", (user_id, today))
        conn.commit()

def get_free_quota_left(user_id):
    used = get_daily_usage(user_id)
    return max(0, 20 - used)

def is_free_quota_available(user_id):
    return get_free_quota_left(user_id) > 0

def use_free_quota(user_id):
    if is_free_quota_available(user_id):
        increment_daily_usage(user_id)
        return True
    return False

def get_daily_tts_usage(user_id):
    today = date.today().isoformat()
    with db_lock:
        row = cursor.execute("SELECT count FROM tts_daily_usage WHERE user_id=? AND usage_date=?", (user_id, today)).fetchone()
        return row[0] if row else 0

def increment_daily_tts_usage(user_id):
    today = date.today().isoformat()
    with db_lock:
        cursor.execute('''
            INSERT INTO tts_daily_usage (user_id, usage_date, count) VALUES (?, ?, 1)
            ON CONFLICT(user_id, usage_date) DO UPDATE SET count = count + 1
        ''', (user_id, today))
        conn.commit()

def get_free_tts_quota_left(user_id):
    used = get_daily_tts_usage(user_id)
    return max(0, 5 - used)

def is_free_tts_available(user_id):
    return get_free_tts_quota_left(user_id) > 0

def use_free_tts_quota(user_id):
    if is_free_tts_available(user_id):
        increment_daily_tts_usage(user_id)
        return True
    return False

# ========== ФУНКЦИИ ДЛЯ РЕЖИМА ПОИСКА ==========

def get_internet_search_enabled(user_id):
    with db_lock:
        r = cursor.execute("SELECT internet_search_enabled FROM users WHERE user_id=?", (user_id,)).fetchone()
        return r[0] if r else 0

def set_internet_search_enabled(user_id, enabled):
    with db_lock:
        cursor.execute("UPDATE users SET internet_search_enabled=? WHERE user_id=?", (int(enabled), user_id))
        conn.commit()

# ========== ФУНКЦИИ ДЛЯ УПРАВЛЕНИЯ ЦЕНАМИ МОДЕЛЕЙ ==========

def save_aitunnel_price(model_id, price_type, value):
    with db_lock:
        cursor.execute(
            "INSERT OR REPLACE INTO settings (key, value) VALUES (?, ?)",
            (f"aitunnel_price_{model_id}_{price_type}", str(value))
        )
        conn.commit()

def load_aitunnel_prices():
    global AITUNNEL_PRICES
    rows = cursor.execute("SELECT key, value FROM settings WHERE key LIKE 'aitunnel_price_%'").fetchall()
    if rows:
        for key, value in rows:
            parts = key.replace('aitunnel_price_', '').split('_')
            if len(parts) == 2:
                model_id, price_type = parts
                if model_id not in AITUNNEL_PRICES:
                    AITUNNEL_PRICES[model_id] = {}
                AITUNNEL_PRICES[model_id][price_type] = float(value)
    else:
        for model_id, prices in AITUNNEL_PRICES.items():
            for price_type, value in prices.items():
                save_aitunnel_price(model_id, price_type, value)

def get_profit_multiplier():
    return float(get_setting('profit_multiplier', PROFIT_MULTIPLIER))

def set_profit_multiplier(value):
    global PROFIT_MULTIPLIER
    set_setting('profit_multiplier', str(value))
    PROFIT_MULTIPLIER = float(value)

def get_media_base_price(service):
    return float(get_setting(f'media_base_price_{service}', MEDIA_BASE_PRICES.get(service, 0)))

def set_media_base_price(service, value):
    set_setting(f'media_base_price_{service}', str(value))
    if 'MEDIA_BASE_PRICES' in globals():
        MEDIA_BASE_PRICES[service] = float(value)

def get_setting(key, default=None):
    with db_lock:
        r = cursor.execute("SELECT value FROM settings WHERE key=?", (key,)).fetchone()
        return r[0] if r else default

def set_setting(key, value):
    with db_lock:
        cursor.execute("INSERT OR REPLACE INTO settings (key, value) VALUES (?, ?)", (key, value))
        conn.commit()

# ========== ЗАГРУЗКА НАСТРОЕК ==========

def load_all_settings():
    global PRICES, AVAILABLE_MODELS, MODEL_PRICES, DEFAULT_TTS_MODEL, DEFAULT_VIDEO_MODEL, DEFAULT_IMAGE_MODEL, ADMIN_IDS, USD_TO_RUB, STAR_TO_RUB, AITUNNEL_PRICES, PROFIT_MULTIPLIER, MEDIA_BASE_PRICES

    rows = cursor.execute("SELECT service, price FROM prices").fetchall()
    if rows:
        PRICES.clear()
        for s, p in rows:
            PRICES[s] = p
    else:
        for s, p in DEFAULT_PRICES.items():
            cursor.execute("INSERT OR REPLACE INTO prices (service, price) VALUES (?,?)", (s, p))
            PRICES[s] = p
        conn.commit()

    for key, default_price in DEFAULT_PRICES.items():
        if key not in PRICES:
            cursor.execute("INSERT OR REPLACE INTO prices (service, price) VALUES (?,?)", (key, default_price))
            PRICES[key] = default_price
    conn.commit()

    code_models = AVAILABLE_MODELS.copy()
    code_prices = MODEL_PRICES.copy()

    rows = cursor.execute("SELECT model_id, display_name, price FROM text_models").fetchall()
    if rows:
        AVAILABLE_MODELS.clear()
        MODEL_PRICES.clear()
        for mid, name, price in rows:
            AVAILABLE_MODELS[name] = mid
            MODEL_PRICES[mid] = price
        for name, mid in code_models.items():
            if mid not in AVAILABLE_MODELS.values():
                price = code_prices.get(mid, 0)
                cursor.execute("INSERT OR REPLACE INTO text_models (model_id, display_name, price) VALUES (?,?,?)", (mid, name, price))
                AVAILABLE_MODELS[name] = mid
                MODEL_PRICES[mid] = price
        conn.commit()
    else:
        for name, mid in code_models.items():
            price = code_prices.get(mid, 0)
            cursor.execute("INSERT OR REPLACE INTO text_models (model_id, display_name, price) VALUES (?,?,?)", (mid, name, price))
        conn.commit()
        AVAILABLE_MODELS = code_models
        MODEL_PRICES = code_prices

    rows = cursor.execute("SELECT model_id, display_name FROM tts_models").fetchall()
    if rows:
        AVAILABLE_TTS_MODELS.clear()
        for mid, name in rows: AVAILABLE_TTS_MODELS[mid] = name
    else:
        for mid, name in AVAILABLE_TTS_MODELS.items():
            cursor.execute("INSERT OR REPLACE INTO tts_models (model_id, display_name) VALUES (?,?)", (mid, name))

    rows = cursor.execute("SELECT model_id, display_name FROM video_models").fetchall()
    if rows:
        AVAILABLE_VIDEO_MODELS.clear()
        for mid, name in rows: AVAILABLE_VIDEO_MODELS[mid] = name
    else:
        for mid, name in AVAILABLE_VIDEO_MODELS.items():
            cursor.execute("INSERT OR REPLACE INTO video_models (model_id, display_name) VALUES (?,?)", (mid, name))

    rows = cursor.execute("SELECT model_id, display_name FROM image_models").fetchall()
    if rows:
        AVAILABLE_IMAGE_MODELS.clear()
        for mid, name in rows: AVAILABLE_IMAGE_MODELS[mid] = name
    else:
        for mid, name in AVAILABLE_IMAGE_MODELS.items():
            cursor.execute("INSERT OR REPLACE INTO video_models (model_id, display_name) VALUES (?,?)", (mid, name))

    settings = dict(cursor.execute("SELECT key, value FROM settings").fetchall())
    if 'tts_model' in settings: DEFAULT_TTS_MODEL = settings['tts_model']
    else: cursor.execute("INSERT OR REPLACE INTO settings (key, value) VALUES (?,?)", ('tts_model', DEFAULT_TTS_MODEL))
    if 'video_model' in settings: DEFAULT_VIDEO_MODEL = settings['video_model']
    else: cursor.execute("INSERT OR REPLACE INTO settings (key, value) VALUES (?,?)", ('video_model', DEFAULT_VIDEO_MODEL))
    if 'image_model' in settings: DEFAULT_IMAGE_MODEL = settings['image_model']
    else: cursor.execute("INSERT OR REPLACE INTO settings (key, value) VALUES (?,?)", ('image_model', DEFAULT_IMAGE_MODEL))
    if 'usd_rate' in settings:
        USD_TO_RUB = float(settings['usd_rate'])
    else:
        cursor.execute("INSERT OR REPLACE INTO settings (key, value) VALUES (?,?)", ('usd_rate', str(USD_TO_RUB)))
    if 'star_rate' in settings:
        STAR_TO_RUB = float(settings['star_rate'])
    else:
        cursor.execute("INSERT OR REPLACE INTO settings (key, value) VALUES (?,?)", ('star_rate', str(STAR_TO_RUB)))

    load_aitunnel_prices()
    PROFIT_MULTIPLIER = get_profit_multiplier()
    for service in MEDIA_BASE_PRICES.keys():
        MEDIA_BASE_PRICES[service] = get_media_base_price(service)

    required_gift_settings = {
        'gift_id': '5170145012310081615'
    }
    for key, default_value in required_gift_settings.items():
        if get_setting(key) is None:
            set_setting(key, default_value)
            print(f"✅ Установлена настройка {key} = {default_value}")

    conn.commit()

    admins = [r[0] for r in cursor.execute("SELECT user_id FROM users WHERE is_admin=1").fetchall()]
    for a in ADMIN_IDS:
        if a not in admins:
            cursor.execute("UPDATE users SET is_admin=1 WHERE user_id=?", (a,))
    conn.commit()
    ADMIN_IDS = [r[0] for r in cursor.execute("SELECT user_id FROM users WHERE is_admin=1").fetchall()]

load_all_settings()

def get_usd_rate():
    return USD_TO_RUB

def set_usd_rate(rate):
    global USD_TO_RUB
    with db_lock:
        cursor.execute("INSERT OR REPLACE INTO settings (key, value) VALUES (?, ?)", ('usd_rate', str(rate)))
        conn.commit()
        USD_TO_RUB = rate

def get_star_rate():
    return STAR_TO_RUB

def set_star_rate(rate):
    global STAR_TO_RUB
    with db_lock:
        cursor.execute("INSERT OR REPLACE INTO settings (key, value) VALUES (?, ?)", ('star_rate', str(rate)))
        conn.commit()
        STAR_TO_RUB = rate

# ========== ФУНКЦИИ ПОЛЬЗОВАТЕЛЕЙ ==========

def register_or_update_user(message):
    uid = message.from_user.id
    username = message.from_user.username or ""
    first_name = message.from_user.first_name or ""
    with db_lock:
        cursor.execute('''
            INSERT INTO users (user_id, username, first_name, reg_date, balance)
            VALUES (?, ?, ?, ?, 0)
            ON CONFLICT(user_id) DO UPDATE SET
                username = excluded.username,
                first_name = excluded.first_name
        ''', (uid, username, first_name, datetime.now().isoformat()))
        conn.commit()

def get_balance(uid):
    with db_lock:
        r = cursor.execute("SELECT balance FROM users WHERE user_id=?", (uid,)).fetchone()
        return r[0] if r else 0

def update_balance(uid, amount, reason=""):
    log_error(f"update_balance: uid={uid}, amount={amount}, reason={reason}")
    with db_lock:
        cursor.execute("UPDATE users SET balance = COALESCE(balance, 0) + ? WHERE user_id=?", (amount, uid))
        conn.commit()
        if amount != 0:
            type_ = 'replenish' if amount > 0 else 'spend'
            cursor.execute(
                "INSERT INTO transactions (user_id, amount, type, description, timestamp) VALUES (?,?,?,?,?)",
                (uid, amount, type_, reason, datetime.now().isoformat())
            )
            conn.commit()

def is_admin(user_id):
    return user_id in ADMIN_IDS

def get_text_model(uid):
    with db_lock:
        r = cursor.execute("SELECT text_model FROM users WHERE user_id=?", (uid,)).fetchone()
        model = r[0] if r and r[0] else DEFAULT_MODEL
        if model not in AVAILABLE_MODELS.values():
            model = DEFAULT_MODEL
            cursor.execute("UPDATE users SET text_model=? WHERE user_id=?", (model, uid))
            conn.commit()
        return model

def set_text_model(uid, model):
    with db_lock:
        cursor.execute("UPDATE users SET text_model=? WHERE user_id=?", (model, uid))
        conn.commit()

def get_voice_enabled(uid):
    with db_lock:
        r = cursor.execute("SELECT voice_enabled FROM users WHERE user_id=?", (uid,)).fetchone()
        return r[0] if r else 0

def set_voice_enabled(uid, enabled):
    with db_lock:
        cursor.execute("UPDATE users SET voice_enabled=? WHERE user_id=?", (int(enabled), uid))
        conn.commit()

def get_tts_voice(uid):
    with db_lock:
        r = cursor.execute("SELECT tts_voice FROM users WHERE user_id=?", (uid,)).fetchone()
        return r[0] if r and r[0] else DEFAULT_TTS_VOICE

def set_tts_voice(uid, voice):
    with db_lock:
        cursor.execute("UPDATE users SET tts_voice=? WHERE user_id=?", (voice, uid))
        conn.commit()

def is_blocked(uid):
    with db_lock:
        r = cursor.execute("SELECT blocked FROM users WHERE user_id=?", (uid,)).fetchone()
        return r and r[0] == 1

def block_user(uid):
    with db_lock:
        cursor.execute("UPDATE users SET blocked=1 WHERE user_id=?", (uid,))
        conn.commit()

def unblock_user(uid):
    with db_lock:
        cursor.execute("UPDATE users SET blocked=0 WHERE user_id=?", (uid,))
        conn.commit()

def set_in_chat(uid, value):
    with db_lock:
        cursor.execute("UPDATE users SET in_chat=? WHERE user_id=?", (int(value), uid))
        conn.commit()

def get_in_chat(uid):
    with db_lock:
        r = cursor.execute("SELECT in_chat FROM users WHERE user_id=?", (uid,)).fetchone()
        return r[0] if r else 0

def set_last_admin_message(uid, message):
    with db_lock:
        cursor.execute("UPDATE users SET last_admin_message=? WHERE user_id=?", (message, uid))
        conn.commit()

def get_last_admin_message(uid):
    with db_lock:
        r = cursor.execute("SELECT last_admin_message FROM users WHERE user_id=?", (uid,)).fetchone()
        return r[0] if r else ""

def set_language(uid, lang):
    with db_lock:
        cursor.execute("UPDATE users SET language=? WHERE user_id=?", (lang, uid))
        conn.commit()

def get_language(uid):
    with db_lock:
        r = cursor.execute("SELECT language FROM users WHERE user_id=?", (uid,)).fetchone()
        return r[0] if r else 'ru'

def set_mode(uid, mode):
    with db_lock:
        cursor.execute("UPDATE users SET mode=? WHERE user_id=?", (mode, uid))
        conn.commit()

def get_mode(uid):
    with db_lock:
        r = cursor.execute("SELECT mode FROM users WHERE user_id=?", (uid,)).fetchone()
        return r[0] if r else 'normal'

def log_query(uid, query, answer):
    with db_lock:
        cursor.execute("INSERT INTO queries_log (user_id, query, answer, timestamp) VALUES (?,?,?,?)",
                       (uid, query, answer[:500], datetime.now().isoformat()))
        conn.commit()

def add_to_history(chat_id, role, content):
    with db_lock:
        cursor.execute("INSERT INTO conversations (chat_id, role, content, timestamp) VALUES (?,?,?,?)",
                       (chat_id, role, content, datetime.now().isoformat()))
        conn.commit()

def get_history(chat_id):
    with db_lock:
        rows = cursor.execute("SELECT role, content FROM conversations WHERE chat_id=? ORDER BY timestamp DESC LIMIT ?",
                              (chat_id, MAX_HISTORY)).fetchall()
        return list(reversed(rows))

def clear_history(chat_id):
    with db_lock:
        cursor.execute("DELETE FROM conversations WHERE chat_id=?", (chat_id,))
        conn.commit()

def get_all_transactions(limit=100):
    with db_lock:
        return cursor.execute(
            "SELECT user_id, amount, type, description, timestamp FROM transactions ORDER BY timestamp DESC LIMIT ?",
            (limit,)
        ).fetchall()

# ========== ФУНКЦИИ ДЛЯ РАБОТЫ С МОДЕЛЯМИ ==========

def get_tts_model():
    return get_setting('tts_model', DEFAULT_TTS_MODEL)

def set_tts_model(model_id):
    set_setting('tts_model', model_id)
    global DEFAULT_TTS_MODEL
    DEFAULT_TTS_MODEL = model_id

def get_video_model():
    return get_setting('video_model', DEFAULT_VIDEO_MODEL)

def set_video_model(model_id):
    set_setting('video_model', model_id)
    global DEFAULT_VIDEO_MODEL
    DEFAULT_VIDEO_MODEL = model_id

def get_image_model():
    return get_setting('image_model', DEFAULT_IMAGE_MODEL)

def set_image_model(model_id):
    set_setting('image_model', model_id)
    global DEFAULT_IMAGE_MODEL
    DEFAULT_IMAGE_MODEL = model_id

def add_text_model(display_name, model_id, price_kopecks):
    with db_lock:
        cursor.execute("INSERT OR REPLACE INTO text_models (model_id, display_name, price) VALUES (?, ?, ?)",
                       (model_id, display_name, price_kopecks))
        conn.commit()
    AVAILABLE_MODELS[display_name] = model_id
    MODEL_PRICES[model_id] = price_kopecks

def add_tts_model(model_id, display_name):
    with db_lock:
        cursor.execute("INSERT OR REPLACE INTO tts_models (model_id, display_name) VALUES (?, ?)",
                       (model_id, display_name))
        conn.commit()
    AVAILABLE_TTS_MODELS[model_id] = display_name

def add_video_model(model_id, display_name):
    with db_lock:
        cursor.execute("INSERT OR REPLACE INTO video_models (model_id, display_name) VALUES (?, ?)",
                       (model_id, display_name))
        conn.commit()
    AVAILABLE_VIDEO_MODELS[model_id] = display_name

def add_image_model(model_id, display_name):
    with db_lock:
        cursor.execute("INSERT OR REPLACE INTO image_models (model_id, display_name) VALUES (?, ?)",
                       (model_id, display_name))
        conn.commit()
    AVAILABLE_IMAGE_MODELS[model_id] = display_name

# ========== РАБОТА С ФАЙЛАМИ ==========

def download_file(file_id, max_retries=3):
    for attempt in range(max_retries):
        try:
            file_info = bot.get_file(file_id)
            url = f'https://api.telegram.org/file/bot{TOKEN}/{file_info.file_path}'
            resp = requests.get(url, timeout=30)
            resp.raise_for_status()
            return resp.content
        except Exception as e:
            log_error(f"Download file error (attempt {attempt+1}/{max_retries}): {e}")
            if attempt == max_retries - 1:
                return None
            time.sleep(2)
    return None

def download_voice(file_id):
    return download_file(file_id)

def extract_text_from_pdf(pdf_bytes):
    if PdfReader is None:
        return "[Библиотека PyPDF2 не установлена]"
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
            tmp.write(pdf_bytes); tmp_path = tmp.name
        reader = PdfReader(tmp_path)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text[:2000]
    except Exception as e:
        log_error(f"PDF extraction error: {e}")
        return f"[Ошибка чтения PDF: {e}]"
    finally:
        if tmp_path and os.path.exists(tmp_path): os.unlink(tmp_path)

def extract_text_from_docx(docx_bytes):
    if Document is None:
        return "[Библиотека python-docx не установлена]"
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp:
            tmp.write(docx_bytes); tmp_path = tmp.name
        doc = Document(tmp_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text[:2000]
    except Exception as e:
        log_error(f"DOCX extraction error: {e}")
        return f"[Ошибка чтения DOCX: {e}]"
    finally:
        if tmp_path and os.path.exists(tmp_path): os.unlink(tmp_path)

def extract_text_from_txt(txt_bytes):
    try:
        return txt_bytes.decode('utf-8')[:2000]
    except:
        return "[Не удалось прочитать текстовый файл]"

def extract_text_from_xlsx(xlsx_bytes):
    if openpyxl is None:
        return "[Библиотека openpyxl не установлена]"
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
            tmp.write(xlsx_bytes); tmp_path = tmp.name
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        text = ""
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    text += row_text + "\n"
        return text[:2000]
    except Exception as e:
        log_error(f"XLSX extraction error: {e}")
        return f"[Ошибка чтения XLSX: {e}]"
    finally:
        if tmp_path and os.path.exists(tmp_path): os.unlink(tmp_path)

def extract_text_from_xls(xls_bytes):
    if xlrd is None:
        return "[Библиотека xlrd не установлена]"
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xls') as tmp:
            tmp.write(xls_bytes); tmp_path = tmp.name
        wb = xlrd.open_workbook(tmp_path)
        text = ""
        for sheet in wb.sheets():
            for row in range(sheet.nrows):
                row_values = sheet.row_values(row)
                row_text = " ".join([str(cell) for cell in row_values if cell != ""])
                if row_text.strip():
                    text += row_text + "\n"
        return text[:2000]
    except Exception as e:
        log_error(f"XLS extraction error: {e}")
        return f"[Ошибка чтения XLS: {e}]"
    finally:
        if tmp_path and os.path.exists(tmp_path): os.unlink(tmp_path)

def extract_text_from_pptx(pptx_bytes):
    if Presentation is None:
        return "[Библиотека python-pptx не установлена]"
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            tmp.write(pptx_bytes); tmp_path = tmp.name
        prs = Presentation(tmp_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text[:2000]
    except Exception as e:
        log_error(f"PPTX extraction error: {e}")
        return f"[Ошибка чтения PPTX: {e}]"
    finally:
        if tmp_path and os.path.exists(tmp_path): os.unlink(tmp_path)

def extract_text_from_rtf(rtf_bytes):
    if rtf_to_text is None:
        return "[Библиотека striprtf не установлена]"
    try:
        rtf_str = rtf_bytes.decode('utf-8', errors='ignore')
        text = rtf_to_text(rtf_str)
        return text[:2000]
    except Exception as e:
        log_error(f"RTF extraction error: {e}")
        return f"[Ошибка чтения RTF: {e}]"

def split_message_simple(text, max_length=4096):
    if len(text) <= max_length:
        return [text]
    return [text[i:i+max_length] for i in range(0, len(text), max_length)]

def send_long_message(chat_id, text, reply_to_message_id=None, force_file_threshold=5000):
    """
    Отправляет длинное сообщение.
    Если длина текста превышает force_file_threshold (по умолчанию 9000),
    отправляет как файл. Иначе разбивает на части по 4096.
    """
    if len(text) > force_file_threshold:
        send_as_file(chat_id, text, reply_to_message_id=reply_to_message_id)
        return

    parts = split_message_simple(text, 4096)
    for idx, part in enumerate(parts):
        try:
            if idx == 0 and reply_to_message_id:
                bot.send_message(chat_id, part, reply_to_message_id=reply_to_message_id, parse_mode="HTML")
            else:
                bot.send_message(chat_id, part, parse_mode="HTML")
            time.sleep(0.3)
        except Exception as e:
            log_error(f"Send long message error: {e}")

def send_as_file(chat_id, text, filename="response.txt", caption=None, reply_to_message_id=None):
    """Отправляет длинный текст как текстовый файл."""
    with tempfile.NamedTemporaryFile(delete=False, suffix='.txt', mode='w', encoding='utf-8') as tmp:
        tmp.write(text)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'rb') as f:
            bot.send_document(
                chat_id, 
                f, 
                visible_file_name=filename,
                caption=caption,
                reply_to_message_id=reply_to_message_id,
                parse_mode="HTML"
            )
    except Exception as e:
        log_error(f"Error sending file: {e}")
        # Если не удалось отправить файл, пробуем разбить на части
        send_long_message(chat_id, text, reply_to_message_id)
    finally:
        os.unlink(tmp_path)


# ========== ФУНКЦИИ ДЛЯ ДЛИННЫХ ОТВЕТОВ ==========

def publish_to_my_site(text, user_id):
    """Публикует текст на вашем сайте, возвращает полный URL или None при ошибке."""
    try:
        response = requests.post(
            f"{LONG_RESPONSE_API}/api/create",
            json={'user_id': user_id, 'text': text},
            timeout=10
        )
        if response.status_code == 200:
            data = response.json()
            relative_url = data.get('url')
            if relative_url:
                full_url = LONG_RESPONSE_API + relative_url
                log_error(f"Полный URL: {full_url}")
                return full_url
        else:
            log_error(f"Ошибка публикации на сайте: {response.status_code} - {response.text}")
    except Exception as e:
        log_error(f"Исключение при публикации на сайте: {e}")
    return None

def handle_long_response(chat_id, text, uid, reply_to_message_id=None):
    """Показывает кнопки с выбором способа получения длинного ответа."""
    site_url = publish_to_my_site(text, uid)
    if site_url:
        response_id = site_url.split('/view/')[-1]
        markup = types.InlineKeyboardMarkup(row_width=2)
        buttons = [
            types.InlineKeyboardButton("📖 Открыть в Telegram", url=site_url),
            types.InlineKeyboardButton("🌐 Открыть в браузере", url=site_url),
            types.InlineKeyboardButton("📨 Прислать в чат", callback_data=f"long_text:{response_id}"),
            types.InlineKeyboardButton("📎 Прислать файлом", callback_data=f"long_file:{response_id}")
        ]
        markup.add(*buttons)
        bot.send_message(
            chat_id,
            "⚠️ Ответ слишком длинный. Выберите, как его получить:",
            reply_markup=markup,
            reply_to_message_id=reply_to_message_id
        )
    else:
        send_long_message(chat_id, text, reply_to_message_id=reply_to_message_id)

# ========== AITUNNEL ТЕКСТ ==========

def ask_aitunnel(chat_id, user_text, image_bytes=None, video_bytes=None, force_search=False):
    uid = chat_id if isinstance(chat_id, int) else chat_id
    model = get_text_model(uid)
    history = get_history(chat_id)
    mode = get_mode(uid)
    system_prompt = SYSTEM_PROMPT_ADULT if mode == 'adult' else SYSTEM_PROMPT_NORMAL

    search_enabled = force_search or (get_internet_search_enabled(uid) == 1)
    actual_model = SEARCH_MODEL if search_enabled else model

    model_prices = AITUNNEL_PRICES.get(actual_model)
    if not model_prices:
        log_error(f"Цена для модели {actual_model} не найдена в AITUNNEL_PRICES. Списание не произведено.")
        return get_text(uid, 'ai_error', error=f"Цена для модели {actual_model} не настроена администратором.")

    free_used = False
    if is_free_quota_available(uid) and actual_model in ['gemini-2.5-flash-lite', 'gpt-5-nano']:
        use_free_quota(uid)
        free_used = True
        max_possible_cost_kopecks = 0
    else:
        estimated_input = 5000
        estimated_output = 10000
        max_possible_cost_kopecks = (estimated_input * model_prices['input'] + estimated_output * model_prices['output']) * get_profit_multiplier()
        if get_balance(uid) < max_possible_cost_kopecks:
            return get_text(uid, 'insufficient_funds',
                            balance=get_balance(uid)/100,
                            price=max_possible_cost_kopecks/100)

    messages = [{"role": "system", "content": system_prompt}]
    for role, content in history:
        messages.append({"role": role, "content": content})

    user_content = []
    if user_text:
        user_content.append({"type": "text", "text": user_text})
    if image_bytes:
        b64 = base64.b64encode(image_bytes).decode('utf-8')
        user_content.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}})
    if video_bytes:
        b64 = base64.b64encode(video_bytes).decode('utf-8')
        user_content.append({"type": "video_url", "video_url": {"url": f"data:video/mp4;base64,{b64}"}})

    if len(user_content) == 1 and user_content[0]["type"] == "text":
        messages.append({"role": "user", "content": user_text})
    else:
        messages.append({"role": "user", "content": user_content})

    try:
        response = aitunnel_client.chat.completions.create(
            model=actual_model,
            messages=messages,
            temperature=0.7,
            max_tokens=32000,
            timeout=180
        )

        answer = response.choices[0].message.content
        prompt_tokens = response.usage.prompt_tokens
        completion_tokens = response.usage.completion_tokens

        if not free_used:
            cost_input = prompt_tokens * model_prices['input']
            cost_output = completion_tokens * model_prices['output']
            total_cost = (cost_input + cost_output) * get_profit_multiplier()
            total_cost_kopecks = int(total_cost) + (1 if total_cost % 1 > 0 else 0)

            if get_balance(uid) < total_cost_kopecks:
                return get_text(uid, 'insufficient_funds',
                                balance=get_balance(uid)/100,
                                price=total_cost_kopecks/100)

            update_balance(uid, -total_cost_kopecks,
                           reason=f"Токены {prompt_tokens}+{completion_tokens} для {actual_model}")
            log_error(f"Списано {total_cost_kopecks/100:.2f}₽ за {prompt_tokens}+{completion_tokens} токенов (модель {actual_model})")
        else:
            log_error(f"Бесплатный запрос для {uid}: использовано {prompt_tokens}+{completion_tokens} токенов.")

        add_to_history(chat_id, "user", f"[Медиа] {user_text}" if (image_bytes or video_bytes) else user_text)
        add_to_history(chat_id, "assistant", answer)
        return format_gpt_answer(answer)

    except Exception as e:
        log_error(f"AITUNNEL chat error: {e}")
        return get_text(uid, 'ai_error', error=str(e))

# ========== GEN-API ВИДЕО ==========

def generate_video(prompt, duration, image_bytes=None):
    url = "https://api.gen-api.ru/api/v1/networks/ltx-2"
    headers = {"Authorization": f"Bearer {GENAPI_KEY}"}

    data = {
        "prompt": prompt,
        "mode": "fast",
        "duration": str(duration),
        "resolution": "1080p",
        "aspect_ratio": "16:9",
        "fps": "25",
    }

    files = {}
    if image_bytes:
        files["image_url"] = ("image.jpg", image_bytes, "image/jpeg")

    max_retries = 3
    retry_delay = 5

    for attempt in range(max_retries):
        try:
            log_error(f"LTX request attempt {attempt+1}/{max_retries}")
            log_error(f"LTX request data: {data}")
            log_error(f"LTX files present: {bool(files)}")

            response = requests.post(url, headers=headers, data=data, files=files, timeout=45)
            if response.status_code != 200:
                log_error(f"LTX response body: {response.text}")
            response.raise_for_status()

            result = response.json()
            request_id = result.get("request_id")
            if not request_id:
                log_error("No request_id in response")
                return None

            log_error(f"Video generation started, request_id: {request_id}")

            for status_attempt in range(60):
                time.sleep(5)
                status_url = f"https://api.gen-api.ru/api/v1/request/get/{request_id}"
                status_resp = requests.get(status_url, headers={"Authorization": f"Bearer {GENAPI_KEY}"}, timeout=30)
                if status_resp.status_code != 200:
                    log_error(f"Status check failed: {status_resp.status_code} - {status_resp.text}")
                status_resp.raise_for_status()
                status_data = status_resp.json()

                if status_data.get("status") == "success":
                    log_error("Video generation completed")
                    video_url = None
                    if status_data.get("result") and isinstance(status_data["result"], list) and len(status_data["result"]) > 0:
                        video_url = status_data["result"][0]
                    elif status_data.get("output"):
                        video_url = status_data["output"]
                    elif status_data.get("video_url"):
                        video_url = status_data["video_url"]
                    elif status_data.get("url"):
                        video_url = status_data["url"]

                    if not video_url:
                        log_error("No video URL found in success response")
                        return None

                    log_error(f"Downloading video from: {video_url}")
                    video_resp = requests.get(video_url, timeout=60)
                    video_resp.raise_for_status()
                    return video_resp.content

                elif status_data.get("status") == "failed":
                    log_error("Video generation failed")
                    return None
            else:
                log_error("Video generation timeout")
                return None

        except requests.exceptions.ConnectTimeout as e:
            log_error(f"Connection timeout (attempt {attempt+1}): {e}")
            if attempt == max_retries - 1:
                log_error("Max retries exceeded, giving up")
                return None
            time.sleep(retry_delay)
        except requests.exceptions.ConnectionError as e:
            log_error(f"Connection error (attempt {attempt+1}): {e}")
            if attempt == max_retries - 1:
                log_error("Max retries exceeded, giving up")
                return None
            time.sleep(retry_delay)
        except Exception as e:
            log_error(f"Video generation error: {traceback.format_exc()}")
            return None

    return None

# ========== GEN-API ИЗОБРАЖЕНИЯ ==========

def create_generation_task(prompt, num_images=1, output_format='png', aspect_ratio='default', translate_input=True):
    headers = {'Authorization': f'Bearer {GENAPI_KEY}', 'Content-Type': 'application/json', 'Accept': 'application/json'}
    data = {'prompt': prompt, 'translate_input': translate_input, 'num_images': num_images, 'output_format': output_format, 'aspect_ratio': aspect_ratio}
    try:
        r = requests.post(f"{GENAPI_BASE_URL}/networks/nano-banana", json=data, headers=headers, timeout=30)
        r.raise_for_status()
        return r.json().get('request_id')
    except Exception as e:
        log_error(f"Create generation task error: {e}")
        return None

def create_edit_task_multiple(prompt, images_bytes_list, num_images=1, output_format='png', aspect_ratio='default'):
    headers = {'Authorization': f'Bearer {GENAPI_KEY}', 'Accept': 'application/json'}
    data = {'prompt': prompt, 'num_images': str(num_images), 'output_format': output_format, 'aspect_ratio': aspect_ratio}
    files = []
    for idx, img_bytes in enumerate(images_bytes_list):
        mime = get_mime_type_from_bytes(img_bytes)
        ext = mime.split('/')[1]
        f = io.BytesIO(img_bytes)
        f.name = f'image_{idx}.{ext}'
        files.append(('image_urls[]', (f.name, f, mime)))
    try:
        r = requests.post(f"{GENAPI_BASE_URL}/networks/nano-banana", data=data, files=files, headers=headers, timeout=30)
        r.raise_for_status()
        return r.json().get('request_id')
    except Exception as e:
        log_error(f"Create edit task multiple error: {e}")
        return None
    finally:
        for _, (_, f, _) in files:
            f.close()

def get_task_result(request_id, max_attempts=60, delay=5):
    url = f"{GENAPI_BASE_URL}/request/get/{request_id}"
    headers = {'Authorization': f'Bearer {GENAPI_KEY}', 'Accept': 'application/json'}
    for _ in range(max_attempts):
        try:
            r = requests.get(url, headers=headers, timeout=30)
            r.raise_for_status()
            data = r.json()
            if data.get('status') == 'success':
                img_url = data.get('result', [None])[0] or data.get('full_response',[{}])[0].get('url')
                if img_url:
                    img = requests.get(img_url, timeout=60)
                    if img.status_code == 200: return img.content
            elif data.get('status') == 'failed':
                return None
            time.sleep(delay)
        except Exception as e:
            log_error(f"Get task result error: {e}")
            time.sleep(delay)
    return None

def generate_aitunnel_image(prompt, size):
    aspect = {'1024x1024':'1:1','1024x1792':'9:16','1792x1024':'16:9'}.get(size, 'default')
    rid = create_generation_task(prompt=prompt, num_images=1, output_format='png', aspect_ratio=aspect)
    return get_task_result(rid) if rid else None

def edit_aitunnel_image(images_bytes_list, prompt, size):
    aspect = {'1024x1024':'1:1','1024x1792':'9:16','1792x1024':'16:9'}.get(size, 'default')
    rid = create_edit_task_multiple(prompt=prompt, images_bytes_list=images_bytes_list, num_images=1, output_format='png', aspect_ratio=aspect)
    return get_task_result(rid) if rid else None

# ========== GEN-API TTS ==========

def create_tts_task(text, voice, stability=1, similarity_boost=1, style=1):
    url = f"{GENAPI_BASE_URL}/networks/tts-turbo"
    headers = {
        'Authorization': f'Bearer {GENAPI_KEY}',
        'Content-Type': 'application/json',
        'Accept': 'application/json'
    }
    payload = {
        "text": text,
        "voice": voice,
        "stability": stability,
        "similarity_boost": similarity_boost,
        "style": style,
        "is_sync": False
    }
    try:
        r = requests.post(url, json=payload, headers=headers, timeout=30)
        r.raise_for_status()
        data = r.json()
        return data.get('request_id')
    except Exception as e:
        log_error(f"Create TTS task error: {e}")
        return None

def get_tts_task_result(request_id, max_attempts=60, delay=5):
    url = f"{GENAPI_BASE_URL}/request/get/{request_id}"
    headers = {'Authorization': f'Bearer {GENAPI_KEY}', 'Accept': 'application/json'}
    for attempt in range(max_attempts):
        try:
            r = requests.get(url, headers=headers, timeout=30)
            r.raise_for_status()
            data = r.json()
            if data.get('status') == 'success':
                output = data.get('output')
                if output:
                    if isinstance(output, list) and len(output) > 0:
                        output = output[0]
                    if isinstance(output, dict) and 'url' in output:
                        output = output['url']
                    if isinstance(output, str):
                        audio_resp = requests.get(output, timeout=60)
                        if audio_resp.status_code == 200:
                            return audio_resp.content
                if data.get('result'):
                    result = data['result']
                    if isinstance(result, list) and len(result) > 0:
                        result = result[0]
                    if isinstance(result, str):
                        audio_resp = requests.get(result, timeout=60)
                        if audio_resp.status_code == 200:
                            return audio_resp.content
                if data.get('response_type') == 'base64' and data.get('output'):
                    return base64.b64decode(data['output'])
            elif data.get('status') == 'failed':
                log_error(f"TTS task failed: {data}")
                return None
        except Exception as e:
            log_error(f"Get TTS task result error: {e}")
        time.sleep(delay)
    log_error(f"TTS task {request_id} timeout")
    return None

def generate_tts_genapi(text, voice):
    rid = create_tts_task(text, voice)
    if not rid:
        return None
    return get_tts_task_result(rid)

# ========== ОТПРАВКА РЕЗУЛЬТАТОВ ==========

def generate_and_send(message, prompt, size=DEFAULT_IMAGE_SIZE):
    uid = message.from_user.id
    chat_id = message.chat.id
    base_price = get_media_base_price('generate')
    price = int(base_price * get_profit_multiplier() * 100)
    balance = get_balance(uid)
    if balance < price:
        bot.send_message(chat_id, get_text(uid, 'insufficient_funds', balance=balance/100, price=price/100), parse_mode="HTML")
        return
    thinking = bot.send_message(chat_id, get_text(uid, 'generating_image'), parse_mode="HTML")
    try:
        result = generate_aitunnel_image(prompt, size)
        if not result:
            bot.edit_message_text(get_text(uid, 'image_failed'), chat_id, thinking.message_id, parse_mode="HTML")
            return
        caption = f"🎨 <b>Запрос:</b> <code>{html.escape(prompt)}</code>"
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
            tmp.write(result); tmp_path = tmp.name
        with open(tmp_path, 'rb') as f:
            bot.send_photo(chat_id, f, caption=caption, timeout=600, parse_mode="HTML")
        os.unlink(tmp_path)
        update_balance(uid, -price, reason=f"Генерация изображения: {prompt[:50]}")
        bot.delete_message(chat_id, thinking.message_id)
    except Exception as e:
        log_error(f"Generate and send error: {e}")
        bot.edit_message_text(f"❌ Ошибка: {e}", chat_id, thinking.message_id, parse_mode="HTML")

def process_edit_step_multiple(chat_id, uid, prompt, size, images_bytes_list):
    base_price = get_media_base_price('generate')
    price = int(base_price * get_profit_multiplier() * 100)
    if get_balance(uid) < price:
        bot.send_message(chat_id, get_text(uid, 'insufficient_funds', balance=get_balance(uid)/100, price=price/100), parse_mode="HTML")
        return
    thinking = bot.send_message(chat_id, get_text(uid, 'editing_photo'), parse_mode="HTML")
    try:
        result = edit_aitunnel_image(images_bytes_list, prompt, size)
        if not result:
            bot.edit_message_text(get_text(uid, 'edit_failed'), chat_id, thinking.message_id, parse_mode="HTML")
            return
        caption = f'<tg-emoji emoji-id="{PREMIUM_EMOJI["edit"]}">✏️</tg-emoji> {escape_html(prompt)}\n📐 {escape_html(size)}\n📸 Использовано фото: {len(images_bytes_list)}'
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
            tmp.write(result); tmp_path = tmp.name
        with open(tmp_path, 'rb') as f:
            bot.send_photo(chat_id, f, caption=caption, timeout=600, parse_mode="HTML")
        os.unlink(tmp_path)
        update_balance(uid, -price, reason=f"Редактирование фото ({len(images_bytes_list)} шт): {prompt[:50]}")
        bot.delete_message(chat_id, thinking.message_id)
    except Exception as e:
        log_error(f"Edit error: {e}")
        bot.edit_message_text(f"❌ Ошибка: {e}", chat_id, thinking.message_id, parse_mode="HTML")

def process_edit_step(message, prompt, size=DEFAULT_IMAGE_SIZE):
    uid = message.from_user.id
    chat_id = message.chat.id
    if not message.photo:
        bot.send_message(chat_id, get_text(uid, 'edit_no_photo'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
        return
    photo = message.photo[-1]
    image_bytes = download_file(photo.file_id)
    if not image_bytes:
        bot.send_message(chat_id, get_text(uid, 'failed_load_photo'), parse_mode="HTML")
        return
    process_edit_step_multiple(chat_id, uid, prompt, size, [image_bytes])

def process_tts_step(message):
    uid = message.from_user.id
    chat_id = message.chat.id
    text = message.text.strip()
    if text.lower() in ["отмена", "cancel"]:
        bot.send_message(chat_id, get_text(uid, 'canceled'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
        return
    if not text:
        bot.send_message(chat_id, get_text(uid, 'empty_prompt'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
        return

    free_used = False
    if is_free_tts_available(uid):
        use_free_tts_quota(uid)
        price = 0
        free_used = True
    else:
        base_tts_price = get_media_base_price('tts_per_1000_chars')
        char_count = len(text)
        price = int((char_count / 1000) * base_tts_price * get_profit_multiplier() * 100)
        price = max(price, 1)
        if get_balance(uid) < price:
            bot.send_message(chat_id, get_text(uid, 'insufficient_funds', balance=get_balance(uid)/100, price=price/100), parse_mode="HTML")
            return

    if len(text) > 3000:
        text = text[:3000] + "..."
        bot.send_message(chat_id, get_text(uid, 'text_too_long'), parse_mode="HTML")

    thinking = bot.send_message(chat_id, get_text(uid, 'generating_tts'), parse_mode="HTML")
    voice = get_tts_voice(uid)

    try:
        audio_bytes = generate_tts_genapi(text, voice)
        if not audio_bytes:
            bot.edit_message_text(get_text(uid, 'tts_failed'), chat_id, thinking.message_id, parse_mode="HTML")
            return

        with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as tmp:
            tmp.write(audio_bytes); tmp_path = tmp.name
        caption = f'<tg-emoji emoji-id="{PREMIUM_EMOJI["tts"]}">🔊</tg-emoji> Текст: {escape_html(text[:200])}...' if len(text) > 200 else f'Текст: {escape_html(text)}'
        with open(tmp_path, 'rb') as f:
            bot.send_audio(chat_id, f, caption=caption, title="Синтезированная речь", performer="KeN4kk_AI", timeout=600, parse_mode="HTML")
        os.unlink(tmp_path)
        if not free_used:
            update_balance(uid, -price, reason=f"TTS: {text[:50]}")
        bot.delete_message(chat_id, thinking.message_id)
    except Exception as e:
        log_error(f"TTS process error: {e}")
        bot.edit_message_text(f"❌ Ошибка: {e}", chat_id, thinking.message_id, parse_mode="HTML")

# ========== ИНТЕРНЕТ ПОИСК ==========

def search_internet(query, uid):
    price = PRICES.get('search', 0)
    if get_balance(uid) < price:
        return get_text(uid, 'insufficient_funds', balance=get_balance(uid)/100, price=price/100)
    try:
        messages = [{"role": "system", "content": "You are a helpful assistant that searches the internet and provides accurate information."},
                    {"role": "user", "content": query}]
        response = aitunnel_client.chat.completions.create(model=SEARCH_MODEL, messages=messages, temperature=0.7, max_tokens=4000)
        answer = response.choices[0].message.content
        update_balance(uid, -price, reason=f"Поиск: {query[:50]}")
        return format_gpt_answer(answer)
    except Exception as e:
        log_error(f"Search error: {e}")
        return get_text(uid, 'search_failed')

# ========== SHAZAM ==========

def recognize_music_by_file(file_bytes, file_ext):
    url = f"{SHAZAM_API_URL}/recognize"
    headers = {"Authorization": f"Bearer {SHAZAM_API_KEY}"}
    files = {"file": (f"audio{file_ext}", file_bytes)}
    try:
        r = requests.post(url, headers=headers, files=files, timeout=30)
        if r.status_code == 200:
            return r.json().get("uuid")
        elif r.status_code == 400:
            return "INVALID_URL"
        else:
            log_error(f"Shazam error {r.status_code}: {r.text}")
            return None
    except Exception as e:
        log_error(f"Shazam exception: {e}")
        return None

def recognize_music_by_url(audio_url):
    url = f"{SHAZAM_API_URL}/recognize"
    headers = {"Authorization": f"Bearer {SHAZAM_API_KEY}"}
    data = {"url": audio_url}
    try:
        r = requests.post(url, headers=headers, data=data, timeout=30)
        if r.status_code == 200:
            return r.json().get("uuid")
        elif r.status_code == 400:
            return "INVALID_URL"
        else:
            log_error(f"Shazam error {r.status_code}")
            return None
    except Exception as e:
        log_error(f"Shazam exception: {e}")
        return None

def get_recognition_result(uuid):
    url = f"{SHAZAM_API_URL}/results/{uuid}"
    headers = {"Authorization": f"Bearer {SHAZAM_API_KEY}"}
    try:
        r = requests.post(url, headers=headers, timeout=30)
        if r.status_code == 200:
            return r.json()
        else:
            log_error(f"Shazam result error {r.status_code}")
            return None
    except Exception as e:
        log_error(f"Shazam result exception: {e}")
        return None

def poll_recognition(uuid, chat_id, thinking_message_id, bot_instance):
    for _ in range(30):
        time.sleep(3)
        res = get_recognition_result(uuid)
        if not res: continue
        status = res.get("status")
        if status == "completed":
            return res.get("results")
        elif status == "failed":
            return None
    return None

def format_recognition_results(results):
    if not results or not isinstance(results, list):
        return "❌ Не удалось распознать трек."
    text = "🔍 Найденные треки:\n\n"
    for idx, item in enumerate(results, 1):
        timecode = item.get("timecode", "?")
        track = item.get("track", {})
        title = track.get("title", "Неизвестно")
        subtitle = track.get("subtitle", "Неизвестно")
        text += f"{idx}. <code>{escape_html(timecode)}</code> – <b>{escape_html(title)}</b> – {escape_html(subtitle)}\n"
    return text

# ========== TIKTOK DOWNLOADER ==========

def is_valid_tiktok_url(url):
    tiktok_domains = ['vm.tiktok.com', 'vt.tiktok.com', 'www.tiktok.com', 'tiktok.com']
    try:
        parsed = urllib.parse.urlparse(url)
        return any(domain in parsed.netloc for domain in tiktok_domains)
    except:
        return False

def get_tiktok_video_tikwm(url):
    try:
        api_url = "https://www.tikwm.com/api/"
        payload = {
            "url": url,
            "count": 12,
            "cursor": 0,
            "web": 1,
            "hd": 1
        }
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        }
        response = requests.post(api_url, data=payload, headers=headers, timeout=15)
        if response.status_code == 200:
            data = response.json()
            if data.get('code') == 0:
                video_data = data.get('data', {})
                video_url = video_data.get('play', '')
                if video_url and not video_url.startswith('http'):
                    video_url = f"https://www.tikwm.com{video_url}"

                stats = video_data.get('stats', {})
                return {
                    'video_url': video_url,
                    'author': video_data.get('author', {}).get('nickname', 'Неизвестно'),
                    'description': video_data.get('title', ''),
                    'music': video_data.get('music_info', {}).get('title', ''),
                    'likes': stats.get('diggCount'),
                    'comments': stats.get('commentCount'),
                    'shares': stats.get('shareCount'),
                    'views': stats.get('playCount'),
                    'duration': video_data.get('duration')
                }
    except Exception as e:
        log_error(f"TikWM API error: {e}")
    return None

def get_tiktok_video_tikdown(url):
    try:
        api_url = "https://tikdown.org/api"
        payload = {"url": url}
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
            "Content-Type": "application/json"
        }
        response = requests.post(api_url, json=payload, headers=headers, timeout=15)
        if response.status_code == 200:
            data = response.json()
            if data.get('success'):
                return {
                    'video_url': data.get('videoUrl', ''),
                    'author': data.get('author', 'Неизвестно'),
                    'description': data.get('description', ''),
                    'music': data.get('music', '')
                }
    except Exception as e:
        log_error(f"TikDown API error: {e}")
    return None

def get_tiktok_video_info(url):
    info = get_tiktok_video_tikwm(url)
    if info and info.get('video_url'):
        return info
    info = get_tiktok_video_tikdown(url)
    if info and info.get('video_url'):
        return info
    return None

def download_tiktok_file(file_url, file_path):
    try:
        response = requests.get(file_url, stream=True, timeout=30)
        if response.status_code == 200:
            content_length = response.headers.get('content-length')
            if content_length and int(content_length) > 50 * 1024 * 1024:
                return False
            with open(file_path, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)
            return True
    except Exception as e:
        log_error(f"Download TikTok file error: {e}")
    return False

# ========== CRYPTO BOT FUNCTIONS ==========

def crypto_api_request(method, params=None):
    url = f"{CRYPTO_API_URL}/{method}"
    headers = {'Crypto-Pay-API-Token': CRYPTO_TOKEN, 'Content-Type': 'application/json'}
    try:
        r = requests.post(url, json=params, headers=headers, timeout=30)
        r.raise_for_status()
        data = r.json()
        if data.get('ok'):
            return data['result']
        else:
            log_error(f"Crypto API error: {data}")
            return None
    except Exception as e:
        log_error(f"Crypto API request error: {e}")
        return None

def create_crypto_invoice(amount_rub):
    amount_usdt = amount_rub / get_usd_rate()
    params = {
        'asset': 'USDT',
        'amount': f"{amount_usdt:.2f}",
        'description': f'Пополнение баланса на {amount_rub} руб.',
        'paid_btn_name': 'viewItem',
        'paid_btn_url': f'https://t.me/{BOT_NAME}',
        'allow_comments': False,
        'allow_anonymous': False,
        'expires_in': 300
    }
    return crypto_api_request('createInvoice', params)

def get_crypto_invoice_status(invoice_id):
    res = crypto_api_request('getInvoices', {'invoice_ids': [invoice_id]})
    return res['items'][0]['status'] if res and res.get('items') else None

def check_crypto_payment(inv_id, uid, rub, cid):
    time.sleep(5)
    start = time.time()
    while time.time() - start < 300:
        status = get_crypto_invoice_status(inv_id)
        if status == 'paid':
            update_balance(uid, rub*100, f"Crypto {rub} руб.")
            with db_lock:
                cursor.execute("UPDATE crypto_invoices SET status='paid' WHERE invoice_id=?", (inv_id,))
                conn.commit()
            bot.send_message(cid, get_text(uid, 'crypto_payment_success', amount=rub))
            return
        elif status == 'expired':
            bot.send_message(cid, get_text(uid, 'crypto_payment_failed'))
            return
        time.sleep(5)
    bot.send_message(cid, get_text(uid, 'crypto_payment_failed'))

# ========== ЗАЯВКИ НА КАРТУ ==========

def create_card_payment(user_id, amount_rub, first_name, last_name):
    comment = str(random.randint(10**9, 10**10-1))
    with db_lock:
        cursor.execute(
            "INSERT INTO card_payments (user_id, amount_rub, first_name, last_name, comment, created_at) VALUES (?,?,?,?,?,?)",
            (user_id, amount_rub, first_name, last_name, comment, datetime.now().isoformat())
        )
        conn.commit()
        return cursor.lastrowid, comment

def get_pending_card_requests():
    with db_lock:
        return cursor.execute("SELECT id, user_id, amount_rub, first_name, last_name, comment, created_at FROM card_payments WHERE status='pending' ORDER BY created_at DESC").fetchall()

def approve_card_payment(payment_id, admin_id):
    log_error(f"approve_card_payment called with payment_id={payment_id}, admin_id={admin_id}")
    with db_lock:
        try:
            row = cursor.execute("SELECT user_id, amount_rub FROM card_payments WHERE id=? AND status='pending'", (payment_id,)).fetchone()
            if not row:
                return False, "Заявка не найдена или уже обработана"
            user_id, rub = row
            update_balance(user_id, rub*100, reason=f"Пополнение картой (заявка #{payment_id})")
            cursor.execute("UPDATE card_payments SET status='approved' WHERE id=?", (payment_id,))
            conn.commit()
            return True, user_id
        except Exception as e:
            log_error(f"approve_card_payment: EXCEPTION: {traceback.format_exc()}")
            return False, f"Ошибка БД: {e}"

def reject_card_payment(payment_id):
    with db_lock:
        cursor.execute("UPDATE card_payments SET status='rejected' WHERE id=?", (payment_id,))
        conn.commit()
        return True

# ========== ТИКЕТЫ ==========

tickets = {}
def create_ticket(user_id, username, text, photo_id=None):
    ticket_id = str(uuid.uuid4())[:8]
    for admin_id in ADMIN_IDS:
        try:
            if photo_id:
                msg = bot.send_photo(admin_id, photo_id,
                    caption=f"<b>📨 Запрос в поддержку!</b>\n\n👤 Пользователь: @{username}\n🆔 ID: {user_id}\n📝 Проблема: {text}",
                    parse_mode="HTML")
            else:
                msg = bot.send_message(admin_id,
                    f"<b>📨 Запрос в поддержку!</b>\n\n👤 Пользователь: @{username}\n🆔 ID: {user_id}\n📝 Проблема: {text}",
                    parse_mode="HTML")
            tickets[msg.message_id] = {'user_id': user_id, 'text': text, 'photo_id': photo_id}
            markup = types.InlineKeyboardMarkup().add(
                types.InlineKeyboardButton("💬 Ответить", callback_data=f"reply_ticket:{user_id}:{msg.message_id}"),
                types.InlineKeyboardButton("❌ Скрыть", callback_data=f"close_ticket:{msg.message_id}")
            )
            bot.edit_message_reply_markup(admin_id, msg.message_id, reply_markup=markup)
        except Exception as e:
            log_error(f"Ошибка отправки тикета админу {admin_id}: {e}")

# ========== РЕФЕРАЛЬНАЯ СИСТЕМА ==========

def generate_referral_link(user_id):
    return f"https://t.me/{BOT_NAME}?start=ref_{user_id}"

def register_referral(referrer_id, referred_id, username, first_name):
    log_error(f"register_referral called: referrer={referrer_id}, referred={referred_id}, user={username}")
    with db_lock:
        # Проверяем, не зарегистрирован ли уже этот пользователь
        existing = cursor.execute("SELECT id FROM referrals WHERE referred_id=?", (referred_id,)).fetchone()
        if existing:
            log_error(f"Реферал {referred_id} уже зарегистрирован ранее")
            return False
        # Проверяем, не совпадает ли реферер с рефералом
        if referrer_id == referred_id:
            log_error(f"Попытка самореферала: {referrer_id}")
            return False
        # Проверяем, существует ли реферер в таблице users
        user_exists = cursor.execute("SELECT user_id FROM users WHERE user_id=?", (referrer_id,)).fetchone()
        if not user_exists:
            log_error(f"Реферер {referrer_id} не найден в базе пользователей")
            return False
        cursor.execute('''
            INSERT INTO referrals (referrer_id, referred_id, username, first_name, joined_date, status)
            VALUES (?, ?, ?, ?, ?, 'pending')
        ''', (referrer_id, referred_id, username, first_name, datetime.now().isoformat()))
        conn.commit()
        log_error(f"Реферал зарегистрирован: {referrer_id} -> {referred_id}")
        return True

def get_referral_count(user_id, status='pending'):
    with db_lock:
        row = cursor.execute("SELECT COUNT(*) FROM referrals WHERE referrer_id=? AND status=?", (user_id, status)).fetchone()
        return row[0] if row else 0

def get_referrals_list(user_id, status='pending'):
    with db_lock:
        rows = cursor.execute("SELECT referred_id, username, first_name, joined_date FROM referrals WHERE referrer_id=? AND status=?", (user_id, status)).fetchall()
        return rows

def send_gift_to_user(user_id, gift_file_id, caption):
    try:
        bot.send_photo(user_id, gift_file_id, caption=caption)
        return True
    except Exception as e:
        log_error(f"Ошибка отправки подарка пользователю {user_id}: {e}")
        return False

def send_gift_direct(chat_id: int, gift_id: int) -> bool:
    url = f"https://api.telegram.org/bot{TOKEN}/sendGift"
    payload = {
        'chat_id': chat_id,
        'gift_id': gift_id
    }
    try:
        response = requests.post(url, data=payload, timeout=10)
        result = response.json()
        if result.get('ok'):
            log_error(f"sendGift: успешно отправлен подарок {gift_id} пользователю {chat_id}")
            return True
        else:
            log_error(f"sendGift: ошибка API: {result}")
            return False
    except Exception as e:
        log_error(f"sendGift: исключение: {e}")
        return False

# ========== ПЕРЕВОДЫ ==========

def get_models_list():
    models_list = ""
    for name, mid in AVAILABLE_MODELS.items():
        models_list += f"• {name}\n"
    return models_list.strip()

LANGUAGES = {
    'ru': {
        'code': 'ru', 'name': 'Русский',
        'choose_language': '🌐 Выберите язык:',
        'language_changed': '✅ Язык изменён на русский.',
        'welcome': '<tg-emoji emoji-id="{welcome_star_id}">✨</tg-emoji> Привет!\nЭтот бот даёт вам доступ к лучшим нейросетям для создания текста, изображений и видео.\n\nДоступны модели:\n{models_list}\n\nБесплатно: Gemini 2.5 Flash Lite и GPT-5 Nano (20 запросов в день).\n\nЧатбот умеет:\n• Писать и переводить тексты <tg-emoji emoji-id="5260512129240276089">📝</tg-emoji>\n• Генерировать картинки и видео <tg-emoji emoji-id="{razmer_id}">📐</tg-emoji><tg-emoji emoji-id="{viborvideo_id}">🎥</tg-emoji>\n• Работать с документами <tg-emoji emoji-id="{documents_id}">📄</tg-emoji>\n• Писать и править код <tg-emoji emoji-id="{code_id}">💻</tg-emoji>\n• Решать математические задачи <tg-emoji emoji-id="{math_id}">🧮</tg-emoji>\n• Редактировать и распознавать фото <tg-emoji emoji-id="{photo_edit_id}">🖼️</tg-emoji>\n• Озвучивать текст и распознавать аудио <tg-emoji emoji-id="{audio_id}">🎧</tg-emoji>',
        'main_menu': '<tg-emoji emoji-id="{main_menu_id}">◀️</tg-emoji> Главное меню:',
        'generate': '<tg-emoji emoji-id="{generate_id}">🎨</tg-emoji> Сгенерировать',
        'edit': '<tg-emoji emoji-id="{edit_id}">✏️</tg-emoji> Редактировать фото',
        'video': '<tg-emoji emoji-id="{video_id}">🎥</tg-emoji> Генерация видео',
        'search': '<tg-emoji emoji-id="{search_id}">🔍</tg-emoji> Интернет поиск',
        'recognize': '<tg-emoji emoji-id="{recognize_id}">🎵</tg-emoji> Распознать музыку',
        'new_chat': '<tg-emoji emoji-id="{new_chat_id}">🆕</tg-emoji> Новый чат',
        'info': '<tg-emoji emoji-id="{info_id}">ℹ️</tg-emoji> Инфо',
        'contacts': '<tg-emoji emoji-id="{contacts_id}">📞</tg-emoji> Контакты',
        'voice': '<tg-emoji emoji-id="{voice_id}">🎤</tg-emoji> Голос {status}',
        'model': '<tg-emoji emoji-id="{model_id}">🤖</tg-emoji> Модель',
        'tts': '<tg-emoji emoji-id="{tts_id}">🔊</tg-emoji> Синтезатор',
        'help': '<tg-emoji emoji-id="{help_id}">❓</tg-emoji> Помощь',
        'topup': '<tg-emoji emoji-id="{topup_id}">💳</tg-emoji> Пополнить',
        'settings': '<tg-emoji emoji-id="{settings_id}">⚙️</tg-emoji> Настройки',
        'tts_voices': '<tg-emoji emoji-id="{tts_voice_emoji_id}">🎤</tg-emoji> Голоса',
        'tts_create': '<tg-emoji emoji-id="{tts_id}">🔊</tg-emoji> Создать',
        'choose_size': '<tg-emoji emoji-id="{razmer_id}">📐</tg-emoji> Выберите размер:',
        'choose_video_duration': '⏱ Выберите длительность видео (цена за генерацию):',
        'choose_mode': 'Выберите режим генерации:',
        'text_mode': '📝 Текст',
        'photo_mode': '🖼 Фото',
        'enter_prompt': '<tg-emoji emoji-id="{successrazmer_id}">✅</tg-emoji> Выбран размер {size} (цена {price:.2f} ₽). Введите описание:',
        'enter_video_text': 'Введите описание видео:',
        'send_photo_for_video': 'Отправьте фото для оживления (можно с подписью):',
        'enter_search_query': '<tg-emoji emoji-id="{search_input_id}">🔍</tg-emoji> Чтобы начать поиск, отправьте в чат ваш запрос 👇',
        'enter_tts_text': '<tg-emoji emoji-id="{tts_input_id}">🎤</tg-emoji> Введите текст для озвучивания:',
        'generating_image': '<tg-emoji emoji-id="{generate_id}">🎨</tg-emoji> Генерирую изображение...',
        'generating_video': '<tg-emoji emoji-id="{video_id}">🎥</tg-emoji> Генерирую видео (это может занять до минуты)...',
        'generating_tts': '<tg-emoji emoji-id="{tts_id}">🔊</tg-emoji> Синтезирую речь...',
        'searching': '<tg-emoji emoji-id="{search_id}">🔍</tg-emoji> Ищу в интернете...',
        'image_failed': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Не удалось сгенерировать изображение.',
        'video_failed': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Не удалось сгенерировать видео.',
        'tts_failed': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Сервис синтеза речи временно недоступен.',
        'search_failed': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Не удалось выполнить поиск.',
        'insufficient_funds': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Недостаточно средств.\nБаланс: {balance} ₽\nНужно: {price} ₽',
        'edit_no_photo': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Отправьте изображение с подписью.',
        'edit_failed': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Не удалось отредактировать.',
        'save_error': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Ошибка сохранения изображения на сервере.',
        'send_error': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Ошибка отправки: {error}',
        'transcribe_failed': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Не удалось распознать голос.',
        'ai_error': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Ошибка при обращении к AITUNNEL: {error}',
        'empty_response': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Пустой ответ от нейросети.',
        'canceled': '<tg-emoji emoji-id="{canceled_id}">❌</tg-emoji> Отменено.',
        'empty_prompt': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Пусто.',
        'blocked': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Вы заблокированы.',
        'no_rights': '⛔ Нет прав.',
        'history_cleared': '<tg-emoji emoji-id="{success_id}">✅</tg-emoji> История очищена.',
        'message_sent': '<tg-emoji emoji-id="{success_id}">✅</tg-emoji> Сообщение отправлено создателю.',
        'replied': '<tg-emoji emoji-id="{success_id}">✅</tg-emoji> Отправлено {target}.',
        'reply_error': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Ошибка: {error}',
        'thanks_review': 'Спасибо за отзыв!',
        'chat_ended': 'Чат завершён. Оцените качество обслуживания:',
        'rate_prompt': 'Вы выбрали {rating} {stars}. Напишите комментарий (или отправьте /skip чтобы пропустить):',
        'settings_menu': '<tg-emoji emoji-id="{settings_id}">⚙️</tg-emoji> <b>Настройки</b>\nВыберите параметр для изменения:',
        'language_settings': '<tg-emoji emoji-id="{settings_id}">🌐</tg-emoji> Язык',
        'current_language': 'Текущий язык: {lang}',
        'choose_language_prompt': '🌐 Выберите язык:',
        'language_updated': '<tg-emoji emoji-id="{success_id}">✅</tg-emoji> Язык изменён на {lang}.',
        'internet_search_toggle': '🔍 Режим поиска: {status}',
        'internet_search_on': 'вкл',
        'internet_search_off': 'выкл',
        'send_audio_or_link': '<tg-emoji emoji-id="{music_id}">🎵</tg-emoji> Отправьте аудиофайл или ссылку на музыку:',
        'invalid_url': '❌ Неверная ссылка. Отправьте прямую ссылку на аудио.',
        'recognition_result': '🔍 Результат распознавания:\n{result}',
        'recognizing_music': '🎵 Распознаю музыку...',
        'recognition_failed': '❌ Не удалось распознать музыку.',
        'crypto_enter_amount_rub': '<tg-emoji emoji-id="{crypto_id}">💵</tg-emoji> Введите сумму в рублях (минимум 30):',
        'crypto_invoice_created': '✅ Счёт создан!\n\nСсылка для оплаты: {pay_url}\n\nПосле оплаты нажмите кнопку ниже, чтобы проверить статус.',
        'crypto_check_payment': '🔄 Проверить оплату',
        'crypto_payment_success': '✅ Оплата прошла успешно! Баланс пополнен на {amount} руб.',
        'crypto_payment_pending': '⏳ Оплата ещё не поступила. Попробуйте позже.',
        'crypto_payment_failed': '❌ Счёт не оплачен или истёк.',
        'crypto_error': '❌ Ошибка при создании счёта. Попробуйте позже.',
        'stars_enter_amount': '<tg-emoji emoji-id="{stars_id}">⭐️</tg-emoji> Введите сумму в рублях, на которую хотите пополнить баланс (минимум 30 ₽).\nКурс: 1 звезда = {star_rate:.2f} ₽',
        'stars_invalid_amount': '❌ Сумма должна быть не менее 30 ₽ и целым числом.',
        'stars_invoice_title': 'Пополнение баланса',
        'stars_invoice_description': 'Пополнение на {amount_rub} ₽ ({stars} звёзд)',
        'card_enter_amount': '<tg-emoji emoji-id="{topup_id}">💳</tg-emoji> Введите сумму пополнения в рублях (минимум 30):',
        'card_enter_first_name': 'Введите ваше имя (как на карте):',
        'card_enter_last_name': 'Введите вашу фамилию:',
        'card_payment_details': '<tg-emoji emoji-id="{topup_id}">💳</tg-emoji> <b>Реквизиты для перевода</b>\n\n<tg-emoji emoji-id="{id_id}">🆔</tg-emoji> Карта: <code>{BANK_CARD}</code>\n<tg-emoji emoji-id="{id_id}">🆔</tg-emoji> ID пользователя: <code>{user_id}</code>\n<tg-emoji emoji-id="{balance_id}">💸</tg-emoji> Сумма: {amount} ₽\n<tg-emoji emoji-id="{profile_id}">👤</tg-emoji> Имя: {first_name}\n<tg-emoji emoji-id="{profile_id}">👤</tg-emoji> Фамилия: {last_name}\n<tg-emoji emoji-id="{contacts_id}">💬</tg-emoji> Комментарий: <code>{comment}</code>\n\nПосле перевода ожидайте подтверждения администратором.',
        'card_request_saved': '✅ Заявка сохранена. После подтверждения администратором баланс будет начислен.',
        'admin_card_requests': '📋 Заявки на карту',
        'card_request_pending': '🕒 Ожидает',
        'card_request_approved': '✅ Подтверждено',
        'card_request_rejected': '❌ Отклонено',
        'card_approve': '✅ Подтвердить',
        'card_reject': '❌ Отклонить',
        'card_approved': '✅ Заявка #{id} подтверждена. Баланс пользователя пополнен на {amount} ₽.',
        'card_rejected': '❌ Заявка #{id} отклонена.',
        'ticket_ask_problem': '<tg-emoji emoji-id="{help_emoji_id}">🧑‍💻</tg-emoji> Пожалуйста, расскажите о вашей проблеме. Если у вас есть возможность, прикрепите скриншот:',
        'ticket_created': '<b>✉️ Ваша заявка успешно создана.</b>\nОжидайте ответа администрации.',
        'ticket_user_reply_prompt': 'Введите ваш ответ для сотрудника поддержки:',
        'ticket_reply_helpful': 'Спасибо за обратную связь! Если возникнут вопросы, обращайтесь снова.',
        'news_channel': '<tg-emoji emoji-id="{news_id}">📰</tg-emoji> Наш канал: https://t.me/ken4kk_news',
        'choose_model': '<tg-emoji emoji-id="{choose_model_id}">🤖</tg-emoji> Выберите модель:',
        'tts_menu': '<tg-emoji emoji-id="{tts_menu_id}">🔊</tg-emoji> Меню синтезатора',
        'topup_menu': '<tg-emoji emoji-id="{topup_id}">💳</tg-emoji> Выберите способ:',
        'profile_menu': '👤 Профиль',
        'back_to_main': '◀️ Главное меню',
        'multiple_photos_edit': '📸 Получено несколько фото. Для редактирования используйте команду /edit и следуйте инструкциям.',
        'group_only_commands': 'В группах я отвечаю только на команды, начинающиеся с /',
        'send_tiktok_link': '📱 Отправьте ссылку на TikTok видео:',
        'editing_photo': '<tg-emoji emoji-id="{edit_id}">✏️</tg-emoji> Редактирую...',
        'text_too_long': '⚠️ Текст слишком длинный, первые 3000 символов.',
        'failed_load_photo': '❌ Не удалось загрузить фото.',
        'failed_load_voice': '❌ Не удалось загрузить голосовое.',
        'history_cleared_toast': 'История очищена',
        'docs_title': '📄 Актуальные цены',
        'docs_generate': '🎨 Генерация изображения:',
        'docs_edit': '✏️ Редактирование фото:',
        'docs_video': '🎥 Видео (за секунду):',
        'docs_tts': '🔊 TTS (за 1000 симв.):',
        'docs_voice': '🎤 Голосовой ввод:',
        'docs_search': '🔍 Интернет-поиск:',
        'docs_recognize': '🎵 Распознавание музыки:',
        'docs_tiktok': '📱 TikTok:',
        'docs_free': 'бесплатно',
        'analyzing': '🤔 Анализирую...',
        'reading_document': '🤔 Читаю документ...',
        'recognizing_voice': '<tg-emoji emoji-id="{voice_id}">🎤</tg-emoji> Распознаю...',
        'recognized_text': '📝 Распознано: {text}',
        'thinking_again': '🤔 Отвечаю...',
    },
    'en': {
        'code': 'en', 'name': 'English',
        'choose_language': '🌐 Choose language:',
        'language_changed': '✅ Language changed to English.',
        'welcome': '<tg-emoji emoji-id="{welcome_star_id}">✨</tg-emoji> Hello!\nThis bot gives you access to the best neural networks for text, images and videos.\n\nAvailable models:\n{models_list}\n\nFree: Gemini 2.5 Flash Lite and GPT-5 Nano (20 requests per day).\n\nChatbot can:\n• Write and translate texts <tg-emoji emoji-id="5260512129240276089">📝</tg-emoji>\n• Generate images and videos <tg-emoji emoji-id="{razmer_id}">📐</tg-emoji><tg-emoji emoji-id="{viborvideo_id}">🎥</tg-emoji>\n• Work with documents <tg-emoji emoji-id="{documents_id}">📄</tg-emoji>\n• Write and edit code <tg-emoji emoji-id="{code_id}">💻</tg-emoji>\n• Solve math problems <tg-emoji emoji-id="{math_id}">🧮</tg-emoji>\n• Edit and recognize photos <tg-emoji emoji-id="{photo_edit_id}">🖼️</tg-emoji>\n• Text-to-speech and audio recognition <tg-emoji emoji-id="{audio_id}">🎧</tg-emoji>',
        'main_menu': '<tg-emoji emoji-id="{main_menu_id}">◀️</tg-emoji> Main menu:',
        'generate': '<tg-emoji emoji-id="{generate_id}">🎨</tg-emoji> Generate',
        'edit': '<tg-emoji emoji-id="{edit_id}">✏️</tg-emoji> Edit photo',
        'video': '<tg-emoji emoji-id="{video_id}">🎥</tg-emoji> Generate video',
        'search': '<tg-emoji emoji-id="{search_id}">🔍</tg-emoji> Internet search',
        'recognize': '<tg-emoji emoji-id="{recognize_id}">🎵</tg-emoji> Recognize music',
        'new_chat': '<tg-emoji emoji-id="{new_chat_id}">🆕</tg-emoji> New chat',
        'info': '<tg-emoji emoji-id="{info_id}">ℹ️</tg-emoji> Info',
        'contacts': '<tg-emoji emoji-id="{contacts_id}">📞</tg-emoji> Contacts',
        'voice': '<tg-emoji emoji-id="{voice_id}">🎤</tg-emoji> Voice {status}',
        'model': '<tg-emoji emoji-id="{model_id}">🤖</tg-emoji> Model',
        'tts': '<tg-emoji emoji-id="{tts_id}">🔊</tg-emoji> TTS',
        'help': '<tg-emoji emoji-id="{help_id}">❓</tg-emoji> Help',
        'topup': '<tg-emoji emoji-id="{topup_id}">💳</tg-emoji> Top up',
        'settings': '<tg-emoji emoji-id="{settings_id}">⚙️</tg-emoji> Settings',
        'tts_voices': '<tg-emoji emoji-id="{tts_voice_emoji_id}">🎤</tg-emoji> Voices',
        'tts_create': '<tg-emoji emoji-id="{tts_id}">🔊</tg-emoji> Create',
        'choose_size': '<tg-emoji emoji-id="{razmer_id}">📐</tg-emoji> Choose size:',
        'choose_video_duration': '⏱ Choose video duration (price per generation):',
        'choose_mode': 'Choose generation mode:',
        'text_mode': '📝 Text',
        'photo_mode': '🖼 Photo',
        'enter_prompt': '<tg-emoji emoji-id="{successrazmer_id}">✅</tg-emoji> Selected size {size} (price {price:.2f} ₽). Enter description:',
        'enter_video_text': 'Enter video description:',
        'send_photo_for_video': 'Send a photo to animate (can have caption):',
        'enter_search_query': '<tg-emoji emoji-id="{search_input_id}">🔍</tg-emoji> To start searching, send your query 👇',
        'enter_tts_text': '<tg-emoji emoji-id="{tts_input_id}">🎤</tg-emoji> Enter text for speech synthesis:',
        'generating_image': '<tg-emoji emoji-id="{generate_id}">🎨</tg-emoji> Generating image...',
        'generating_video': '<tg-emoji emoji-id="{video_id}">🎥</tg-emoji> Generating video (may take up to a minute)...',
        'generating_tts': '<tg-emoji emoji-id="{tts_id}">🔊</tg-emoji> Synthesizing speech...',
        'searching': '<tg-emoji emoji-id="{search_id}">🔍</tg-emoji> Searching the internet...',
        'image_failed': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Failed to generate image.',
        'video_failed': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Failed to generate video.',
        'tts_failed': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Speech synthesis service temporarily unavailable.',
        'search_failed': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Search failed.',
        'insufficient_funds': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Insufficient funds.\nBalance: {balance} ₽\nNeed: {price} ₽',
        'edit_no_photo': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Send an image with caption.',
        'edit_failed': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Failed to edit.',
        'save_error': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Error saving image on server.',
        'send_error': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Send error: {error}',
        'transcribe_failed': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Failed to recognize voice.',
        'ai_error': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Error accessing AITUNNEL: {error}',
        'empty_response': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Empty response from neural network.',
        'canceled': '<tg-emoji emoji-id="{canceled_id}">❌</tg-emoji> Canceled.',
        'empty_prompt': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Empty.',
        'blocked': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> You are blocked.',
        'no_rights': '⛔ No rights.',
        'history_cleared': '<tg-emoji emoji-id="{success_id}">✅</tg-emoji> History cleared.',
        'message_sent': '<tg-emoji emoji-id="{success_id}">✅</tg-emoji> Message sent to creator.',
        'replied': '<tg-emoji emoji-id="{success_id}">✅</tg-emoji> Sent to {target}.',
        'reply_error': '<tg-emoji emoji-id="{warning_id}">❌</tg-emoji> Error: {error}',
        'thanks_review': 'Thank you for your feedback!',
        'chat_ended': 'Chat ended. Rate the quality of service:',
        'rate_prompt': 'You chose {rating} {stars}. Write a comment (or send /skip to skip):',
        'settings_menu': '<tg-emoji emoji-id="{settings_id}">⚙️</tg-emoji> <b>Settings</b>\nSelect a parameter to change:',
        'language_settings': '<tg-emoji emoji-id="{settings_id}">🌐</tg-emoji> Language',
        'current_language': 'Current language: {lang}',
        'choose_language_prompt': '🌐 Choose language:',
        'language_updated': '<tg-emoji emoji-id="{success_id}">✅</tg-emoji> Language changed to {lang}.',
        'internet_search_toggle': '🔍 Search mode: {status}',
        'internet_search_on': 'on',
        'internet_search_off': 'off',
        'send_audio_or_link': '<tg-emoji emoji-id="{music_id}">🎵</tg-emoji> Send an audio file or a link to music:',
        'invalid_url': '❌ Invalid URL. Send a direct audio link.',
        'recognition_result': '🔍 Recognition result:\n{result}',
        'recognizing_music': '🎵 Recognizing music...',
        'recognition_failed': '❌ Failed to recognize music.',
        'crypto_enter_amount_rub': '<tg-emoji emoji-id="{crypto_id}">💵</tg-emoji> Enter amount in rubles (minimum 30):',
        'crypto_invoice_created': '✅ Invoice created!\n\nPayment link: {pay_url}\n\nAfter payment, click the button below to check status.',
        'crypto_check_payment': '🔄 Check payment',
        'crypto_payment_success': '✅ Payment successful! Balance topped up by {amount} ₽.',
        'crypto_payment_pending': '⏳ Payment not yet received. Try later.',
        'crypto_payment_failed': '❌ Invoice not paid or expired.',
        'crypto_error': '❌ Error creating invoice. Try later.',
        'stars_enter_amount': '<tg-emoji emoji-id="{stars_id}">⭐️</tg-emoji> Enter amount in rubles you want to top up (minimum 30 ₽).\nRate: 1 star = {star_rate:.2f} ₽',
        'stars_invalid_amount': '❌ Amount must be at least 30 ₽ and an integer.',
        'stars_invoice_title': 'Top up balance',
        'stars_invoice_description': 'Top up {amount_rub} ₽ ({stars} stars)',
        'card_enter_amount': '<tg-emoji emoji-id="{topup_id}">💳</tg-emoji> Enter top-up amount in rubles (minimum 30):',
        'card_enter_first_name': 'Enter your first name (as on card):',
        'card_enter_last_name': 'Enter your last name:',
        'card_payment_details': '<tg-emoji emoji-id="{topup_id}">💳</tg-emoji> <b>Transfer details</b>\n\n<tg-emoji emoji-id="{id_id}">🆔</tg-emoji> Card: <code>{BANK_CARD}</code>\n<tg-emoji emoji-id="{id_id}">🆔</tg-emoji> User ID: <code>{user_id}</code>\n<tg-emoji emoji-id="{balance_id}">💸</tg-emoji> Amount: {amount} ₽\n<tg-emoji emoji-id="{profile_id}">👤</tg-emoji> First name: {first_name}\n<tg-emoji emoji-id="{profile_id}">👤</tg-emoji> Last name: {last_name}\n<tg-emoji emoji-id="{contacts_id}">💬</tg-emoji> Comment: <code>{comment}</code>\n\nAfter transfer, wait for admin confirmation.',
        'card_request_saved': '✅ Request saved. Balance will be credited after admin confirmation.',
        'admin_card_requests': '📋 Card requests',
        'card_request_pending': '🕒 Pending',
        'card_request_approved': '✅ Approved',
        'card_request_rejected': '❌ Rejected',
        'card_approve': '✅ Approve',
        'card_reject': '❌ Reject',
        'card_approved': '✅ Request #{id} approved. User balance topped up by {amount} ₽.',
        'card_rejected': '❌ Request #{id} rejected.',
        'ticket_ask_problem': '<tg-emoji emoji-id="{help_emoji_id}">🧑‍💻</tg-emoji> Please describe your problem. If possible, attach a screenshot:',
        'ticket_created': '<b>✉️ Your request has been created successfully.</b>\nWait for admin response.',
        'ticket_user_reply_prompt': 'Enter your reply for the support staff:',
        'ticket_reply_helpful': 'Thank you for your feedback! If you have any questions, feel free to contact again.',
        'news_channel': '<tg-emoji emoji-id="{news_id}">📰</tg-emoji> Our channel: https://t.me/ken4kk_news',
        'choose_model': '<tg-emoji emoji-id="{choose_model_id}">🤖</tg-emoji> Choose model:',
        'tts_menu': '<tg-emoji emoji-id="{tts_menu_id}">🔊</tg-emoji> TTS menu',
        'topup_menu': '<tg-emoji emoji-id="{topup_id}">💳</tg-emoji> Choose method:',
        'profile_menu': '👤 Profile',
        'back_to_main': '◀️ Main menu',
        'multiple_photos_edit': '📸 Received multiple photos. For editing, use the /edit command and follow the instructions.',
        'group_only_commands': 'In groups, I only respond to commands starting with /',
        'send_tiktok_link': '📱 Send a TikTok video link:',
        'editing_photo': '<tg-emoji emoji-id="{edit_id}">✏️</tg-emoji> Editing...',
        'text_too_long': '⚠️ Text is too long, first 3000 characters.',
        'failed_load_photo': '❌ Failed to load photo.',
        'failed_load_voice': '❌ Failed to load voice message.',
        'history_cleared_toast': 'History cleared',
        'docs_title': '📄 Current Prices',
        'docs_generate': '🎨 Generate image:',
        'docs_edit': '✏️ Edit photo:',
        'docs_video': '🎥 Video (per second):',
        'docs_tts': '🔊 TTS (per 1000 chars):',
        'docs_voice': '🎤 Voice input:',
        'docs_search': '🔍 Internet search:',
        'docs_recognize': '🎵 Recognize music:',
        'docs_tiktok': '📱 TikTok:',
        'docs_free': 'free',
        'analyzing': '🤔 Analyzing...',
        'reading_document': '🤔 Reading document...',
        'recognizing_voice': '<tg-emoji emoji-id="{voice_id}">🎤</tg-emoji> Recognizing...',
        'recognized_text': '📝 Recognized: {text}',
        'thinking_again': '🤔 Thinking...',
    }
}

THEMES = {'default': {'name': 'Классическая', 'emoji': '🔵'}}

def get_text(user_id, key, **kwargs):
    with db_lock:
        row = cursor.execute("SELECT language, theme, mode FROM users WHERE user_id=?", (user_id,)).fetchone()
        if row:
            lang = row[0] or 'ru'
            theme = row[1] or 'default'
            mode = row[2] or 'normal'
        else:
            lang = 'ru'; theme = 'default'; mode = 'normal'
    text = LANGUAGES.get(lang, LANGUAGES['ru']).get(key, key)
    mode_names = {'ru': {'normal': 'Обычный', 'adult': 'Без цензуры'}, 'en': {'normal': 'Normal', 'adult': 'Uncensored'}}
    # Исправление: если lang нет в mode_names, используем русский как запасной

    mode_dict = mode_names.get(lang, mode_names['ru'])
    format_dict = {
        'generate_id': PREMIUM_EMOJI.get('generate', ''),
        'edit_id': PREMIUM_EMOJI.get('edit', ''),
        'video_id': PREMIUM_EMOJI.get('video', ''),
        'viborvideo_id': PREMIUM_EMOJI.get('viborvideo', ''),
        'search_id': PREMIUM_EMOJI.get('search', ''),
        'recognize_id': PREMIUM_EMOJI.get('recognize', ''),
        'new_chat_id': PREMIUM_EMOJI.get('new_chat', ''),
        'info_id': PREMIUM_EMOJI.get('info', ''),
        'contacts_id': PREMIUM_EMOJI.get('contacts', ''),
        'voice_id': PREMIUM_EMOJI.get('voice', ''),
        'model_id': PREMIUM_EMOJI.get('model', ''),
        'tts_id': PREMIUM_EMOJI.get('tts', ''),
        'help_id': PREMIUM_EMOJI.get('help', ''),
        'topup_id': PREMIUM_EMOJI.get('topup', ''),
        'settings_id': PREMIUM_EMOJI.get('settings', ''),
        'main_menu_id': PREMIUM_EMOJI.get('main_menu', ''),
        'canceled_id': PREMIUM_EMOJI.get('canceled', ''),
        'success_id': PREMIUM_EMOJI.get('success', ''),
        'warning_id': PREMIUM_EMOJI.get('warning', ''),
        'sparkles_id': PREMIUM_EMOJI.get('sparkles', ''),
        'successrazmer_id': PREMIUM_EMOJI.get('successrazmer', ''),
        'razmer_id': PREMIUM_EMOJI.get('razmer', ''),
        'profile_id': PREMIUM_EMOJI.get('profile', ''),
        'purchases_id': PREMIUM_EMOJI.get('purchases', ''),
        'balance_id': PREMIUM_EMOJI.get('balance', ''),
        'tts_voice_emoji_id': PREMIUM_EMOJI.get('tts_voice_emoji', ''),
        'id_id': PREMIUM_EMOJI.get('id', ''),
        'ima_id': PREMIUM_EMOJI.get('ima', ''),
        'popolnenie_id': PREMIUM_EMOJI.get('popolnenie', ''),
        'stars_id': PREMIUM_EMOJI.get('stars', ''),
        'crypto_id': PREMIUM_EMOJI.get('crypto', ''),
        'help_emoji_id': PREMIUM_EMOJI.get('help_emoji', ''),
        'tts_menu_id': PREMIUM_EMOJI.get('tts_menu', ''),
        'choose_model_id': PREMIUM_EMOJI.get('choose_model', ''),
        'news_id': PREMIUM_EMOJI.get('news', ''),
        'music_id': PREMIUM_EMOJI.get('music', ''),
        'tts_input_id': PREMIUM_EMOJI.get('tts_input', ''),
        'search_input_id': PREMIUM_EMOJI.get('search_input', ''),
        'documents_id': PREMIUM_EMOJI.get('documents', ''),
        'code_id': PREMIUM_EMOJI.get('code', ''),
        'math_id': PREMIUM_EMOJI.get('math', ''),
        'photo_edit_id': PREMIUM_EMOJI.get('photo_edit', ''),
        'audio_id': PREMIUM_EMOJI.get('audio', ''),
        'welcome_star_id': PREMIUM_EMOJI.get('welcome_star', ''),
        'time_reset_id': PREMIUM_EMOJI.get('time_reset', ''),
        'subscribe_id': PREMIUM_EMOJI.get('subscribe', ''),
        'star_rate': get_star_rate(),
        'theme_emoji': THEMES.get(theme, THEMES['default'])['emoji'],
        'status': '', 'balance': 0, 'price': 0, 'size': '', 'target': '', 'error': '', 'code': '', 'rating': '', 'stars': '',
        'lang': get_language_name(lang), 'name': '', 'bot_name': BOT_NAME, 'creator': CREATOR_TAG, 'result': '', 'mode': mode_dict.get(mode, mode),
        'BANK_CARD': BANK_CARD,
        'models_list': get_models_list(),
    }
    format_dict['user_id'] = user_id
    for k, v in kwargs.items():
        if k == 'result':
            format_dict[k] = v
        elif isinstance(v, str):
            format_dict[k] = escape_html(v)
        else:
            format_dict[k] = v
    return text.format(**format_dict)

def get_language_name(lang_code):
    return LANGUAGES.get(lang_code, LANGUAGES['ru'])['name']

def get_theme_name(theme_code):
    return THEMES.get(theme_code, THEMES['default'])['name']

bot = telebot.TeleBot(TOKEN)
bot.timeout = 60
bot.long_polling_timeout = 120
user_states = {}

pending_media_groups = {}
media_group_lock = threading.Lock()

def process_media_group(media_group_id, uid, chat_id, state_info):
    with media_group_lock:
        if media_group_id not in pending_media_groups:
            return
        data = pending_media_groups.pop(media_group_id)
        images_bytes = data['images']
    if state_info:
        size, prompt = state_info
        process_edit_step_multiple(chat_id, uid, prompt, size, images_bytes)
    else:
        bot.send_message(chat_id, get_text(uid, 'multiple_photos_edit'))

# ========== КЛАВИАТУРЫ ==========

def show_profile(uid, chat_id, from_user):
    username = from_user.username
    first_name = from_user.first_name
    if username:
        user_display = f"@{username}"
    else:
        user_display = first_name or str(uid)
    try:
        bal = get_balance(uid) / 100
        spent = cursor.execute("SELECT SUM(amount) FROM transactions WHERE user_id=? AND amount<0", (uid,)).fetchone()[0] or 0
        purchases = -spent / 100
        free_left = get_free_quota_left(uid)
        free_tts_left = get_free_tts_quota_left(uid)

        now = datetime.now()
        next_reset = datetime.combine(date.today() + timedelta(days=1), datetime.min.time())
        delta = next_reset - now
        hours, remainder = divmod(delta.seconds, 3600)
        minutes = remainder // 60
        time_until_reset = f"{hours} ч {minutes} мин"

        reset_text = ""
        if free_left == 0 or free_tts_left == 0:
            reset_text = f'\n<tg-emoji emoji-id="{PREMIUM_EMOJI["time_reset"]}">⏳</tg-emoji> Восстановление через: {time_until_reset}'

        text = (f"➖➖➖➖➖➖➖➖➖➖\n"
                f"<tg-emoji emoji-id=\"{PREMIUM_EMOJI['profile']}\">ℹ️</tg-emoji> Личный кабинет:\n"
                f"<tg-emoji emoji-id=\"{PREMIUM_EMOJI['ima']}\">👤</tg-emoji> Имя: {escape_html(user_display)}\n"
                f"<tg-emoji emoji-id=\"{PREMIUM_EMOJI['id']}\">🆔</tg-emoji> ID: <code>{uid}</code>\n"
                f"<tg-emoji emoji-id=\"{PREMIUM_EMOJI['purchases']}\">🎁</tg-emoji> Покупок на сумму: {purchases:.2f} руб.\n"
                f"<tg-emoji emoji-id=\"{PREMIUM_EMOJI['balance']}\">💸</tg-emoji> Баланс: {bal:.2f} руб.\n"
                f"<tg-emoji emoji-id=\"{PREMIUM_EMOJI['sparkles']}\">✨</tg-emoji> Бесплатных текстовых запросов сегодня: {free_left}\n"
                f"<tg-emoji emoji-id=\"{PREMIUM_EMOJI['audio']}\">🎧</tg-emoji> Бесплатных TTS запросов сегодня: {free_tts_left}{reset_text}\n"
                f"<tg-emoji emoji-id=\"{PREMIUM_EMOJI['popolnenie']}\">📧</tg-emoji> Для пополнения баланса нажмите на кнопку ниже!\n"
                f"➖➖➖➖➖➖➖➖➖➖")
        markup = types.InlineKeyboardMarkup(row_width=2)
        markup.add(
            types.InlineKeyboardButton("💳 Пополнить", callback_data="topup_menu"),
            types.InlineKeyboardButton("◀️ Назад", callback_data="back_to_main")
        )
        bot.send_message(chat_id, text, parse_mode="HTML", reply_markup=markup)
    except Exception as e:
        log_error(f"Profile error: {e}")
        bot.send_message(chat_id, "❌ Произошла ошибка при загрузке профиля.")

def main_inline_keyboard(uid):
    voice_enabled = get_voice_enabled(uid)
    voice_status = "✅" if voice_enabled else "❌"
    search_enabled = get_internet_search_enabled(uid)
    search_status = get_text(uid, 'internet_search_on') if search_enabled else get_text(uid, 'internet_search_off')
    lang = get_language(uid)

    markup = types.InlineKeyboardMarkup(row_width=2)

    base_url = "https://ken4kk-app.ru"
    params = {
        'model': get_text_model(uid),
        'search': '1' if get_internet_search_enabled(uid) else '0',
        'voice': get_tts_voice(uid),
        'lang': get_language(uid),
        'mode': get_mode(uid),
    }
    query_string = '&'.join([f"{k}={v}" for k, v in params.items()])
    full_url = f"{base_url}?{query_string}"

    mini_app_button = types.InlineKeyboardButton(
        "🚀 Запустить Mini App" if lang == 'ru' else "🚀 Launch Mini App",
        web_app=types.WebAppInfo(url=full_url)
    )
    markup.row(mini_app_button)

    if lang == 'ru':
        buttons = [
            ("🎨 Сгенерировать", "generate"),
            ("✏️ Редактировать фото", "edit"),
            ("🎥 Генерация видео", "video"),
            (f"🔍 Режим поиска: {search_status}", "toggle_internet_search"),
            ("🎵 Распознать музыку", "recognize_music"),
            ("🎵 TikTok", "tiktok_downloader"),
            ("🆕 Новый чат", "new_chat"),
            ("📰 Новости", "news"),
            ("🎁 Подарки за друзей", "referrals"),
            ("📄 Документация", "docs"),
            ("📜 Пользовательское соглашение", "user_agreement"),
            (f"🎤 Голос {voice_status}", "toggle_voice"),
            ("🤖 Модель", "choose_model"),
            ("🔊 Синтезатор", "tts_menu"),
            ("❓ Помощь", "help"),
            ("💳 Пополнить", "topup_menu"),
            ("⚙️ Настройки", "settings_menu"),
        ]
    else:
        buttons = [
            ("🎨 Generate", "generate"),
            ("✏️ Edit photo", "edit"),
            ("🎥 Generate video", "video"),
            (f"🔍 Search mode: {search_status}", "toggle_internet_search"),
            ("🎵 Recognize music", "recognize_music"),
            ("🎵 TikTok", "tiktok_downloader"),
            ("🆕 New chat", "new_chat"),
            ("📰 News", "news"),
            ("🎁 Referral Gifts", "referrals"),
            ("📄 Documentation", "docs"),
            ("📜 User Agreement", "user_agreement"),
            (f"🎤 Voice {voice_status}", "toggle_voice"),
            ("🤖 Model", "choose_model"),
            ("🔊 TTS", "tts_menu"),
            ("❓ Help", "help"),
            ("💳 Top up", "topup_menu"),
            ("⚙️ Settings", "settings_menu"),
        ]

    if is_admin(uid):
        buttons.append(("🔧 Admin panel", "admin_panel"))

    markup.add(*[types.InlineKeyboardButton(t, callback_data=c) for t, c in buttons])

    profile_text = "👤 Профиль" if lang == 'ru' else "👤 Profile"
    markup.add(types.InlineKeyboardButton(profile_text, callback_data="profile"))

    return markup

def model_choice_keyboard(uid):
    current = get_text_model(uid)
    markup = types.InlineKeyboardMarkup(row_width=1)
    for name, mid in AVAILABLE_MODELS.items():
        text = f"✅ {name}" if mid == current else name
        markup.add(types.InlineKeyboardButton(text, callback_data=f"set_model_{mid}"))
    lang = get_language(uid)
    back_text = "◀️ Главное меню" if lang == 'ru' else "◀️ Main menu"
    markup.add(types.InlineKeyboardButton(back_text, callback_data="back_to_main"))
    return markup

def topup_menu_keyboard(uid):
    lang = get_language(uid)
    markup = types.InlineKeyboardMarkup(row_width=2)
    markup.add(
        types.InlineKeyboardButton("⭐ Звёзды" if lang == 'ru' else "⭐ Stars", callback_data="stars_topup"),
        types.InlineKeyboardButton("💳 Карта" if lang == 'ru' else "💳 Card", callback_data="card_payment"),
        types.InlineKeyboardButton("🤖 Крипто бот" if lang == 'ru' else "🤖 Crypto Bot", callback_data="crypto_topup"),
        types.InlineKeyboardButton("◀️ Главное меню" if lang == 'ru' else "◀️ Main menu", callback_data="back_to_main")
    )
    return markup

def tts_menu_keyboard(uid):
    lang = get_language(uid)
    markup = types.InlineKeyboardMarkup(row_width=1)
    markup.add(
        types.InlineKeyboardButton("🎤 Голоса" if lang == 'ru' else "🎤 Voices", callback_data="choose_tts_voice"),
        types.InlineKeyboardButton("🔊 Создать" if lang == 'ru' else "🔊 Create", callback_data="tts_synthesize"),
        types.InlineKeyboardButton("◀️ Главное меню" if lang == 'ru' else "◀️ Main menu", callback_data="back_to_main")
    )
    return markup

def tts_voice_keyboard(uid):
    current = get_tts_voice(uid)
    markup = types.InlineKeyboardMarkup(row_width=1)
    for voice, desc in AVAILABLE_TTS_VOICES.items():
        text = f"{'✅ ' if voice == current else ''}{voice} — {desc}"
        markup.add(types.InlineKeyboardButton(text, callback_data=f"set_tts_voice_{voice}"))
    lang = get_language(uid)
    back_text = "◀️ Назад" if lang == 'ru' else "◀️ Back"
    markup.add(types.InlineKeyboardButton(back_text, callback_data="tts_menu"))
    return markup

def contacts_menu_keyboard(uid):
    lang = get_language(uid)
    markup = types.InlineKeyboardMarkup(row_width=1)
    markup.add(
        types.InlineKeyboardButton("📩 Написать в ЛС" if lang == 'ru' else "📩 Write to PM", url=f"https://t.me/KeN4kk_n1"),
        types.InlineKeyboardButton("📞 Контакты" if lang == 'ru' else "📞 Contacts", callback_data="contact_via_bot"),
        types.InlineKeyboardButton("◀️ Главное меню" if lang == 'ru' else "◀️ Main menu", callback_data="back_to_main")
    )
    return markup

def image_size_keyboard(uid, action):
    markup = types.InlineKeyboardMarkup(row_width=1)
    if action == 'generate' or action == 'edit':
        base_price = get_media_base_price('generate')
        price_per_image = base_price * get_profit_multiplier()
        for size, display in IMAGE_SIZES.items():
            price_display = f"{display} — {price_per_image:.2f} ₽"
            markup.add(types.InlineKeyboardButton(price_display, callback_data=f"size_{action}_{size}"))
    else:
        for size, display in IMAGE_SIZES.items():
            markup.add(types.InlineKeyboardButton(display, callback_data=f"size_{action}_{size}"))
    lang = get_language(uid)
    cancel_text = "◀️ Отмена" if lang == 'ru' else "◀️ Cancel"
    markup.add(types.InlineKeyboardButton(cancel_text, callback_data="back_to_main"))
    return markup

def video_duration_keyboard(uid):
    markup = types.InlineKeyboardMarkup(row_width=2)
    base_price_per_sec = get_media_base_price('video_per_second')
    multiplier = get_profit_multiplier()
    for dur in VIDEO_DURATIONS:
        price = dur * base_price_per_sec * multiplier
        markup.add(types.InlineKeyboardButton(f"{dur} сек — {price:.2f} ₽", callback_data=f"vdur_{dur}"))
    lang = get_language(uid)
    cancel_text = "◀️ Отмена" if lang == 'ru' else "◀️ Cancel"
    markup.add(types.InlineKeyboardButton(cancel_text, callback_data="back_to_main"))
    return markup

def video_mode_keyboard(uid, duration):
    lang = get_language(uid)
    text_mode = "📝 Текст" if lang == 'ru' else "📝 Text"
    photo_mode = "🖼 Фото" if lang == 'ru' else "🖼 Photo"
    cancel_text = "◀️ Отмена" if lang == 'ru' else "◀️ Cancel"
    markup = types.InlineKeyboardMarkup(row_width=2)
    markup.add(
        types.InlineKeyboardButton(text_mode, callback_data=f"vmode_text_{duration}"),
        types.InlineKeyboardButton(photo_mode, callback_data=f"vmode_photo_{duration}"),
        types.InlineKeyboardButton(cancel_text, callback_data="back_to_main")
    )
    return markup

def rating_keyboard(uid):
    markup = types.InlineKeyboardMarkup(row_width=5)
    buttons = [types.InlineKeyboardButton("⭐", callback_data=f"rate_{i}") for i in range(1,6)]
    markup.add(*buttons)
    return markup

def language_keyboard(uid):
    markup = types.InlineKeyboardMarkup(row_width=2)
    for code, data in LANGUAGES.items():
        markup.add(types.InlineKeyboardButton(data['name'], callback_data=f"set_lang_{code}"))
    lang = get_language(uid)
    back_text = "◀️ Назад" if lang == 'ru' else "◀️ Back"
    markup.add(types.InlineKeyboardButton(back_text, callback_data="settings_menu"))
    return markup

def settings_menu_keyboard(uid):
    lang = get_language(uid)
    markup = types.InlineKeyboardMarkup(row_width=1)
    markup.add(
        types.InlineKeyboardButton("🌐 Язык" if lang == 'ru' else "🌐 Language", callback_data="settings_language"),
        types.InlineKeyboardButton("◀️ Главное меню" if lang == 'ru' else "◀️ Main menu", callback_data="back_to_main")
    )
    return markup

def admin_panel_keyboard(uid):
    lang = get_language(uid)
    if lang == 'ru':
        buttons = [
            ("📋 Список пользователей", "admin_list"),
            ("💰 Изменить баланс (+/-)", "admin_balance"),
            ("📢 Рассылка (как бот)", "admin_broadcast_as_bot"),
            ("📢 Рассылка (с подписью)", "admin_broadcast"),
            ("🚫 Заблокировать", "admin_block"),
            ("✅ Разблокировать", "admin_unblock"),
            ("🔚 Завершить чат", "admin_end_chat"),
            ("📊 Статистика", "admin_stats"),
            ("⚙️ Управление ценами", "admin_prices_management"),
            ("💰 Курсы валют", "admin_rates"),
            ("📜 Логи пользователя", "admin_user_logs"),
            ("🤖 Управление текстовыми моделями", "admin_text_models"),
            ("🎤 Управление TTS моделями", "admin_tts_models"),
            ("🎥 Управление видео моделями", "admin_video_models"),
            ("🖼 Управление моделями изображений", "admin_image_models"),
            ("👑 Управление админами", "admin_admins"),
            ("📜 История платежей", "admin_transactions"),
            ("📋 Заявки на карту", "admin_card_requests"),
            ("🎁 Заявки на подарки", "admin_gift_requests"),
            ("🔙 Назад", "back_to_main")
        ]
    else:
        buttons = [
            ("📋 User list", "admin_list"),
            ("💰 Change balance (+/-)", "admin_balance"),
            ("📢 Broadcast (as bot)", "admin_broadcast_as_bot"),
            ("📢 Broadcast (with signature)", "admin_broadcast"),
            ("🚫 Block", "admin_block"),
            ("✅ Unblock", "admin_unblock"),
            ("🔚 End chat", "admin_end_chat"),
            ("📊 Statistics", "admin_stats"),
            ("⚙️ Manage prices", "admin_prices_management"),
            ("💰 Exchange rates", "admin_rates"),
            ("📜 User logs", "admin_user_logs"),
            ("🤖 Manage text models", "admin_text_models"),
            ("🎤 Manage TTS models", "admin_tts_models"),
            ("🎥 Manage video models", "admin_video_models"),
            ("🖼 Manage image models", "admin_image_models"),
            ("👑 Manage admins", "admin_admins"),
            ("📜 Transaction history", "admin_transactions"),
            ("📋 Card requests", "admin_card_requests"),
            ("🎁 Gift requests", "admin_gift_requests"),
            ("🔙 Back", "back_to_main")
        ]
    markup = types.InlineKeyboardMarkup(row_width=2)
    markup.add(*[types.InlineKeyboardButton(t, callback_data=c) for t, c in buttons])
    return markup

def referrals_keyboard(uid):
    lang = get_language(uid)
    if lang == 'ru':
        buttons = [
            ("📤 Отправить заявку", "send_referral_request"),
            ("◀️ Главное меню", "back_to_main")
        ]
    else:
        buttons = [
            ("📤 Send request", "send_referral_request"),
            ("◀️ Main menu", "back_to_main")
        ]
    markup = types.InlineKeyboardMarkup(row_width=1)
    markup.add(*[types.InlineKeyboardButton(t, callback_data=c) for t, c in buttons])
    return markup

def text_models_keyboard(uid):
    markup = types.InlineKeyboardMarkup(row_width=2)
    markup.add(
        types.InlineKeyboardButton("📋 Список", callback_data="admin_text_list"),
        types.InlineKeyboardButton("➕ Добавить", callback_data="admin_text_add"),
        types.InlineKeyboardButton("✏️ Изменить цену", callback_data="admin_text_edit_price"),
        types.InlineKeyboardButton("❌ Удалить", callback_data="admin_text_delete"),
        types.InlineKeyboardButton("◀️ Назад", callback_data="admin_panel")
    )
    return markup

def tts_models_keyboard(uid):
    markup = types.InlineKeyboardMarkup(row_width=2)
    markup.add(
        types.InlineKeyboardButton("📋 Список", callback_data="admin_tts_list"),
        types.InlineKeyboardButton("➕ Добавить", callback_data="admin_tts_add"),
        types.InlineKeyboardButton("❌ Удалить", callback_data="admin_tts_delete"),
        types.InlineKeyboardButton("◀️ Назад", callback_data="admin_panel")
    )
    return markup

def video_models_keyboard(uid):
    markup = types.InlineKeyboardMarkup(row_width=2)
    markup.add(
        types.InlineKeyboardButton("📋 Список", callback_data="admin_video_list"),
        types.InlineKeyboardButton("➕ Добавить", callback_data="admin_video_add"),
        types.InlineKeyboardButton("❌ Удалить", callback_data="admin_video_delete"),
        types.InlineKeyboardButton("◀️ Назад", callback_data="admin_panel")
    )
    return markup

def image_models_keyboard(uid):
    markup = types.InlineKeyboardMarkup(row_width=2)
    markup.add(
        types.InlineKeyboardButton("📋 Список", callback_data="admin_image_list"),
        types.InlineKeyboardButton("➕ Добавить", callback_data="admin_image_add"),
        types.InlineKeyboardButton("❌ Удалить", callback_data="admin_image_delete"),
        types.InlineKeyboardButton("◀️ Назад", callback_data="admin_panel")
    )
    return markup

def admin_admins_keyboard(uid):
    markup = types.InlineKeyboardMarkup(row_width=2)
    markup.add(
        types.InlineKeyboardButton("➕ Добавить админа", callback_data="admin_add_admin"),
        types.InlineKeyboardButton("➖ Удалить админа", callback_data="admin_remove_admin"),
        types.InlineKeyboardButton("📋 Список админов", callback_data="admin_list_admins"),
        types.InlineKeyboardButton("◀️ Назад", callback_data="admin_panel")
    )
    return markup

def price_edit_keyboard(uid):
    markup = types.InlineKeyboardMarkup(row_width=2)
    for service in PRICES.keys():
        markup.add(types.InlineKeyboardButton(service, callback_data=f"admin_setprice_{service}"))
    markup.add(types.InlineKeyboardButton("◀️ Назад", callback_data="admin_panel"))
    return markup

def card_requests_keyboard(uid, payments):
    markup = types.InlineKeyboardMarkup(row_width=1)
    for pid, uid_user, rub, fn, ln, comment, ts in payments:
        text = f"#{pid} | {rub}₽ | {fn} {ln} | {ts[:10]}"
        markup.add(types.InlineKeyboardButton(text, callback_data=f"card_request_{pid}"))
    markup.add(types.InlineKeyboardButton("◀️ Назад", callback_data="admin_panel"))
    return markup

def card_request_action_keyboard(pid):
    markup = types.InlineKeyboardMarkup(row_width=2)
    markup.add(
        types.InlineKeyboardButton("✅ Подтвердить", callback_data=f"card_approve_{pid}"),
        types.InlineKeyboardButton("❌ Отклонить", callback_data=f"card_reject_{pid}"),
        types.InlineKeyboardButton("◀️ Назад", callback_data="admin_card_requests")
    )
    return markup

def admin_prices_management_keyboard(uid):
    lang = get_language(uid)
    if lang == 'ru':
        buttons = [
            ("📈 Коэффициент наценки", "admin_edit_profit_multiplier"),
            ("🤖 Цены текстовых моделей (AITUNNEL)", "admin_edit_aitunnel_prices"),
            ("🖼 Базовая цена изображения", "admin_edit_media_generate"),
            ("🎥 Базовая цена видео (за сек)", "admin_edit_media_video"),
            ("🔊 Базовая цена TTS (за 1000 симв.)", "admin_edit_media_tts"),
            ("🎤 Цена распознавания голоса", "admin_edit_media_voice"),
            ("◀️ Назад", "admin_panel")
        ]
    else:
        buttons = [
            ("📈 Profit Multiplier", "admin_edit_profit_multiplier"),
            ("🤖 AITUNNEL Text Model Prices", "admin_edit_aitunnel_prices"),
            ("🖼 Base Image Price", "admin_edit_media_generate"),
            ("🎥 Base Video Price (per sec)", "admin_edit_media_video"),
            ("🔊 Base TTS Price (per 1000 chars)", "admin_edit_media_tts"),
            ("🎤 Voice Recognition Price", "admin_edit_media_voice"),
            ("◀️ Back", "admin_panel")
        ]
    markup = types.InlineKeyboardMarkup(row_width=1)
    markup.add(*[types.InlineKeyboardButton(t, callback_data=c) for t, c in buttons])
    return markup

def aitunnel_price_edit_keyboard(uid):
    markup = types.InlineKeyboardMarkup(row_width=1)
    for model_id in AITUNNEL_PRICES.keys():
        display_name = model_id
        for name, mid in AVAILABLE_MODELS.items():
            if mid == model_id:
                display_name = name
                break
        markup.add(types.InlineKeyboardButton(display_name, callback_data=f"aem_{model_id}"))
    markup.add(types.InlineKeyboardButton("➕ Добавить модель", callback_data="admin_add_aitunnel_model"))
    markup.add(types.InlineKeyboardButton("◀️ Назад", callback_data="admin_prices_management"))
    return markup

def aitunnel_model_price_edit_keyboard(uid, model_id):
    markup = types.InlineKeyboardMarkup(row_width=2)
    markup.add(
        types.InlineKeyboardButton("🔹 Input price", callback_data=f"asp_{model_id}_input"),
        types.InlineKeyboardButton("🔸 Output price", callback_data=f"asp_{model_id}_output")
    )
    markup.add(types.InlineKeyboardButton("◀️ Назад", callback_data="admin_edit_aitunnel_prices"))
    return markup

# ========== КОМАНДЫ ==========

@bot.message_handler(commands=['start'], func=lambda m: m.chat.type == 'private')
@subscription_required
def start_private(message):
    uid = message.from_user.id
    register_or_update_user(message)

    parts = message.text.split()
    if len(parts) > 1 and parts[1].startswith('ref_'):
        try:
            referrer_id = int(parts[1].replace('ref_', ''))
            if referrer_id != uid:
                username = message.from_user.username or ""
                first_name = message.from_user.first_name or ""
                if register_referral(referrer_id, uid, username, first_name):
                    # Отправляем уведомление рефереру
                    try:
                        bot.send_message(referrer_id, 
                            f"🎉 По вашей ссылке зарегистрировался новый пользователь: @{username} ({first_name})\n"
                            f"Осталось приглашений до подарка: {5 - get_referral_count(referrer_id, 'pending')}")
                    except Exception as e:
                        log_error(f"Не удалось уведомить реферера {referrer_id}: {e}")
        except ValueError:
            log_error(f"Неверный формат реферального ID: {parts[1]}")
        except Exception as e:
            log_error(f"Ошибка при обработке реферала: {e}")

    if is_blocked(uid):
        bot.send_message(message.chat.id, get_text(uid, 'blocked'), parse_mode="HTML")
        return
    text = get_text(uid, 'welcome')
    bot.send_message(message.chat.id, text, parse_mode="HTML", reply_markup=main_inline_keyboard(uid))

@bot.message_handler(commands=['start'], func=lambda m: m.chat.type in ['group', 'supergroup'])
def start_group(message):
    uid = message.from_user.id
    register_or_update_user(message)
    if is_blocked(uid):
        bot.reply_to(message, get_text(uid, 'blocked'), parse_mode="HTML")
        return
    text = get_text(uid, 'welcome') + "\n\n" + get_text(uid, 'group_only_commands')
    bot.send_message(message.chat.id, text, parse_mode="HTML", reply_markup=main_inline_keyboard(uid))

@bot.message_handler(commands=['help'])
@subscription_required
def help_command(message):
    uid = message.from_user.id
    lang = get_language(uid)
    if lang == 'ru':
        help_text = (
            "📚 <b>Справка</b>\n\n"
            "<b>🔹 Команды в личном чате:</b>\n"
            f"• /start – регистрация и приветствие\n"
            f"• /help – эта справка\n"
            f"• /profile – ваш профиль (баланс, статистика)\n"
            f"• /settings – настройки (язык)\n"
            f"• /model – выбор модели нейросети\n"
            f"• /generate – генерация изображения ({PRICES['generate']/100:.2f} ₽)\n"
            f"• /edit – редактировать фото ({PRICES['edit']/100:.2f} ₽)\n"
            f"• /video – генерация видео (от {min(PRICES['video_6s'], PRICES['video_10s'])/100:.2f} ₽)\n"
            f"• /tts – синтез речи ({PRICES['tts']/100:.2f} ₽)\n"
            f"• /voice – голосовой ввод ({PRICES['voice']/100:.2f} ₽)\n"
            f"• /search – интернет-поиск ({PRICES['search']/100:.2f} ₽)\n"
            f"• /recognize – распознать музыку (бесплатно)\n"
            f"• /tiktok – скачать TikTok видео (бесплатно)\n"
            f"• /new_chat – очистить историю\n"
            f"• /info – информация о боте\n"
            f"• /topup – пополнить баланс\n"
            f"• /contacts – связь с создателем\n\n"
            "<b>🔸 Команды в группах:</b>\n"
            f"• /ask – задать вопрос боту\n"
            f"• /generate – генерация изображения ({PRICES['generate']/100:.2f} ₽)\n"
            f"• /edit – редактировать фото ({PRICES['edit']/100:.2f} ₽)\n"
            f"• /video – генерация видео (от {min(PRICES['video_6s'], PRICES['video_10s'])/100:.2f} ₽)\n"
            f"• /tts – синтез речи ({PRICES['tts']/100:.2f} ₽)\n"
            f"• /search – интернет-поиск ({PRICES['search']/100:.2f} ₽)\n"
            f"• /recognize – распознать музыку (бесплатно)\n"
            f"• /tiktok – скачать TikTok видео (бесплатно)\n"
        )
    else:
        help_text = (
            "📚 <b>Help</b>\n\n"
            "<b>🔹 Commands in private chat:</b>\n"
            f"• /start – registration and welcome\n"
            f"• /help – this help\n"
            f"• /profile – your profile (balance, stats)\n"
            f"• /settings – settings (language)\n"
            f"• /model – choose AI model\n"
            f"• /generate – generate image ({PRICES['generate']/100:.2f} ₽)\n"
            f"• /edit – edit photo ({PRICES['edit']/100:.2f} ₽)\n"
            f"• /video – generate video (from {min(PRICES['video_6s'], PRICES['video_10s'])/100:.2f} ₽)\n"
            f"• /tts – text-to-speech ({PRICES['tts']/100:.2f} ₽)\n"
            f"• /voice – voice input ({PRICES['voice']/100:.2f} ₽)\n"
            f"• /search – internet search ({PRICES['search']/100:.2f} ₽)\n"
            f"• /recognize – recognize music (free)\n"
            f"• /tiktok – download TikTok video (free)\n"
            f"• /new_chat – clear history\n"
            f"• /info – bot info\n"
            f"• /topup – top up balance\n"
            f"• /contacts – contact creator\n\n"
            "<b>🔸 Commands in groups:</b>\n"
            f"• /ask – ask a question\n"
            f"• /generate – generate image ({PRICES['generate']/100:.2f} ₽)\n"
            f"• /edit – edit photo ({PRICES['edit']/100:.2f} ₽)\n"
            f"• /video – generate video (from {min(PRICES['video_6s'], PRICES['video_10s'])/100:.2f} ₽)\n"
            f"• /tts – text-to-speech ({PRICES['tts']/100:.2f} ₽)\n"
            f"• /search – internet search ({PRICES['search']/100:.2f} ₽)\n"
            f"• /recognize – recognize music (free)\n"
            f"• /tiktok – download TikTok video (free)\n"
        )
    bot.reply_to(message, help_text, parse_mode="HTML")

@bot.message_handler(commands=['users'])
@subscription_required
def users_count(message):
    if not is_admin(message.from_user.id):
        bot.reply_to(message, get_text(message.from_user.id, 'no_rights'), parse_mode="HTML")
        return
    total = cursor.execute("SELECT COUNT(*) FROM users").fetchone()[0]
    bot.reply_to(message, f"👥 Всего пользователей: <b>{total}</b>", parse_mode="HTML")

@bot.message_handler(commands=['settings'])
@subscription_required
def settings_command(message):
    uid = message.from_user.id
    bot.send_message(uid, get_text(uid, 'settings_menu'), reply_markup=settings_menu_keyboard(uid), parse_mode="HTML")

@bot.message_handler(commands=['model'])
@subscription_required
def model_command(message):
    uid = message.from_user.id
    text = get_text(uid, 'choose_model')
    bot.send_message(uid, text, reply_markup=model_choice_keyboard(uid), parse_mode="HTML")

@bot.message_handler(commands=['profile'])
@subscription_required
def profile_command(message):
    uid = message.from_user.id
    show_profile(uid, message.chat.id, message.from_user)

@bot.message_handler(commands=['search'])
@subscription_required
def search_command(message):
    uid = message.from_user.id
    if is_blocked(uid):
        bot.reply_to(message, get_text(uid, 'blocked'), parse_mode="HTML")
        return
    user_states[uid] = "awaiting_search_query"
    markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Отмена", callback_data="back_to_main"))
    bot.reply_to(message, get_text(uid, 'enter_search_query'), reply_markup=markup, parse_mode="HTML")

@bot.message_handler(commands=['tiktok'])
@subscription_required
def tiktok_command(message):
    uid = message.from_user.id
    if is_blocked(uid):
        bot.reply_to(message, get_text(uid, 'blocked'), parse_mode="HTML")
        return
    user_states[uid] = "awaiting_tiktok_url"
    markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Отмена", callback_data="back_to_main"))
    bot.reply_to(message, get_text(uid, 'send_tiktok_link'), reply_markup=markup, parse_mode="HTML")

@bot.message_handler(commands=['aibalance'])
@subscription_required
def aitunnel_balance(message):
    if not is_admin(message.from_user.id):
        bot.reply_to(message, get_text(message.from_user.id, 'no_rights'), parse_mode="HTML")
        return
    try:
        r = requests.get(f"{AITUNNEL_URL}/aitunnel/balance", headers={"Authorization": f"Bearer {AITUNNEL_KEY}"})
        if r.status_code == 200:
            data = r.json()
            bot.reply_to(message, f"💰 Баланс AITUNNEL: <b>{data.get('balance','?')}</b> руб.", parse_mode="HTML")
        else:
            bot.reply_to(message, f"❌ Ошибка {r.status_code}")
    except Exception as e:
        bot.reply_to(message, f"❌ Ошибка: {e}")

@bot.message_handler(commands=['genbalance'])
@subscription_required
def genapi_balance(message):
    if not is_admin(message.from_user.id):
        bot.reply_to(message, get_text(message.from_user.id, 'no_rights'), parse_mode="HTML")
        return
    try:
        r = requests.get("https://api.gen-api.ru/api/v1/balance", headers={"Authorization": f"Bearer {GENAPI_KEY}"})
        if r.status_code == 200:
            data = r.json()
            bot.reply_to(message, f"💰 Баланс Gen-API: <b>{data.get('balance','?')}</b> руб.", parse_mode="HTML")
        else:
            bot.reply_to(message, f"❌ Ошибка {r.status_code}")
    except Exception as e:
        bot.reply_to(message, f"❌ Ошибка: {e}")

@bot.message_handler(commands=['admin'])
@subscription_required
def admin_panel(message):
    if not is_admin(message.from_user.id):
        bot.reply_to(message, get_text(message.from_user.id, 'no_rights'), parse_mode="HTML")
        return
    bot.send_message(message.chat.id, "🔧 <b>Админ-панель</b>", parse_mode="HTML", reply_markup=admin_panel_keyboard(message.from_user.id))

@bot.message_handler(commands=['addbalance'])
@subscription_required
def add_balance(message):
    if not is_admin(message.from_user.id):
        bot.reply_to(message, get_text(message.from_user.id, 'no_rights'), parse_mode="HTML")
        return
    try:
        parts = message.text.split()
        if len(parts) != 3:
            bot.reply_to(message, "Использование: /addbalance <user_id> <сумма в рублях> (положительная или отрицательная)")
            return
        target = int(parts[1])
        rub = float(parts[2])
        kop = int(round(rub * 100))
        update_balance(target, kop, reason="admin_add via command")
        bot.reply_to(message, f"✅ Баланс пользователя {target} изменён на {rub:.2f} ₽. Теперь: {get_balance(target)/100:.2f} ₽")
    except Exception as e:
        bot.reply_to(message, f"❌ Ошибка: {e}")

@bot.message_handler(commands=['block'])
@subscription_required
def block_user_cmd(message):
    if not is_admin(message.from_user.id):
        bot.reply_to(message, get_text(message.from_user.id, 'no_rights'), parse_mode="HTML")
        return
    parts = message.text.split()
    if len(parts) != 2:
        bot.reply_to(message, "Использование: /block <user_id>")
        return
    try:
        target = int(parts[1])
        block_user(target)
        bot.reply_to(message, f"✅ Пользователь {target} заблокирован.")
    except Exception as e:
        bot.reply_to(message, f"❌ Ошибка: {e}")

@bot.message_handler(commands=['setgiftid'])
@subscription_required
def set_gift_id(message):
    if not is_admin(message.from_user.id):
        bot.reply_to(message, get_text(message.from_user.id, 'no_rights'), parse_mode="HTML")
        return
    parts = message.text.split()
    if len(parts) != 2:
        bot.reply_to(message, "Использование: /setgiftid <число>")
        return
    try:
        gift_id = int(parts[1])
        set_setting('gift_id', str(gift_id))
        bot.reply_to(message, f"✅ ID подарка установлен: {gift_id}")
    except ValueError:
        bot.reply_to(message, "❌ Введите целое число.")

@bot.message_handler(commands=['unblock'])
@subscription_required
def unblock_user_cmd(message):
    if not is_admin(message.from_user.id):
        bot.reply_to(message, get_text(message.from_user.id, 'no_rights'), parse_mode="HTML")
        return
    parts = message.text.split()
    if len(parts) != 2:
        bot.reply_to(message, "Использование: /unblock <user_id>")
        return
    try:
        target = int(parts[1])
        unblock_user(target)
        bot.reply_to(message, f"✅ Пользователь {target} разблокирован.")
    except Exception as e:
        bot.reply_to(message, f"❌ Ошибка: {e}")

@bot.message_handler(commands=['broadcast'])
@subscription_required
def broadcast(message):
    if not is_admin(message.from_user.id):
        bot.reply_to(message, get_text(message.from_user.id, 'no_rights'), parse_mode="HTML")
        return
    text = message.text.replace('/broadcast', '', 1).strip()
    if not text:
        bot.reply_to(message, "Введите текст для рассылки.")
        return
    users = cursor.execute("SELECT user_id FROM users").fetchall()
    sent = failed = 0
    status = bot.send_message(message.chat.id, "⏳ Рассылка...")
    for (target,) in users:
        try:
            bot.send_message(target, f"📢 <b>Сообщение от администратора:</b>\n\n{escape_html(text)}", parse_mode="HTML")
            sent += 1
        except:
            failed += 1
        time.sleep(0.05)
    bot.edit_message_text(f"✅ Отправлено: {sent}, не удалось: {failed}", message.chat.id, status.message_id)

@bot.message_handler(commands=['userinfo'])
@subscription_required
def userinfo(message):
    if not is_admin(message.from_user.id):
        bot.reply_to(message, get_text(message.from_user.id, 'no_rights'), parse_mode="HTML")
        return
    parts = message.text.split()
    if len(parts) != 2:
        bot.reply_to(message, "Использование: /userinfo <user_id>")
        return
    try:
        user_id = int(parts[1])
        row = cursor.execute("SELECT username, first_name, balance, blocked, reg_date, text_model, voice_enabled, language, mode FROM users WHERE user_id=?", (user_id,)).fetchone()
        if not row:
            bot.reply_to(message, "Пользователь не найден.")
            return
        uname, fname, bal, blocked, reg, model, voice, lang, mode = row
        safe_uname = escape_html(uname); safe_fname = escape_html(fname); safe_model = escape_html(model)
        mode_name = get_text(user_id, 'mode_normal') if mode == 'normal' else get_text(user_id, 'mode_adult')
        text = f"👤 Информация о пользователе {user_id}\nUsername: @{safe_uname}\nИмя: {safe_fname}\nДата регистрации: {reg[:10]}\nЗаблокирован: {'Да' if blocked else 'Нет'}\nБаланс: {bal/100:.2f} ₽\nМодель: {safe_model}\nГолос: {'Вкл' if voice else 'Выкл'}\nЯзык: {lang}\nРежим: {mode_name}"
        bot.reply_to(message, text)
    except Exception as e:
        bot.reply_to(message, f"❌ Ошибка: {e}")

# ========== ГРУППОВЫЕ КОМАНДЫ ==========

def is_group_message(message): return message.chat.type in ['group','supergroup']
def is_private(message): return message.chat.type == 'private'

@bot.message_handler(commands=['ask'], func=is_group_message)
def ask_in_group(message):
    uid = message.from_user.id
    if is_blocked(uid):
        bot.reply_to(message, get_text(uid, 'blocked'), parse_mode="HTML")
        return
    question = message.text.replace('/ask', '', 1).strip()
    if not question:
        bot.reply_to(message, "Введите вопрос после /ask", parse_mode="HTML")
        return
    thinking = bot.reply_to(message, "🤔 Думаю...")
    answer = ask_aitunnel(message.chat.id, question)
    log_query(uid, question, answer)
    if not answer:
        bot.edit_message_text(get_text(uid, 'empty_response'), message.chat.id, thinking.message_id, parse_mode="HTML")
        return
    bot.delete_message(thinking.chat.id, thinking.message_id)
    if len(answer) > MAX_MESSAGE_LENGTH:
        handle_long_response(message.chat.id, answer, uid, reply_to_message_id=message.message_id)
    else:
        send_long_message(message.chat.id, answer, reply_to_message_id=message.message_id)

@bot.message_handler(content_types=['photo'], func=lambda m: is_group_message(m) and m.caption and m.caption.startswith('/ask'))
def handle_group_photo(message):
    uid = message.from_user.id
    if is_blocked(uid):
        bot.reply_to(message, get_text(uid, 'blocked'), parse_mode="HTML")
        return
    caption = message.caption
    question = caption.replace('/ask', '', 1).strip() or "Что на изображении?"
    image_bytes = download_file(message.photo[-1].file_id)
    if not image_bytes:
        bot.reply_to(message, get_text(uid, 'failed_load_photo'), parse_mode="HTML")
        return
    thinking = bot.reply_to(message, get_text(uid, 'analyzing'), parse_mode="HTML")
    answer = ask_aitunnel(message.chat.id, question, image_bytes)
    log_query(uid, f"[Фото] {question}", answer)
    if not answer:
        bot.edit_message_text(get_text(uid, 'empty_response'), message.chat.id, thinking.message_id, parse_mode="HTML")
        return
    bot.delete_message(thinking.chat.id, thinking.message_id)
    if len(answer) > MAX_MESSAGE_LENGTH:
        handle_long_response(message.chat.id, answer, uid, reply_to_message_id=message.message_id)
    else:
        send_long_message(message.chat.id, answer, reply_to_message_id=message.message_id)

@bot.message_handler(commands=['generate'], func=is_group_message)
def group_generate(message):
    uid = message.from_user.id
    if is_blocked(uid):
        bot.reply_to(message, get_text(uid, 'blocked'), parse_mode="HTML")
        return
    prompt = message.text.replace('/generate', '', 1).strip()
    if not prompt:
        bot.reply_to(message, "Введите описание после /generate", parse_mode="HTML")
        return
    generate_and_send(message, prompt, DEFAULT_IMAGE_SIZE)

@bot.message_handler(commands=['edit'], func=is_group_message)
def group_edit(message):
    uid = message.from_user.id
    if is_blocked(uid):
        bot.reply_to(message, get_text(uid, 'blocked'), parse_mode="HTML")
        return
    prompt = message.text.replace('/edit', '', 1).strip()
    if not prompt:
        bot.reply_to(message, "Введите описание изменений после /edit, например: /edit добавить шляпу")
        return
    user_states[uid] = f"awaiting_group_edit:{prompt}"
    bot.reply_to(message, "📷 Отправьте изображение, которое нужно отредактировать (можно с подписью, но описание уже задано).")

@bot.message_handler(commands=['video'], func=is_group_message)
def group_video(message):
    uid = message.from_user.id
    if is_blocked(uid):
        bot.reply_to(message, get_text(uid, 'blocked'), parse_mode="HTML")
        return
    prompt = message.text.replace('/video', '', 1).strip()
    if not prompt:
        bot.reply_to(message, "Введите описание после /video, например: /video кот играет с мячом")
        return
    user_states[uid] = f"awaiting_group_video_text:{prompt}"
    bot.reply_to(message, "Видео будет сгенерировано с длительностью 6 секунд. Ожидайте...")
    base_price_per_sec = get_media_base_price('video_per_second')
    price = int(6 * base_price_per_sec * get_profit_multiplier() * 100)
    if get_balance(uid) < price:
        bot.reply_to(message, get_text(uid, 'insufficient_funds', balance=get_balance(uid)/100, price=price/100), parse_mode="HTML")
        return
    thinking = bot.reply_to(message, get_text(uid, 'generating_video'), parse_mode="HTML")
    result = generate_video(prompt, 6, image_bytes=None)
    if not result:
        bot.edit_message_text(get_text(uid, 'video_failed'), message.chat.id, thinking.message_id, parse_mode="HTML")
        return
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.mp4') as tmp:
            tmp.write(result); tmp_path = tmp.name
        with open(tmp_path,'rb') as f:
            bot.send_video(message.chat.id, f, caption=f"🎥 {escape_html(prompt)}", timeout=600, reply_to_message_id=message.message_id)
        os.unlink(tmp_path)
        update_balance(uid, -price, reason=f"Видео в группе: {prompt[:50]}")
        bot.delete_message(message.chat.id, thinking.message_id)
    except Exception as e:
        bot.edit_message_text(f"❌ Ошибка отправки видео: {e}", message.chat.id, thinking.message_id)

@bot.message_handler(commands=['tts'], func=is_group_message)
def group_tts(message):
    uid = message.from_user.id
    if is_blocked(uid):
        bot.reply_to(message, get_text(uid, 'blocked'), parse_mode="HTML")
        return
    text = message.text.replace('/tts', '', 1).strip()
    if not text:
        bot.reply_to(message, "Введите текст после /tts", parse_mode="HTML")
        return
    if len(text) > 3000:
        text = text[:3000] + "..."
        bot.reply_to(message, get_text(uid, 'text_too_long'), parse_mode="HTML")

    free_used = False
    if is_free_tts_available(uid):
        use_free_tts_quota(uid)
        price = 0
        free_used = True
    else:
        base_tts_price = get_media_base_price('tts_per_1000_chars')
        char_count = len(text)
        price = int((char_count / 1000) * base_tts_price * get_profit_multiplier() * 100)
        price = max(price, 1)
        if get_balance(uid) < price:
            bot.reply_to(message, get_text(uid, 'insufficient_funds', balance=get_balance(uid)/100, price=price/100), parse_mode="HTML")
            return

    thinking = bot.reply_to(message, get_text(uid, 'generating_tts'), parse_mode="HTML")
    voice = get_tts_voice(uid)
    try:
        audio_bytes = generate_tts_genapi(text, voice)
        if not audio_bytes:
            bot.edit_message_text(get_text(uid, 'tts_failed'), chat_id=message.chat.id, message_id=thinking.message_id, parse_mode="HTML")
            return
        with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as tmp:
            tmp.write(audio_bytes); tmp_path = tmp.name
        caption = f'<tg-emoji emoji-id="{PREMIUM_EMOJI["tts"]}">🔊</tg-emoji> Текст: {escape_html(text[:200])}...' if len(text) > 200 else f'Текст: {escape_html(text)}'
        with open(tmp_path, 'rb') as f:
            bot.send_audio(message.chat.id, f, caption=caption, title="Синтезированная речь", performer="KeN4kk_AI", timeout=600, parse_mode="HTML", reply_to_message_id=message.message_id)
        os.unlink(tmp_path)
        if not free_used:
            update_balance(uid, -price, reason=f"TTS группа: {text[:50]}")
        bot.delete_message(thinking.chat.id, thinking.message_id)
    except Exception as e:
        log_error(f"Group TTS error: {e}")
        bot.edit_message_text(f"❌ Ошибка: {e}", message.chat.id, thinking.message_id, parse_mode="HTML")

@bot.message_handler(commands=['search'], func=is_group_message)
def group_search(message):
    uid = message.from_user.id
    if is_blocked(uid):
        bot.reply_to(message, get_text(uid, 'blocked'), parse_mode="HTML")
        return
    query = message.text.replace('/search', '', 1).strip()
    if not query:
        bot.reply_to(message, "Введите запрос после /search", parse_mode="HTML")
        return
    thinking = bot.reply_to(message, get_text(uid, 'searching'), parse_mode="HTML")
    answer = search_internet(query, uid)
    if not answer:
        bot.edit_message_text(get_text(uid, 'search_failed'), message.chat.id, thinking.message_id, parse_mode="HTML")
        return
    bot.delete_message(thinking.chat.id, thinking.message_id)
    if len(answer) > MAX_MESSAGE_LENGTH:
        handle_long_response(message.chat.id, answer, uid, reply_to_message_id=message.message_id)
    else:
        send_long_message(message.chat.id, answer, reply_to_message_id=message.message_id)

@bot.message_handler(commands=['tiktok'], func=is_group_message)
def group_tiktok(message):
    uid = message.from_user.id
    if is_blocked(uid):
        bot.reply_to(message, get_text(uid, 'blocked'), parse_mode="HTML")
        return
    user_states[uid] = "awaiting_tiktok_url"
    markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="back_to_main"))
    bot.reply_to(message, get_text(uid, 'send_tiktok_link'), reply_markup=markup, parse_mode="HTML")

@bot.message_handler(commands=['recognize'], func=is_group_message)
def group_recognize(message):
    uid = message.from_user.id
    if is_blocked(uid):
        bot.reply_to(message, get_text(uid, 'blocked'), parse_mode="HTML")
        return
    user_states[uid] = "awaiting_music_recognition"
    markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="back_to_main"))
    bot.reply_to(message, get_text(uid, 'send_audio_or_link'), reply_markup=markup, parse_mode="HTML")

# ========== ОБРАБОТЧИК ФОТО В ГРУППЕ ДЛЯ /edit ==========

@bot.message_handler(content_types=['photo'], func=is_group_message)
def handle_group_edit_photo(message):
    uid = message.from_user.id
    if is_blocked(uid):
        bot.reply_to(message, get_text(uid, 'blocked'), parse_mode="HTML")
        return
    if uid in user_states and isinstance(user_states[uid], str) and user_states[uid].startswith("awaiting_group_edit:"):
        prompt = user_states[uid].split(":", 1)[1]
        del user_states[uid]
        process_edit_step(message, prompt, DEFAULT_IMAGE_SIZE)
        return

# ========== ОБРАБОТЧИКИ ПРИВАТНЫХ СООБЩЕНИЙ ==========

@bot.message_handler(func=lambda m: is_private(m) and m.text, content_types=['text'])
@subscription_required
def handle_private_text(message):
    global ADMIN_IDS
    uid = message.from_user.id
    chat_id = message.chat.id

    if get_in_chat(uid):
        forward_to_admin(message, uid, 'text', message.text)
        bot.reply_to(message, get_text(uid, 'message_sent'), parse_mode="HTML")
        return

    if uid in ADMIN_IDS and uid in user_states and isinstance(user_states[uid], str) and user_states[uid].startswith("replying_to_"):
        target = int(user_states[uid].split("_")[2])
        try:
            bot.send_message(target, f"✉️ Ответ от создателя:\n{escape_html(message.text)}")
            set_last_admin_message(target, message.text)
            bot.reply_to(message, get_text(uid, 'replied', target=target), parse_mode="HTML")
        except Exception as e:
            bot.reply_to(message, get_text(uid, 'reply_error', error=str(e)), parse_mode="HTML")
        return

    if uid in ADMIN_IDS and isinstance(user_states.get(uid), dict):
        state = user_states[uid]
        if state.get("state") == "admin_reply_to_ticket":
            target_user_id = state["target_user_id"]
            ticket_msg_id = state["ticket_msg_id"]
            reply_text = message.caption if message.caption else message.text
            reply_photo = message.photo[-1].file_id if message.photo else None
            markup_user = types.InlineKeyboardMarkup().add(
                types.InlineKeyboardButton("💬 Ответить", callback_data="user_reply"),
                types.InlineKeyboardButton("Ответ помог 👍", callback_data="helpful")
            )
            try:
                if reply_photo:
                    bot.send_photo(target_user_id, reply_photo, caption=reply_text, reply_markup=markup_user)
                else:
                    bot.send_message(target_user_id, reply_text, reply_markup=markup_user)
                bot.send_message(chat_id, "✅ Ответ отправлен пользователю.")
            except Exception as e:
                bot.send_message(chat_id, f"❌ Ошибка: {e}")
            if ticket_msg_id in tickets:
                markup = types.InlineKeyboardMarkup().add(
                    types.InlineKeyboardButton("💬 Ответить", callback_data=f"reply_ticket:{target_user_id}:{ticket_msg_id}"),
                    types.InlineKeyboardButton("❌ Скрыть", callback_data=f"close_ticket:{ticket_msg_id}")
                )
                try:
                    bot.edit_message_reply_markup(chat_id, ticket_msg_id, reply_markup=markup)
                except: pass
            del user_states[uid]
            return

    if isinstance(user_states.get(uid), dict):
        state = user_states[uid]
        if state.get("state") == "awaiting_crypto_amount":
            text = message.text.strip()
            if text.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, get_text(uid, 'canceled'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
                del user_states[uid]
                return
            try:
                rub = int(text)
                if rub < 30:
                    bot.send_message(chat_id, "❌ Минимальная сумма 30 рублей.")
                    return
                invoice = create_crypto_invoice(rub)
                if not invoice:
                    bot.send_message(chat_id, get_text(uid, 'crypto_error'))
                    del user_states[uid]
                    return
                with db_lock:
                    cursor.execute(
                        "INSERT INTO crypto_invoices (invoice_id, user_id, amount_rub, status, created_at, pay_url) VALUES (?,?,?,?,?,?)",
                        (invoice['invoice_id'], uid, rub, 'active', datetime.now().isoformat(), invoice['pay_url'])
                    )
                    conn.commit()
                markup = types.InlineKeyboardMarkup().add(
                    types.InlineKeyboardButton(get_text(uid, 'crypto_check_payment'), callback_data=f"crypto_check:{invoice['invoice_id']}"),
                    types.InlineKeyboardButton("◀️ Назад", callback_data="topup_menu")
                )
                bot.send_message(chat_id, get_text(uid, 'crypto_invoice_created', pay_url=invoice['pay_url']), reply_markup=markup)
                threading.Thread(target=check_crypto_payment, args=(invoice['invoice_id'], uid, rub, chat_id), daemon=True).start()
                del user_states[uid]
            except ValueError:
                bot.send_message(chat_id, "❌ Введите целое число.")
            return

        if state.get("state") == "awaiting_stars_amount":
            text = message.text.strip()
            if text.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, get_text(uid, 'canceled'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
                del user_states[uid]
                return
            try:
                rub = int(text)
                if rub < 30:
                    bot.send_message(chat_id, get_text(uid, 'stars_invalid_amount'))
                    return
                star_rate = get_star_rate()
                stars = int((rub + star_rate - 0.0001) // star_rate)
                prices = [types.LabeledPrice(label=get_text(uid, 'stars_invoice_title'), amount=stars)]
                bot.send_invoice(chat_id, title=get_text(uid, 'stars_invoice_title'),
                                 description=get_text(uid, 'stars_invoice_description', amount_rub=rub, stars=stars),
                                 invoice_payload=f"stars_{uid}_{rub}", provider_token="", currency="XTR", prices=prices,
                                 start_parameter="stars_topup")
                del user_states[uid]
            except ValueError:
                bot.send_message(chat_id, get_text(uid, 'stars_invalid_amount'))
            return

        if state.get("state") == "awaiting_card_amount":
            text = message.text.strip()
            if text.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, get_text(uid, 'canceled'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
                del user_states[uid]
                return
            try:
                rub = int(text)
                if rub < 30:
                    bot.send_message(chat_id, get_text(uid, 'card_enter_amount') + " (минимум 30)")
                    return
                user_states[uid] = {"state": "awaiting_card_first_name", "amount": rub}
                bot.send_message(chat_id, get_text(uid, 'card_enter_first_name'))
            except ValueError:
                bot.send_message(chat_id, get_text(uid, 'card_enter_amount') + " (целое число)")
            return

        if state.get("state") == "awaiting_card_first_name":
            first_name = message.text.strip()
            if first_name.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, get_text(uid, 'canceled'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
                del user_states[uid]
                return
            if not first_name:
                bot.send_message(chat_id, get_text(uid, 'card_enter_first_name'))
                return
            user_states[uid] = {"state": "awaiting_card_last_name", "amount": state["amount"], "first_name": first_name}
            bot.send_message(chat_id, get_text(uid, 'card_enter_last_name'))
            return

        if state.get("state") == "awaiting_card_last_name":
            last_name = message.text.strip()
            if last_name.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, get_text(uid, 'canceled'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
                del user_states[uid]
                return
            if not last_name:
                bot.send_message(chat_id, get_text(uid, 'card_enter_last_name'))
                return
            pid, comment = create_card_payment(uid, state["amount"], state["first_name"], last_name)
            text = get_text(uid, 'card_payment_details', amount=state["amount"], first_name=state["first_name"], last_name=last_name, comment=comment)
            bot.send_message(chat_id, text, parse_mode="HTML")
            bot.send_message(chat_id, get_text(uid, 'card_request_saved'))
            user_row = cursor.execute("SELECT username FROM users WHERE user_id=?", (uid,)).fetchone()
            username = user_row[0] if user_row else None
            user_display = f"@{username}" if username else f"ID: {uid}"
            for admin in ADMIN_IDS:
                try:
                    bot.send_message(admin, 
                        f"📋 <b>Новая заявка на карту</b>\n\n"
                        f"🆔 ID пользователя: <code>{uid}</code>\n"
                        f"👤 Пользователь: {user_display}\n"
                        f"💰 Сумма: {state['amount']} ₽\n"
                        f"👤 Имя: {state['first_name']}\n"
                        f"👤 Фамилия: {last_name}\n"
                        f"💬 Комментарий: <code>{comment}</code>\n"
                        f"📅 Время: {datetime.now().strftime('%d.%m.%Y %H:%M')}\n\n"
                        f"Заявка #{pid}",
                        parse_mode="HTML"
                    )
                except Exception as e:
                    log_error(f"Error sending card request to admin {admin}: {e}")
            del user_states[uid]
            return

    if uid in user_states and isinstance(user_states[uid], str):
        state = user_states[uid]

        if state == "admin_waiting_usd_rate":
            rate_str = message.text.strip()
            if rate_str.lower() in ["отмена","cancel"]:
                bot.send_message(chat_id, "❌ Отменено.", reply_markup=admin_panel_keyboard(uid))
                del user_states[uid]
                return
            try:
                rate = float(rate_str)
                set_usd_rate(rate)
                bot.send_message(chat_id, f"✅ Курс USD установлен: {rate:.2f}", reply_markup=admin_panel_keyboard(uid))
            except ValueError:
                bot.send_message(chat_id, "❌ Неверный формат. Введите число, например 80.5")
                return
            del user_states[uid]
            return

        if state == "admin_waiting_star_rate":
            rate_str = message.text.strip()
            if rate_str.lower() in ["отмена","cancel"]:
                bot.send_message(chat_id, "❌ Отменено.", reply_markup=admin_panel_keyboard(uid))
                del user_states[uid]
                return
            try:
                rate = float(rate_str)
                set_star_rate(rate)
                bot.send_message(chat_id, f"✅ Курс звезды установлен: {rate:.2f}", reply_markup=admin_panel_keyboard(uid))
            except ValueError:
                bot.send_message(chat_id, "❌ Неверный формат. Введите число, например 1.5")
                return
            del user_states[uid]
            return

        if state == "admin_waiting_add_admin":
            target_str = message.text.strip()
            if target_str.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, "❌ Отменено.", reply_markup=admin_panel_keyboard(uid))
                del user_states[uid]
                return
            try:
                target = int(target_str)
                user = cursor.execute("SELECT user_id FROM users WHERE user_id=?", (target,)).fetchone()
                if not user:
                    bot.send_message(chat_id, "❌ Пользователь с таким ID не найден.")
                    return
                cursor.execute("UPDATE users SET is_admin=1 WHERE user_id=?", (target,))
                conn.commit()
                if target not in ADMIN_IDS:
                    ADMIN_IDS.append(target)
                bot.send_message(chat_id, f"✅ Пользователь {target} теперь администратор.", reply_markup=admin_panel_keyboard(uid))
            except ValueError:
                bot.send_message(chat_id, "❌ Введите корректный ID.")
            finally:
                del user_states[uid]
            return

        if state == "admin_waiting_remove_admin":
            target_str = message.text.strip()
            if target_str.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, "❌ Отменено.", reply_markup=admin_panel_keyboard(uid))
                del user_states[uid]
                return
            try:
                target = int(target_str)
                if target not in ADMIN_IDS:
                    bot.send_message(chat_id, "❌ Этот пользователь не является администратором.")
                    return
                if target in [6197133464, 7588557213]:
                    bot.send_message(chat_id, "❌ Нельзя удалить главного администратора.")
                    return
                cursor.execute("UPDATE users SET is_admin=0 WHERE user_id=?", (target,))
                conn.commit()
                ADMIN_IDS.remove(target)
                bot.send_message(chat_id, f"✅ Пользователь {target} лишён прав администратора.", reply_markup=admin_panel_keyboard(uid))
            except ValueError:
                bot.send_message(chat_id, "❌ Введите корректный ID.")
            finally:
                del user_states[uid]
            return

        if state == "admin_waiting_userid_for_logs":
            target_str = message.text.strip()
            if target_str.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, "❌ Отменено.", reply_markup=admin_panel_keyboard(uid))
                del user_states[uid]
                return
            try:
                target = int(target_str)
                logs = cursor.execute(
                    "SELECT query, answer, timestamp FROM queries_log WHERE user_id=? ORDER BY timestamp DESC LIMIT 20",
                    (target,)
                ).fetchall()
                if not logs:
                    bot.send_message(chat_id, "📭 У пользователя нет запросов.")
                else:
                    text = f"📜 Последние 20 запросов пользователя {target}:\n\n"
                    for q, a, ts in logs:
                        text += f"🕒 {ts[:19]}\n❓ {q}\n✅ {a[:100]}...\n\n"
                    send_long_message(chat_id, text)
                bot.send_message(chat_id, "Выберите действие:", reply_markup=admin_panel_keyboard(uid))
            except ValueError:
                bot.send_message(chat_id, "❌ Введите корректный ID.")
            finally:
                del user_states[uid]
            return

        if state == "awaiting_ticket_problem":
            text = message.caption if message.caption else message.text
            photo_id = message.photo[-1].file_id if message.photo else None
            create_ticket(uid, message.from_user.username or "None", text, photo_id)
            bot.send_message(chat_id, get_text(uid, 'ticket_created'), parse_mode="HTML")
            del user_states[uid]
            return

        if state == "awaiting_music_recognition":
            text = message.text.strip()
            if text.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, get_text(uid, 'canceled'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
                del user_states[uid]; return
            if re.match(r'^https?://', text):
                thinking = bot.send_message(chat_id, get_text(uid, 'recognizing_music'), parse_mode="HTML")
                uuid_res = recognize_music_by_url(text)
                if uuid_res == "INVALID_URL":
                    bot.edit_message_text(get_text(uid, 'invalid_url'), chat_id, thinking.message_id, parse_mode="HTML")
                    del user_states[uid]; return
                if not uuid_res:
                    bot.edit_message_text(get_text(uid, 'recognition_error'), chat_id, thinking.message_id, parse_mode="HTML")
                    del user_states[uid]; return
                results = poll_recognition(uuid_res, chat_id, thinking.message_id, bot)
                if results:
                    answer = format_recognition_results(results)
                    bot.edit_message_text(get_text(uid, 'recognition_result', result=answer), chat_id, thinking.message_id, parse_mode="HTML")
                else:
                    bot.edit_message_text(get_text(uid, 'recognition_failed'), chat_id, thinking.message_id, parse_mode="HTML")
                del user_states[uid]
                return
            else:
                bot.send_message(chat_id, get_text(uid, 'invalid_url'), parse_mode="HTML")
                return

        if state == "awaiting_search_query":
            query = message.text.strip()
            if query.lower() in ["отмена","cancel"]:
                bot.send_message(chat_id, get_text(uid, 'canceled'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
                del user_states[uid]; return
            if not query:
                bot.send_message(chat_id, get_text(uid, 'empty_prompt'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
                del user_states[uid]; return
            del user_states[uid]
            thinking = bot.send_message(chat_id, get_text(uid, 'searching'), parse_mode="HTML")
            answer = search_internet(query, uid)
            if not answer:
                bot.edit_message_text(get_text(uid, 'search_failed'), chat_id, thinking.message_id, parse_mode="HTML")
                return
            bot.delete_message(chat_id, thinking.message_id)
            if len(answer) > MAX_MESSAGE_LENGTH:
                handle_long_response(chat_id, answer, uid, reply_to_message_id=message.message_id)
            else:
                send_long_message(chat_id, answer, reply_to_message_id=message.message_id)
            return

        if state == "awaiting_tts_text":
            del user_states[uid]
            process_tts_step(message)
            return

        if state.startswith("awaiting_generate_prompt:"):
            size = state.split(":",1)[1]
            del user_states[uid]
            prompt = message.text.strip()
            if prompt.lower() in ["отмена","cancel"]:
                bot.send_message(chat_id, get_text(uid, 'canceled'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
                return
            if not prompt:
                bot.send_message(chat_id, get_text(uid, 'empty_prompt'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
                return
            generate_and_send(message, prompt, size)
            return

        if state.startswith("awaiting_edit_prompt:"):
            size = state.split(":",1)[1]
            prompt = message.text.strip()
            if prompt.lower() in ["отмена","cancel"]:
                bot.send_message(chat_id, get_text(uid, 'canceled'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
                del user_states[uid]; return
            if not prompt:
                bot.send_message(chat_id, get_text(uid, 'empty_prompt'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
                del user_states[uid]; return
            cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="back_to_main"))
            user_states[uid] = f"awaiting_edit_photo:{size}:{prompt}"
            bot.send_message(chat_id, "Теперь отправьте фото (можно несколько в одном альбоме), которое нужно отредактировать.", reply_markup=cancel_markup)
            return

        if state.startswith("awaiting_edit_photo:"):
            pass

        if state.startswith("awaiting_video_text_prompt:"):
            duration = state.split(":",1)[1].lstrip(':')
            del user_states[uid]
            prompt = message.text.strip()
            if prompt.lower() in ["отмена","cancel"]:
                bot.send_message(chat_id, get_text(uid, 'canceled'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
                return
            if not prompt:
                bot.send_message(chat_id, get_text(uid, 'empty_description'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
                return

            base_price_per_sec = get_media_base_price('video_per_second')
            price = int(int(duration) * base_price_per_sec * get_profit_multiplier() * 100)
            if get_balance(uid) < price:
                bot.send_message(chat_id, get_text(uid, 'insufficient_funds', balance=get_balance(uid)/100, price=price/100), parse_mode="HTML")
                return

            thinking = bot.send_message(chat_id, get_text(uid, 'generating_video'), parse_mode="HTML")
            result = generate_video(prompt, int(duration), image_bytes=None)
            if not result:
                bot.edit_message_text(get_text(uid, 'video_failed'), chat_id, thinking.message_id, parse_mode="HTML")
                return
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix='.mp4') as tmp:
                    tmp.write(result); tmp_path = tmp.name
                with open(tmp_path,'rb') as f:
                    bot.send_video(chat_id, f, caption=f"🎥 {escape_html(prompt)}", timeout=600)
                os.unlink(tmp_path)
                update_balance(uid, -price, reason=f"Видео: {prompt[:50]}")
                bot.delete_message(chat_id, thinking.message_id)
            except Exception as e:
                bot.edit_message_text(f"❌ Ошибка отправки видео: {e}", chat_id, thinking.message_id)
            return

        if state == "awaiting_tiktok_url":
            url = message.text.strip()
            if url.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, get_text(uid, 'canceled'), reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
                del user_states[uid]
                return
            if not is_valid_tiktok_url(url):
                bot.send_message(chat_id, "❌ Это не похоже на ссылку TikTok. Попробуйте ещё раз или отправьте /cancel.", parse_mode="HTML")
                return
            waiting_msg = bot.send_message(chat_id, "⏳ Скачиваю видео из TikTok...", parse_mode="HTML")
            try:
                info = get_tiktok_video_info(url)
                if not info or not info.get('video_url'):
                    bot.edit_message_text("❌ Не удалось получить видео. Возможно, ссылка недействительна или видео приватное.", chat_id, waiting_msg.message_id, parse_mode="HTML")
                    del user_states[uid]
                    return

                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.mp4')
                temp_path = temp_file.name
                temp_file.close()

                success = download_tiktok_file(info['video_url'], temp_path)
                if not success:
                    bot.edit_message_text("❌ Ошибка при скачивании видео.", chat_id, waiting_msg.message_id, parse_mode="HTML")
                    os.unlink(temp_path)
                    del user_states[uid]
                    return

                caption = f"📱 TikTok\n👤 Автор: {info.get('author', 'Неизвестно')}"
                if info.get('description'):
                    desc = info['description'][:100] + "..." if len(info['description']) > 100 else info['description']
                    caption += f"\n📝 {desc}"
                if info.get('music'):
                    caption += f"\n🎵 {info['music']}"

                with open(temp_path, 'rb') as f:
                    bot.send_video(chat_id, f, caption=caption, timeout=120, parse_mode="HTML")

                os.unlink(temp_path)
                bot.delete_message(chat_id, waiting_msg.message_id)
                bot.send_message(chat_id, "✅ Готово!", reply_markup=main_inline_keyboard(uid), parse_mode="HTML")
            except Exception as e:
                log_error(f"TikTok download error: {e}")
                bot.edit_message_text(f"❌ Ошибка: {str(e)}", chat_id, waiting_msg.message_id, parse_mode="HTML")
            finally:
                del user_states[uid]
            return

        if state == "admin_waiting_balance_input":
            text = message.text.strip()
            if text.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, "❌ Отменено.", reply_markup=admin_panel_keyboard(uid))
                del user_states[uid]
                return
            try:
                parts = text.replace(',', '.').replace(':', '.').split()
                if len(parts) != 2:
                    raise ValueError("Нужно два значения: ID и сумма")
                target = int(parts[0])
                rub = float(parts[1])
                kop = int(round(rub * 100))
                update_balance(target, kop, reason="admin_edit")
                bot.send_message(chat_id, f"✅ Баланс пользователя {target} изменён на {rub:.2f} ₽. Теперь: {get_balance(target)/100:.2f} ₽", reply_markup=admin_panel_keyboard(uid))
            except Exception as e:
                bot.send_message(chat_id, f"❌ Ошибка: {e}. Используйте формат: ID сумма (например, 12345 100 или 12345 -50.5)", parse_mode="HTML")
                return
            del user_states[uid]
            return

        if state == "admin_waiting_broadcast":
            text = message.text.strip()
            if text.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, "❌ Отменено.", reply_markup=admin_panel_keyboard(uid))
                del user_states[uid]
                return
            if not text:
                bot.send_message(chat_id, "❌ Введите текст для рассылки.", reply_markup=admin_panel_keyboard(uid))
                return
            users = cursor.execute("SELECT user_id FROM users").fetchall()
            sent = 0
            failed = 0
            status_msg = bot.send_message(chat_id, "⏳ Рассылка...")
            for (target,) in users:
                try:
                    bot.send_message(target, f"📢 <b>Сообщение от администратора:</b>\n\n{escape_html(text)}", parse_mode="HTML")
                    sent += 1
                except Exception as e:
                    log_error(f"Broadcast error to {target}: {e}")
                    failed += 1
                time.sleep(0.05)
            bot.edit_message_text(f"✅ Отправлено: {sent}, не удалось: {failed}", chat_id, status_msg.message_id, reply_markup=admin_panel_keyboard(uid))
            del user_states[uid]
            return

        if state == "admin_waiting_broadcast_as_bot":
            text = message.text.strip()
            if text.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, "❌ Отменено.", reply_markup=admin_panel_keyboard(uid))
                del user_states[uid]
                return
            if not text:
                bot.send_message(chat_id, "❌ Введите текст для рассылки.", reply_markup=admin_panel_keyboard(uid))
                return
            users = cursor.execute("SELECT user_id FROM users").fetchall()
            sent = 0
            failed = 0
            status_msg = bot.send_message(chat_id, "⏳ Рассылка...")
            for (target,) in users:
                try:
                    bot.send_message(target, escape_html(text), parse_mode="HTML")
                    sent += 1
                except Exception as e:
                    log_error(f"Broadcast as bot error to {target}: {e}")
                    failed += 1
                time.sleep(0.05)
            bot.edit_message_text(f"✅ Отправлено: {sent}, не удалось: {failed}", chat_id, status_msg.message_id, reply_markup=admin_panel_keyboard(uid))
            del user_states[uid]
            return

        if state == "admin_waiting_profit_multiplier":
            value_str = message.text.strip()
            if value_str.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, "❌ Отменено.", reply_markup=admin_panel_keyboard(uid))
                del user_states[uid]
                return
            try:
                new_value = float(value_str.replace(',', '.'))
                if new_value <= 0:
                    raise ValueError
                set_profit_multiplier(new_value)
                bot.send_message(chat_id, f"✅ Коэффициент наценки изменен на <b>{new_value}</b>.", parse_mode="HTML", reply_markup=admin_panel_keyboard(uid))
            except ValueError:
                bot.send_message(chat_id, "❌ Ошибка: введите положительное число (например, 2.0).")
                return
            del user_states[uid]
            return

        if state.startswith("admin_waiting_aitunnel_price:"):
            parts = state.split(":", 2)
            if len(parts) != 3:
                del user_states[uid]
                return
            model_id, price_type = parts[1], parts[2]
            value_str = message.text.strip().replace(',', '.')
            if value_str.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, "❌ Отменено.", reply_markup=admin_panel_keyboard(uid))
                del user_states[uid]
                return
            try:
                new_value = float(value_str)
                if new_value < 0:
                    raise ValueError
                if model_id not in AITUNNEL_PRICES:
                    AITUNNEL_PRICES[model_id] = {}
                AITUNNEL_PRICES[model_id][price_type] = new_value
                save_aitunnel_price(model_id, price_type, new_value)
                price_name = "входную" if price_type == "input" else "выходную"
                bot.send_message(chat_id, f"✅ {price_name.capitalize()} цена для модели <b>{model_id}</b> изменена на <b>{new_value}</b> коп./токен.", parse_mode="HTML", reply_markup=admin_panel_keyboard(uid))
            except ValueError:
                bot.send_message(chat_id, "❌ Ошибка: введите неотрицательное число (например, 0.00192).")
                return
            del user_states[uid]
            return

        if state.startswith("admin_waiting_media_price:"):
            service = state.split(":", 1)[1]
            value_str = message.text.strip().replace(',', '.')
            if value_str.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, "❌ Отменено.", reply_markup=admin_panel_keyboard(uid))
                del user_states[uid]
                return
            try:
                new_value = float(value_str)
                if new_value <= 0:
                    raise ValueError

                if service == 'voice':
                    price_kop = int(new_value * 100)
                    with db_lock:
                        cursor.execute("INSERT OR REPLACE INTO prices (service, price) VALUES (?, ?)", ('voice', price_kop))
                        conn.commit()
                    PRICES['voice'] = price_kop
                    bot.send_message(chat_id, f"✅ Цена за распознавание голоса изменена на <b>{new_value:.2f}</b> ₽.", parse_mode="HTML", reply_markup=admin_panel_keyboard(uid))
                else:
                    MEDIA_BASE_PRICES[service] = new_value
                    set_media_base_price(service, new_value)
                    service_names = {
                        'generate': 'генерации изображения',
                        'video_per_second': 'видео (за секунду)',
                        'tts_per_1000_chars': 'TTS (за 1000 символов)'
                    }
                    bot.send_message(chat_id, f"✅ Базовая цена для {service_names.get(service, service)} изменена на <b>{new_value}</b> ₽.", parse_mode="HTML", reply_markup=admin_panel_keyboard(uid))
            except ValueError:
                bot.send_message(chat_id, "❌ Ошибка: введите положительное число.")
                return
            del user_states[uid]
            return

        if state == "admin_waiting_aitunnel_new_model_id":
            model_id = message.text.strip()
            if model_id.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, "❌ Отменено.", reply_markup=admin_panel_keyboard(uid))
                del user_states[uid]
                return
            if model_id in AITUNNEL_PRICES:
                bot.send_message(chat_id, "❌ Модель с таким ID уже существует. Введите другой ID или отправьте /cancel.", parse_mode="HTML")
                return
            user_states[uid] = f"admin_waiting_aitunnel_new_model_input:{model_id}"
            bot.send_message(chat_id, f"Введите <b>цену input</b> для модели {model_id} в копейках за токен (например, 0.024):", parse_mode="HTML")
            return

        if state.startswith("admin_waiting_aitunnel_new_model_input:"):
            parts = state.split(":", 1)
            if len(parts) != 2:
                del user_states[uid]
                return
            model_id = parts[1]
            value_str = message.text.strip().replace(',', '.')
            if value_str.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, "❌ Отменено.", reply_markup=admin_panel_keyboard(uid))
                del user_states[uid]
                return
            try:
                input_price = float(value_str)
                if input_price < 0:
                    raise ValueError
            except ValueError:
                bot.send_message(chat_id, "❌ Ошибка: введите неотрицательное число (например, 0.024).")
                return
            user_states[uid] = f"admin_waiting_aitunnel_new_model_output:{model_id}:{input_price}"
            bot.send_message(chat_id, f"Введите <b>цену output</b> для модели {model_id} в копейках за токен (например, 0.192):", parse_mode="HTML")
            return

        if state.startswith("admin_waiting_aitunnel_new_model_output:"):
            parts = state.split(":", 2)
            if len(parts) != 3:
                del user_states[uid]
                return
            model_id, input_price_str = parts[1], parts[2]
            try:
                input_price = float(input_price_str)
            except:
                del user_states[uid]
                return
            value_str = message.text.strip().replace(',', '.')
            if value_str.lower() in ["отмена", "cancel"]:
                bot.send_message(chat_id, "❌ Отменено.", reply_markup=admin_panel_keyboard(uid))
                del user_states[uid]
                return
            try:
                output_price = float(value_str)
                if output_price < 0:
                    raise ValueError
            except ValueError:
                bot.send_message(chat_id, "❌ Ошибка: введите неотрицательное число (например, 0.192).")
                return

            if model_id not in AITUNNEL_PRICES:
                AITUNNEL_PRICES[model_id] = {}
            AITUNNEL_PRICES[model_id]['input'] = input_price
            AITUNNEL_PRICES[model_id]['output'] = output_price
            save_aitunnel_price(model_id, 'input', input_price)
            save_aitunnel_price(model_id, 'output', output_price)

            bot.send_message(chat_id, f"✅ Модель <b>{model_id}</b> успешно добавлена!\nInput: {input_price} коп./токен\nOutput: {output_price} коп./токен", parse_mode="HTML", reply_markup=admin_panel_keyboard(uid))
            text = "🤖 <b>Цены моделей AITUNNEL</b>\n\nВыберите модель для редактирования:"
            markup = aitunnel_price_edit_keyboard(uid)
            bot.send_message(chat_id, text, reply_markup=markup, parse_mode="HTML")
            del user_states[uid]
            return

    if is_blocked(uid):
        bot.send_message(chat_id, get_text(uid, 'blocked'), parse_mode="HTML")
        return

    text = message.text
    if text.startswith('/') or text in [get_text(uid, 'generate'), get_text(uid, 'edit'), get_text(uid, 'video'), get_text(uid, 'search'), get_text(uid, 'recognize'), get_text(uid, 'tts'), get_text(uid, 'new_chat'), get_text(uid, 'info'), get_text(uid, 'contacts'), get_text(uid, 'help'), get_text(uid, 'model'), get_text(uid, 'voice', status=''), get_text(uid, 'topup'), get_text(uid, 'settings'), get_text(uid, 'admin_panel')]:
        return

    thinking = bot.send_message(chat_id, get_text(uid, 'analyzing'), parse_mode="HTML")
    answer = ask_aitunnel(uid, text)
    log_query(uid, text, answer)
    if not answer:
        bot.edit_message_text(get_text(uid, 'empty_response'), chat_id, thinking.message_id, parse_mode="HTML")
        return
    bot.delete_message(chat_id, thinking.message_id)
    if len(answer) > MAX_MESSAGE_LENGTH:
        handle_long_response(chat_id, answer, uid, reply_to_message_id=message.message_id)
    else:
        send_long_message(chat_id, answer, reply_to_message_id=message.message_id)

# ========== ОБРАБОТЧИКИ МЕДИА ==========

@bot.message_handler(content_types=['photo'], func=is_private)
@subscription_required
def handle_private_photo(message):
    uid = message.from_user.id
    chat_id = message.chat.id

    if get_in_chat(uid):
        forward_to_admin(message, uid, 'photo', message.photo[-1].file_id)
        bot.reply_to(message, get_text(uid, 'message_sent'), parse_mode="HTML")
        return

    if uid in ADMIN_IDS and uid in user_states and isinstance(user_states[uid], str) and user_states[uid].startswith("replying_to_"):
        target = int(user_states[uid].split("_")[2])
        try:
            bot.send_photo(target, message.photo[-1].file_id, caption="✉️ Ответ от создателя")
            bot.reply_to(message, get_text(uid, 'replied', target=target), parse_mode="HTML")
        except Exception as e:
            bot.reply_to(message, get_text(uid, 'reply_error', error=str(e)), parse_mode="HTML")
        return

    if uid in user_states and isinstance(user_states[uid], str) and user_states[uid].startswith("awaiting_video_photo:"):
        duration = user_states[uid].split(":",1)[1].lstrip(':')
        del user_states[uid]
        prompt = message.caption or "Оживи это фото"
        image_bytes = download_file(message.photo[-1].file_id)
        if not image_bytes:
            bot.send_message(chat_id, get_text(uid, 'failed_load_photo'), parse_mode="HTML")
            return

        base_price_per_sec = get_media_base_price('video_per_second')
        price = int(int(duration) * base_price_per_sec * get_profit_multiplier() * 100)
        if get_balance(uid) < price:
            bot.send_message(chat_id, get_text(uid, 'insufficient_funds', balance=get_balance(uid)/100, price=price/100), parse_mode="HTML")
            return

        thinking = bot.send_message(chat_id, get_text(uid, 'generating_video'), parse_mode="HTML")
        result = generate_video(prompt, int(duration), image_bytes=image_bytes)
        if not result:
            bot.edit_message_text(get_text(uid, 'video_failed'), chat_id, thinking.message_id, parse_mode="HTML")
            return
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.mp4') as tmp:
                tmp.write(result); tmp_path = tmp.name
            with open(tmp_path,'rb') as f:
                bot.send_video(chat_id, f, caption=f"🎥 {escape_html(prompt)}", timeout=600)
            os.unlink(tmp_path)
            update_balance(uid, -price, reason=f"Видео с фото: {prompt[:50]}")
            bot.delete_message(chat_id, thinking.message_id)
        except Exception as e:
            bot.edit_message_text(f"❌ Ошибка отправки видео: {e}", chat_id, thinking.message_id)
        return

    if uid in user_states and isinstance(user_states[uid], str) and user_states[uid].startswith("awaiting_edit_photo:"):
        parts = user_states[uid].split(":",2)
        if len(parts) == 3:
            size, prompt = parts[1], parts[2]
            if message.media_group_id:
                media_group_id = message.media_group_id
                with media_group_lock:
                    if media_group_id not in pending_media_groups:
                        pending_media_groups[media_group_id] = {
                            'images': [],
                            'uid': uid,
                            'chat_id': chat_id,
                            'state_info': (size, prompt)
                        }
                        img_bytes = download_file(message.photo[-1].file_id)
                        if img_bytes:
                            pending_media_groups[media_group_id]['images'].append(img_bytes)
                        timer = threading.Timer(2.0, process_media_group, args=(media_group_id, uid, chat_id, (size, prompt)))
                        timer.daemon = True
                        timer.start()
                    else:
                        img_bytes = download_file(message.photo[-1].file_id)
                        if img_bytes:
                            pending_media_groups[media_group_id]['images'].append(img_bytes)
                return
            else:
                del user_states[uid]
                image_bytes = download_file(message.photo[-1].file_id)
                if not image_bytes:
                    bot.send_message(chat_id, get_text(uid, 'failed_load_photo'), parse_mode="HTML")
                    return
                process_edit_step_multiple(chat_id, uid, prompt, size, [image_bytes])
                return

    if is_blocked(uid):
        bot.send_message(chat_id, get_text(uid, 'blocked'), parse_mode="HTML")
        return

    image_bytes = download_file(message.photo[-1].file_id)
    if not image_bytes:
        bot.send_message(chat_id, get_text(uid, 'failed_load_photo'), parse_mode="HTML")
        return
    caption = message.caption or "Опиши это изображение."
    thinking = bot.send_message(chat_id, get_text(uid, 'analyzing'), parse_mode="HTML")
    answer = ask_aitunnel(uid, caption, image_bytes)
    log_query(uid, f"[Фото] {caption}", answer)
    if not answer:
        bot.edit_message_text(get_text(uid, 'empty_response'), chat_id, thinking.message_id, parse_mode="HTML")
        return
    bot.delete_message(chat_id, thinking.message_id)
    if len(answer) > MAX_MESSAGE_LENGTH:
        handle_long_response(chat_id, answer, uid, reply_to_message_id=message.message_id)
    else:
        send_long_message(chat_id, answer, reply_to_message_id=message.message_id)

@bot.message_handler(content_types=['video'], func=is_private)
@subscription_required
def handle_private_video(message):
    uid = message.from_user.id
    chat_id = message.chat.id

    if get_in_chat(uid):
        forward_to_admin(message, uid, 'video', message.video.file_id)
        bot.reply_to(message, get_text(uid, 'message_sent'), parse_mode="HTML")
        return

    if is_blocked(uid):
        bot.send_message(chat_id, get_text(uid, 'blocked'), parse_mode="HTML")
        return

    video_bytes = download_file(message.video.file_id)
    if not video_bytes:
        bot.send_message(chat_id, get_text(uid, 'failed_load_photo'), parse_mode="HTML")
        return

    caption = message.caption or "Опиши это видео."
    thinking = bot.send_message(chat_id, get_text(uid, 'analyzing'), parse_mode="HTML")
    answer = ask_aitunnel(uid, caption, video_bytes=video_bytes)
    log_query(uid, f"[Видео] {caption}", answer)
    if not answer:
        bot.edit_message_text(get_text(uid, 'empty_response'), chat_id, thinking.message_id, parse_mode="HTML")
        return
    bot.delete_message(chat_id, thinking.message_id)
    if len(answer) > MAX_MESSAGE_LENGTH:
        handle_long_response(chat_id, answer, uid, reply_to_message_id=message.message_id)
    else:
        send_long_message(chat_id, answer, reply_to_message_id=message.message_id)

@bot.message_handler(content_types=['document'], func=is_private)
@subscription_required
def handle_private_document(message):
    uid = message.from_user.id
    chat_id = message.chat.id

    if get_in_chat(uid):
        forward_to_admin(message, uid, 'document', message.document.file_id)
        bot.reply_to(message, get_text(uid, 'message_sent'), parse_mode="HTML")
        return

    if uid in ADMIN_IDS and uid in user_states and isinstance(user_states[uid], str) and user_states[uid].startswith("replying_to_"):
        target = int(user_states[uid].split("_")[2])
        try:
            bot.send_document(target, message.document.file_id, caption="✉️ Ответ от создателя")
            bot.reply_to(message, get_text(uid, 'replied', target=target), parse_mode="HTML")
        except Exception as e:
            bot.reply_to(message, get_text(uid, 'reply_error', error=str(e)), parse_mode="HTML")
        return

    if is_blocked(uid):
        bot.send_message(chat_id, get_text(uid, 'blocked'), parse_mode="HTML")
        return

    doc = message.document
    ext = os.path.splitext(doc.file_name or "")[1].lower()
    file_bytes = download_file(doc.file_id)
    if not file_bytes:
        bot.send_message(chat_id, get_text(uid, 'failed_load_photo'), parse_mode="HTML")
        return
    extractors = {
        '.pdf': extract_text_from_pdf,
        '.docx': extract_text_from_docx,
        '.txt': extract_text_from_txt,
        '.csv': extract_text_from_txt,
        '.json': extract_text_from_txt,
        '.xml': extract_text_from_txt,
        '.html': extract_text_from_txt,
        '.css': extract_text_from_txt,
        '.js': extract_text_from_txt,
        '.py': extract_text_from_txt,
        '.log': extract_text_from_txt,
        '.ini': extract_text_from_txt,
        '.conf': extract_text_from_txt,
        '.yaml': extract_text_from_txt,
        '.yml': extract_text_from_txt,
        '.md': extract_text_from_txt,
        '.xlsx': extract_text_from_xlsx,
        '.xls': extract_text_from_xls,
        '.pptx': extract_text_from_pptx,
        '.rtf': extract_text_from_rtf,
    }
    if ext in extractors:
        extracted = extractors[ext](file_bytes)
    else:
        bot.reply_to(message, "❌ Неподдерживаемый формат.")
        return

    if extracted.startswith('[Ошибка') or extracted.startswith('[Библиотека'):
        bot.reply_to(message, f"❌ {extracted}")
        return

    caption = message.caption or "Содержимое файла:"
    full_prompt = f"{caption}\n\n{extracted}"
    thinking = bot.send_message(chat_id, get_text(uid, 'reading_document'), parse_mode="HTML")
    answer = ask_aitunnel(uid, full_prompt)
    log_query(uid, f"[Документ] {doc.file_name}", answer)
    if not answer:
        bot.edit_message_text(get_text(uid, 'empty_response'), chat_id, thinking.message_id, parse_mode="HTML")
        return
    bot.delete_message(chat_id, thinking.message_id)
    if len(answer) > MAX_MESSAGE_LENGTH:
        handle_long_response(chat_id, answer, uid, reply_to_message_id=message.message_id)
    else:
        send_long_message(chat_id, answer, reply_to_message_id=message.message_id)

@bot.message_handler(content_types=['voice'])
@subscription_required
def handle_voice(message):
    if not is_private(message): return
    uid = message.from_user.id
    chat_id = message.chat.id

    if get_in_chat(uid):
        forward_to_admin(message, uid, 'voice', message.voice.file_id)
        bot.reply_to(message, get_text(uid, 'message_sent'), parse_mode="HTML")
        return

    if uid in user_states and user_states[uid] == "awaiting_music_recognition":
        voice_bytes = download_voice(message.voice.file_id)
        if not voice_bytes:
            bot.send_message(chat_id, get_text(uid, 'failed_load_voice'), parse_mode="HTML")
            return
        thinking = bot.send_message(chat_id, get_text(uid, 'recognizing_music'), parse_mode="HTML")
        uuid_res = recognize_music_by_file(voice_bytes, '.ogg')
        if uuid_res == "INVALID_URL":
            bot.edit_message_text(get_text(uid, 'invalid_url'), chat_id, thinking.message_id, parse_mode="HTML")
            del user_states[uid]; return
        if not uuid_res:
            bot.edit_message_text(get_text(uid, 'recognition_error'), chat_id, thinking.message_id, parse_mode="HTML")
            del user_states[uid]; return
        results = poll_recognition(uuid_res, chat_id, thinking.message_id, bot)
        if results:
            answer = format_recognition_results(results)
            bot.edit_message_text(get_text(uid, 'recognition_result', result=answer), chat_id, thinking.message_id, parse_mode="HTML")
        else:
            bot.edit_message_text(get_text(uid, 'recognition_failed'), chat_id, thinking.message_id, parse_mode="HTML")
        del user_states[uid]
        return

    if not get_voice_enabled(uid):
        bot.send_message(chat_id, get_text(uid, 'voice', status='❌'), parse_mode="HTML")
        return

    price = PRICES.get('voice',0)
    if get_balance(uid) < price:
        bot.send_message(chat_id, get_text(uid, 'insufficient_funds', balance=get_balance(uid)/100, price=price/100), parse_mode="HTML")
        return

    update_balance(uid, -price, reason="Голосовой ввод")
    thinking = bot.send_message(chat_id, get_text(uid, 'recognizing_voice'), parse_mode="HTML")
    voice_bytes = download_voice(message.voice.file_id)
    if not voice_bytes:
        bot.edit_message_text(get_text(uid, 'failed_load_voice'), chat_id, thinking.message_id, parse_mode="HTML")
        return
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.ogg') as tmp:
            tmp.write(voice_bytes); tmp_path = tmp.name
        with open(tmp_path, 'rb') as audio:
            transcript = aitunnel_client.audio.transcriptions.create(model=TRANSCRIBE_MODEL, file=audio, language="ru")
        os.unlink(tmp_path)
        transcribed = transcript.text
    except:
        bot.edit_message_text(get_text(uid, 'transcribe_failed'), chat_id, thinking.message_id, parse_mode="HTML")
        return

    bot.edit_message_text(
        get_text(uid, 'recognized_text', text=escape_html(transcribed)) + "\n\n" + get_text(uid, 'thinking_again'),
        chat_id, thinking.message_id
    )
    answer = ask_aitunnel(uid, transcribed)
    log_query(uid, f"[Voice] {transcribed}", answer)
    if not answer:
        bot.edit_message_text(get_text(uid, 'empty_response'), chat_id, thinking.message_id, parse_mode="HTML")
        return
    bot.delete_message(chat_id, thinking.message_id)
    if len(answer) > MAX_MESSAGE_LENGTH:
        handle_long_response(chat_id, answer, uid, reply_to_message_id=message.message_id)
    else:
        send_long_message(chat_id, answer, reply_to_message_id=message.message_id)

@bot.message_handler(content_types=['audio','video','sticker'])
@subscription_required
def handle_other_media(message):
    if not is_private(message): return
    uid = message.from_user.id
    chat_id = message.chat.id

    if get_in_chat(uid):
        forward_to_admin(message, uid, 'audio' if message.audio else 'video' if message.video else 'sticker',
                         message.audio.file_id if message.audio else message.video.file_id if message.video else message.sticker.file_id)
        bot.reply_to(message, get_text(uid, 'message_sent'), parse_mode="HTML")
        return

    if uid in user_states and user_states[uid] == "awaiting_music_recognition":
        if message.audio:
            file_bytes = download_file(message.audio.file_id)
            if not file_bytes:
                bot.send_message(chat_id, get_text(uid, 'failed_load_voice'), parse_mode="HTML")
                return
            ext = os.path.splitext(message.audio.file_name or ".mp3")[1]
            thinking = bot.send_message(chat_id, get_text(uid, 'recognizing_music'), parse_mode="HTML")
            uuid_res = recognize_music_by_file(file_bytes, ext)
            if uuid_res == "INVALID_URL":
                bot.edit_message_text(get_text(uid, 'invalid_url'), chat_id, thinking.message_id, parse_mode="HTML")
                del user_states[uid]; return
            if not uuid_res:
                bot.edit_message_text(get_text(uid, 'recognition_error'), chat_id, thinking.message_id, parse_mode="HTML")
                del user_states[uid]; return
            results = poll_recognition(uuid_res, chat_id, thinking.message_id, bot)
            if results:
                answer = format_recognition_results(results)
                bot.edit_message_text(get_text(uid, 'recognition_result', result=answer), chat_id, thinking.message_id, parse_mode="HTML")
            else:
                bot.edit_message_text(get_text(uid, 'recognition_failed'), chat_id, thinking.message_id, parse_mode="HTML")
            del user_states[uid]
            return
        else:
            bot.send_message(chat_id, get_text(uid, 'invalid_url'), parse_mode="HTML")
            return

    bot.reply_to(message, "Я принимаю текст, фото, документы и голосовые.")

def forward_to_admin(message, user_id, content_type, content):
    caption = f"📨 От @{escape_html(message.from_user.username or message.from_user.first_name)} (ID: {user_id})\n"
    markup = types.InlineKeyboardMarkup(row_width=2).add(
        types.InlineKeyboardButton("✏️ Ответить", callback_data=f"reply_to_{user_id}"),
        types.InlineKeyboardButton("🔚 Завершить чат", callback_data=f"end_chat_{user_id}")
    )
    for admin_id in ADMIN_IDS:
        try:
            if content_type == 'text':
                bot.send_message(admin_id, caption + content, reply_markup=markup)
            elif content_type == 'photo':
                bot.send_photo(admin_id, content, caption=caption, reply_markup=markup)
            elif content_type == 'document':
                bot.send_document(admin_id, content, caption=caption, reply_markup=markup)
            elif content_type == 'voice':
                bot.send_voice(admin_id, content, caption=caption, reply_markup=markup)
            elif content_type == 'audio':
                bot.send_audio(admin_id, content, caption=caption, reply_markup=markup)
            elif content_type == 'video':
                bot.send_video(admin_id, content, caption=caption, reply_markup=markup)
            elif content_type == 'sticker':
                bot.send_sticker(admin_id, content, reply_markup=markup)
        except Exception as e:
            log_error(f"Forward to admin {admin_id} error: {e}")

# ========== НОВЫЙ ОБРАБОТЧИК WEB APP DATA ==========

@bot.message_handler(content_types=['web_app_data'])
def handle_web_app_data(message):
    uid = message.from_user.id
    try:
        data = json.loads(message.web_app_data.data)
        log_error(f"Web App data from {uid}: {data}")
        if 'model' in data:
            set_text_model(uid, data['model'])
        if 'search' in data:
            set_internet_search_enabled(uid, data['search'])
        if 'voice' in data:
            set_tts_voice(uid, data['voice'])
        if 'mode' in data:
            set_mode(uid, data['mode'])
        bot.answer_web_app_query(message.web_app_data.query_id, "Настройки сохранены")
    except Exception as e:
        log_error(f"Web App data error: {e}")
        bot.answer_web_app_query(message.web_app_data.query_id, "Ошибка")

# ========== АДМИНСКИЕ ШАГИ (register_next_step_handler) ==========

def admin_edit_balance_step(message):
    uid = message.from_user.id
    if not is_admin(uid): return
    text = message.text.strip()
    if text.lower() == "отмена":
        bot.send_message(message.chat.id, "Отменено.", reply_markup=main_inline_keyboard(uid))
        return
    try:
        parts = text.split()
        target = int(parts[0])
        rub = float(parts[1])
        kop = int(round(rub * 100))
        update_balance(target, kop, reason="admin_edit")
        bot.send_message(message.chat.id, f"✅ Баланс {target} теперь {get_balance(target)/100:.2f} ₽", reply_markup=main_inline_keyboard(uid))
    except:
        bot.send_message(message.chat.id, "❌ Ошибка формата. Используйте: ID сумма (например 12345 100 или 12345 -50)", reply_markup=main_inline_keyboard(uid))

def admin_broadcast_step(message):
    uid = message.from_user.id
    if not is_admin(uid): return
    text = message.text.strip()
    if text.lower() == "отмена":
        bot.send_message(message.chat.id, "Отменено.", reply_markup=main_inline_keyboard(uid))
        return
    users = cursor.execute("SELECT user_id FROM users").fetchall()
    sent = failed = 0
    status = bot.send_message(message.chat.id, "⏳ Рассылка...")
    for (target,) in users:
        try:
            bot.send_message(target, f"📢 <b>Сообщение от администратора:</b>\n\n{escape_html(text)}", parse_mode="HTML")
            sent += 1
        except:
            failed += 1
    bot.edit_message_text(f"✅ Отправлено: {sent}, не удалось: {failed}", message.chat.id, status.message_id, reply_markup=main_inline_keyboard(uid))

def admin_broadcast_as_bot_step(message):
    uid = message.from_user.id
    if not is_admin(uid): return
    text = message.text.strip()
    if text.lower() == "отмена":
        bot.send_message(message.chat.id, "Отменено.", reply_markup=main_inline_keyboard(uid))
        return
    users = cursor.execute("SELECT user_id FROM users").fetchall()
    sent = failed = 0
    status = bot.send_message(message.chat.id, "⏳ Рассылка (как бот)...")
    for (target,) in users:
        try:
            bot.send_message(target, escape_html(text), parse_mode="HTML")
            sent += 1
        except:
            failed += 1
    bot.edit_message_text(f"✅ Отправлено: {sent}, не удалось: {failed}", message.chat.id, status.message_id, reply_markup=main_inline_keyboard(uid))

def admin_block_step(message):
    uid = message.from_user.id
    if not is_admin(uid): return
    text = message.text.strip()
    if text.lower() == "отмена":
        bot.send_message(message.chat.id, "Отменено.", reply_markup=main_inline_keyboard(uid))
        return
    try:
        target = int(text)
        block_user(target)
        bot.send_message(message.chat.id, f"✅ {target} заблокирован.", reply_markup=main_inline_keyboard(uid))
    except:
        bot.send_message(message.chat.id, "❌ Введите ID.", reply_markup=main_inline_keyboard(uid))

def admin_unblock_step(message):
    uid = message.from_user.id
    if not is_admin(uid): return
    text = message.text.strip()
    if text.lower() == "отмена":
        bot.send_message(message.chat.id, "Отменено.", reply_markup=main_inline_keyboard(uid))
        return
    try:
        target = int(text)
        unblock_user(target)
        bot.send_message(message.chat.id, f"✅ {target} разблокирован.", reply_markup=main_inline_keyboard(uid))
    except:
        bot.send_message(message.chat.id, "❌ Введите ID.", reply_markup=main_inline_keyboard(uid))

def admin_end_chat_step(message):
    uid = message.from_user.id
    if not is_admin(uid): return
    try:
        target = int(message.text.strip())
    except:
        bot.send_message(message.chat.id, "❌ Введите корректный ID.", reply_markup=main_inline_keyboard(uid))
        return
    if uid in user_states and isinstance(user_states[uid], str) and user_states[uid].startswith("replying_to_") and int(user_states[uid].split("_")[2]) == target:
        del user_states[uid]
    set_in_chat(target, 0)
    markup = rating_keyboard(target)
    try:
        bot.send_message(target, get_text(target, 'chat_ended'), reply_markup=markup, parse_mode="HTML")
        bot.send_message(message.chat.id, f"✅ Чат с {target} завершён, запрос оценки отправлен.", reply_markup=main_inline_keyboard(uid))
    except Exception as e:
        bot.send_message(message.chat.id, f"❌ Не удалось отправить сообщение пользователю {target}: {e}", reply_markup=main_inline_keyboard(uid))

# ========== ОТЗЫВЫ ==========

REVIEWS_FILE = "reviews.json"
def save_review(user_id, username, first_name, rating_stars, comment, admin_response):
    user_display = f"@{username}" if username else first_name or str(user_id)
    review = {"rating": rating_stars, "user": user_display, "comment": comment or "", "response": admin_response or ""}
    if os.path.exists(REVIEWS_FILE):
        with open(REVIEWS_FILE, "r", encoding="utf-8") as f:
            try:
                reviews = json.load(f)
            except:
                reviews = []
    else:
        reviews = []
    reviews.append(review)
    with open(REVIEWS_FILE, "w", encoding="utf-8") as f:
        json.dump(reviews, f, ensure_ascii=False, indent=2)

# ========== CALLBACK QUERY ==========

@bot.callback_query_handler(func=lambda call: True)
def callback_query(call):
    log_error(f"!!! callback_query ВЫЗВАНА от пользователя {call.from_user.id}")

    uid = call.from_user.id

    if call.message.chat.type == 'private' and not is_admin(uid):
        if not is_subscribed(uid):
            text = (f'<a href="https://t.me/ken4kk_news">KeN4kk_News</a>. '
                    f'Чтобы пользоваться ботом, необходимо подписаться на канал. '
                    f'Мы сделали так, чтобы не допускать ботов и чтобы вы могли получать '
                    f'сообщения об обновлениях и новости через канал/бота.')
            markup = types.InlineKeyboardMarkup().add(
                types.InlineKeyboardButton("📢 Подписаться", url="https://t.me/ken4kk_news")
            )
            try:
                bot.send_message(call.message.chat.id, text, parse_mode="HTML", reply_markup=markup)
                log_error(f"Отправлено сообщение о подписке пользователю {uid}")
            except Exception as e:
                log_error(f"Ошибка при отправке сообщения о подписке пользователю {uid}: {e}")

            try:
                bot.answer_callback_query(call.id)
            except Exception as e:
                log_error(f"Ошибка при ответе на callback (подписка): {e}")

            return

    register_or_update_user(call.message)
    cid = call.message.chat.id
    mid = call.message.message_id

    # ===== ОБРАБОТКА ДЛИННЫХ ОТВЕТОВ =====

    if call.data.startswith("long_text:"):
        response_id = call.data.split(":", 1)[1]
        try:
            resp = requests.get(f"{LONG_RESPONSE_API}/api/get/{response_id}", timeout=10)
            if resp.status_code == 200:
                data = resp.json()
                text = data['text']
                send_long_message(call.message.chat.id, text, reply_to_message_id=call.message.message_id)
            else:
                bot.answer_callback_query(call.id, "❌ Ответ не найден на сервере.")
        except Exception as e:
            log_error(f"Ошибка получения текста с сайта: {e}")
            bot.answer_callback_query(call.id, "❌ Ошибка при загрузке ответа.")
        bot.answer_callback_query(call.id)
        return

    elif call.data.startswith("long_file:"):
        response_id = call.data.split(":", 1)[1]
        try:
            resp = requests.get(f"{LONG_RESPONSE_API}/api/get/{response_id}", timeout=10)
            if resp.status_code == 200:
                data = resp.json()
                text = data['text']
                with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False, encoding='utf-8') as f:
                    f.write(text)
                    temp_path = f.name
                try:
                    with open(temp_path, 'rb') as f:
                        bot.send_document(
                            call.message.chat.id,
                            f,
                            caption="📄 Ваш длинный ответ",
                            reply_to_message_id=call.message.message_id
                        )
                finally:
                    os.unlink(temp_path)
            else:
                bot.answer_callback_query(call.id, "❌ Ответ не найден на сервере.")
        except Exception as e:
            log_error(f"Ошибка получения текста с сайта для файла: {e}")
            bot.answer_callback_query(call.id, "❌ Ошибка при загрузке ответа.")
        bot.answer_callback_query(call.id)
        return

    # ===== ОСТАЛЬНЫЕ ОБРАБОТЧИКИ =====

    if call.data == "generate":
        bot.answer_callback_query(call.id)
        text = get_text(uid, 'choose_size')
        markup = image_size_keyboard(uid, "generate")
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
        except:
            bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")

    elif call.data == "edit":
        bot.answer_callback_query(call.id)
        text = get_text(uid, 'choose_size')
        markup = image_size_keyboard(uid, "edit")
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
        except:
            bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")

    elif call.data == "video":
        bot.answer_callback_query(call.id)
        text = get_text(uid, 'choose_video_duration')
        markup = video_duration_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
        except:
            bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")

    elif call.data == "search":
        bot.answer_callback_query(call.id)
        text = get_text(uid, 'enter_search_query')
        markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Отмена", callback_data="back_to_main"))
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
        except:
            bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")
        user_states[uid] = "awaiting_search_query"

    elif call.data == "recognize_music":
        bot.answer_callback_query(call.id)
        user_states[uid] = "awaiting_music_recognition"
        markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="back_to_main"))
        text = get_text(uid, 'send_audio_or_link')
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
        except:
            bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")

    elif call.data == "tiktok_downloader":
        bot.answer_callback_query(call.id)
        text = get_text(uid, 'send_tiktok_link')
        markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Отмена", callback_data="back_to_main"))
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
        except:
            bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")
        user_states[uid] = "awaiting_tiktok_url"

    elif call.data == "tts_synthesize":
        bot.answer_callback_query(call.id)
        text = get_text(uid, 'enter_tts_text')
        markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Отмена", callback_data="back_to_main"))
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
        except:
            bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")
        user_states[uid] = "awaiting_tts_text"

    elif call.data.startswith("size_"):
        parts = call.data.split("_")
        if len(parts) != 3: return
        action, size = parts[1], parts[2]
        if size not in IMAGE_SIZES: return
        price = get_media_base_price('generate') * get_profit_multiplier()
        user_states[uid] = f"awaiting_{action}_prompt:{size}"
        text = get_text(uid, 'enter_prompt', size=size, price=price)
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="back_to_main"))
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=cancel_markup, parse_mode="HTML")
        except:
            bot.send_message(cid, text, reply_markup=cancel_markup, parse_mode="HTML")
        bot.answer_callback_query(call.id)

    elif call.data.startswith("vdur_"):
        parts = call.data.split("_")
        if len(parts) != 2: return
        duration = parts[1]
        if int(duration) not in VIDEO_DURATIONS:
            bot.answer_callback_query(call.id, "❌ Неверная длительность")
            return
        text = get_text(uid, 'choose_mode')
        markup = video_mode_keyboard(uid, duration)
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
        except:
            bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")
        bot.answer_callback_query(call.id)

    elif call.data.startswith("vmode_text_"):
        parts = call.data.split("_")
        if len(parts) != 3: return
        duration = parts[2]
        user_states[uid] = f"awaiting_video_text_prompt:{duration}"
        text = get_text(uid, 'enter_video_text')
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="back_to_main"))
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=cancel_markup, parse_mode="HTML")
        except:
            bot.send_message(cid, text, reply_markup=cancel_markup, parse_mode="HTML")
        bot.answer_callback_query(call.id)

    elif call.data.startswith("vmode_photo_"):
        parts = call.data.split("_")
        if len(parts) != 3: return
        duration = parts[2]
        user_states[uid] = f"awaiting_video_photo:{duration}"
        text = get_text(uid, 'send_photo_for_video')
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="back_to_main"))
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=cancel_markup, parse_mode="HTML")
        except:
            bot.send_message(cid, text, reply_markup=cancel_markup, parse_mode="HTML")
        bot.answer_callback_query(call.id)

    elif call.data == "new_chat":
        clear_history(uid)
        bot.answer_callback_query(call.id, get_text(uid, 'history_cleared_toast'))
        text = get_text(uid, 'history_cleared')
        markup = main_inline_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
        except:
            bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")

    elif call.data == "info":
        model_id = get_text_model(uid)
        display = next((d for d,m in AVAILABLE_MODELS.items() if m==model_id), model_id)
        voice_status = "✅ Включён" if get_voice_enabled(uid) else "❌ Отключён"
        bal = get_balance(uid)/100
        mode = get_mode(uid)
        mode_name = get_text(uid, 'mode_normal') if mode=='normal' else get_text(uid, 'mode_adult')
        free_left = get_free_quota_left(uid)
        free_tts_left = get_free_tts_quota_left(uid)
        text = (f'<tg-emoji emoji-id="{PREMIUM_EMOJI["info"]}">ℹ️</tg-emoji> <b>{BOT_NAME}</b>\n📌 Версия: {BOT_VERSION}\n🧠 Модель: {escape_html(display)} (цена зависит от длины запроса)\n'
                f"🎤 Голос: {voice_status}\n💰 Баланс: {bal:.2f} ₽\n"
                f"🔄 Режим: {mode_name}\n✨ Бесплатных текстовых запросов: {free_left}\n"
                f"🎧 Бесплатных TTS запросов: {free_tts_left}\n👤 Создатель: {CREATOR_TAG}")
        markup = main_inline_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)

    elif call.data == "contacts":
        text = "📞 Связь с создателем"
        markup = contacts_menu_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)

    elif call.data == "news":
        bot.answer_callback_query(call.id)
        text = get_text(uid, 'news_channel')
        bot.send_message(cid, text,
                         parse_mode="HTML",
                         reply_markup=types.InlineKeyboardMarkup().add(
                             types.InlineKeyboardButton("◀️ Назад", callback_data="back_to_main")
                         ))

    elif call.data == "toggle_voice":
        new = 0 if get_voice_enabled(uid) else 1
        set_voice_enabled(uid, new)
        bot.answer_callback_query(call.id, f"Голос {'включён' if new else 'выключен'}")
        text = get_text(uid, 'main_menu')
        markup = main_inline_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
        except:
            bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")

    elif call.data == "toggle_internet_search":
        new = 0 if get_internet_search_enabled(uid) else 1
        set_internet_search_enabled(uid, new)
        status = "включён" if new else "выключен"
        bot.answer_callback_query(call.id, f"Режим поиска {status}")
        text = get_text(uid, 'main_menu')
        markup = main_inline_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
        except:
            bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")

    elif call.data == "choose_model":
        text = get_text(uid, 'choose_model')
        markup = model_choice_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")
        bot.answer_callback_query(call.id)

    elif call.data.startswith("set_model_"):
        model_id = call.data.replace("set_model_","")
        if model_id in AVAILABLE_MODELS.values():
            set_text_model(uid, model_id)
            name = next((n for n,m in AVAILABLE_MODELS.items() if m==model_id), model_id)
            bot.answer_callback_query(call.id, f"Модель: {name}")
            text = get_text(uid, 'main_menu')
            markup = main_inline_keyboard(uid)
            try:
                bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
            except:
                bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")
        else:
            bot.answer_callback_query(call.id, "Неизвестная модель")

    elif call.data == "help":
        text = get_text(uid, 'ticket_ask_problem')
        markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Назад", callback_data="back_to_main"))
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        user_states[uid] = "awaiting_ticket_problem"
        bot.answer_callback_query(call.id)

    elif call.data == "profile":
        show_profile(uid, cid, call.from_user)
        bot.answer_callback_query(call.id)

    elif call.data == "back_to_main":
        if uid in user_states: del user_states[uid]
        text = get_text(uid, 'main_menu')
        markup = main_inline_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
        except:
            bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")
        bot.answer_callback_query(call.id)

    elif call.data == "topup_menu":
        text = get_text(uid, 'topup_menu')
        markup = topup_menu_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")
        bot.answer_callback_query(call.id)

    elif call.data == "stars_topup":
        text = get_text(uid, 'stars_enter_amount')
        markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Назад", callback_data="topup_menu"))
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        user_states[uid] = {"state": "awaiting_stars_amount"}
        bot.answer_callback_query(call.id)

    elif call.data == "card_payment":
        text = get_text(uid, 'card_enter_amount')
        markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Назад", callback_data="topup_menu"))
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        user_states[uid] = {"state": "awaiting_card_amount"}
        bot.answer_callback_query(call.id)

    elif call.data == "crypto_topup":
        text = get_text(uid, 'crypto_enter_amount_rub')
        markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Назад", callback_data="topup_menu"))
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        user_states[uid] = {"state": "awaiting_crypto_amount"}
        bot.answer_callback_query(call.id)

    elif call.data.startswith("crypto_check:"):
        invoice_id = call.data.split(":",1)[1]
        row = cursor.execute("SELECT user_id, amount_rub, status FROM crypto_invoices WHERE invoice_id=?", (invoice_id,)).fetchone()
        if not row:
            bot.answer_callback_query(call.id, "Счёт не найден")
            return
        user_id, rub, status = row
        if user_id != uid:
            bot.answer_callback_query(call.id, "Это не ваш счёт")
            return
        if status == 'paid':
            bot.answer_callback_query(call.id, "✅ Счёт уже оплачен")
            return
        api_status = get_crypto_invoice_status(invoice_id)
        if api_status == 'paid':
            update_balance(uid, rub*100, f"Crypto {rub} руб.")
            cursor.execute("UPDATE crypto_invoices SET status='paid' WHERE invoice_id=?", (invoice_id,))
            conn.commit()
            text = get_text(uid, 'crypto_payment_success', amount=rub)
            markup = main_inline_keyboard(uid)
            try:
                bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
            except:
                bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")
            bot.answer_callback_query(call.id)
        elif api_status == 'active':
            bot.answer_callback_query(call.id, get_text(uid, 'crypto_payment_pending'), show_alert=True)
        else:
            text = get_text(uid, 'crypto_payment_failed')
            markup = main_inline_keyboard(uid)
            try:
                bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
            except:
                bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")
            bot.answer_callback_query(call.id)

    elif call.data == "tts_menu":
        text = get_text(uid, 'tts_menu')
        markup = tts_menu_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")
        bot.answer_callback_query(call.id)

    elif call.data == "choose_tts_voice":
        text = get_text(uid, 'tts_voices')
        markup = tts_voice_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)

    elif call.data.startswith("set_tts_voice_"):
        voice = call.data.replace("set_tts_voice_","")
        if voice in AVAILABLE_TTS_VOICES:
            set_tts_voice(uid, voice)
            bot.answer_callback_query(call.id, f"Голос: {voice}")
            text = get_text(uid, 'tts_menu')
            markup = tts_menu_keyboard(uid)
            try:
                bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
            except:
                bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        else:
            bot.answer_callback_query(call.id, "Неизвестный голос")

    elif call.data == "contact_via_bot":
        set_in_chat(uid, 1)
        bot.send_message(cid, get_text(uid, 'message_sent'), parse_mode="HTML")
        bot.answer_callback_query(call.id)

    elif call.data == "settings_menu":
        text = get_text(uid, 'settings_menu')
        markup = settings_menu_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)

    elif call.data == "settings_language":
        text = get_text(uid, 'choose_language_prompt')
        markup = language_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)

    elif call.data.startswith("set_lang_"):
        lang = call.data.replace("set_lang_","")
        if lang in LANGUAGES:
            set_language(uid, lang)
            bot.answer_callback_query(call.id, get_text(uid, 'language_updated', lang=get_language_name(lang)))
            text = get_text(uid, 'main_menu')
            markup = main_inline_keyboard(uid)
            try:
                bot.edit_message_text(text, cid, mid, reply_markup=markup, parse_mode="HTML")
            except:
                bot.send_message(cid, text, reply_markup=markup, parse_mode="HTML")
        else:
            bot.answer_callback_query(call.id, "Неверный язык")

    elif call.data == "admin_panel":
        if not is_admin(uid): return
        text = "🔧 <b>Админ-панель</b>"
        markup = admin_panel_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)

    elif call.data == "docs":
        multiplier = get_profit_multiplier()
        generate_price = get_media_base_price('generate') * multiplier
        edit_price = generate_price
        video_price = get_media_base_price('video_per_second') * multiplier
        tts_price = get_media_base_price('tts_per_1000_chars') * multiplier
        voice_price = PRICES.get('voice', 500) / 100
        search_price = PRICES.get('search', 100) / 100
        text = (f"{get_text(uid, 'docs_title')}\n\n"
                f"{get_text(uid, 'docs_generate')} <b>{generate_price:.2f} ₽</b>\n"
                f"{get_text(uid, 'docs_edit')} <b>{edit_price:.2f} ₽</b>\n"
                f"{get_text(uid, 'docs_video')} <b>{video_price:.2f} ₽</b>\n"
                f"{get_text(uid, 'docs_tts')} <b>{tts_price:.2f} ₽</b>\n"
                f"{get_text(uid, 'docs_voice')} <b>{voice_price:.2f} ₽</b>\n"
                f"{get_text(uid, 'docs_search')} <b>{search_price:.2f} ₽</b>\n"
                f"{get_text(uid, 'docs_recognize')} <b>{get_text(uid, 'docs_free')}</b>\n"
                f"{get_text(uid, 'docs_tiktok')} <b>{get_text(uid, 'docs_free')}</b>")
        bot.send_message(cid, text, parse_mode="HTML", reply_markup=main_inline_keyboard(uid))
        bot.answer_callback_query(call.id)

    elif call.data == "user_agreement":
        lang = get_language(uid)
        if lang == 'ru':
            text = "📜 <b>Юридические документы</b>\n\nВыберите документ для просмотра:"
            markup = types.InlineKeyboardMarkup(row_width=1)
            markup.add(
                types.InlineKeyboardButton("📄 Публичная оферта", url="https://telegra.ph/Publichnaya-oferta-Polzovatelskoe-soglashenie-03-07"),
                types.InlineKeyboardButton("🔒 Политика конфиденциальности", url="https://telegra.ph/Politika-Konfidencialnosti-03-07-28"),
                types.InlineKeyboardButton("◀️ Назад", callback_data="back_to_main")
            )
        else:
            text = "📜 <b>Legal Documents</b>\n\nChoose a document to view:"
            markup = types.InlineKeyboardMarkup(row_width=1)
            markup.add(
                types.InlineKeyboardButton("📄 Public Offer", url="https://telegra.ph/Publichnaya-oferta-Polzovatelskoe-soglashenie-03-07"),
                types.InlineKeyboardButton("🔒 Privacy Policy", url="https://telegra.ph/Politika-Konfidencialnosti-03-07-28"),
                types.InlineKeyboardButton("◀️ Back", callback_data="back_to_main")
            )
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)

    elif call.data == "referrals":
        total_invited = get_referral_count(uid, 'pending') + get_referral_count(uid, 'approved') + get_referral_count(uid, 'rejected')
        pending = get_referral_count(uid, 'pending')
        available_for_request = pending
        remaining = max(0, 5 - pending)

        referral_link = generate_referral_link(uid)

        text = (f'<tg-emoji emoji-id="5258165702707125574">✨</tg-emoji> <b>Подарки за друзей</b>\n\n'
                f'<tg-emoji emoji-id="5258362837411045098">🎁</tg-emoji> Пригласи 5 друзей и получи подарок бесплатно.\n'
                f'<i>После жми «Отправить заявку», админы проверят накрутку и вышлют подарок.</i>\n\n'
                f'<tg-emoji emoji-id="5258486128742244085">👥</tg-emoji> Всего приглашённых: <b>{total_invited}</b>\n'
                f'<tg-emoji emoji-id="5260416304224936047">📨</tg-emoji> Доступно для заявки: <b>{available_for_request}</b>\n\n'
                f'<tg-emoji emoji-id="5258461531464539536">🔗</tg-emoji> Твоя ссылка:\n<code>{referral_link}</code>\n\n'
                f'Осталось до подарка: <b>{remaining}</b>')

        markup = referrals_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)

    elif call.data == "send_referral_request":
        pending = get_referral_count(uid, 'pending')
        if pending == 0:
            bot.answer_callback_query(call.id, "❌ У вас нет новых приглашённых для заявки.", show_alert=True)
            return

        referrals = get_referrals_list(uid, 'pending')
        if not referrals:
            bot.answer_callback_query(call.id, "❌ Ошибка при получении списка.", show_alert=True)
            return

        admin_msg = f"📨 <b>Заявка на подарок от пользователя</b>\n\n"
        admin_msg += f"👤 Пользователь: @{call.from_user.username or 'нет'} (ID: {uid})\n"
        admin_msg += f"📊 Количество приглашённых: {len(referrals)}\n\n"
        admin_msg += "<b>Список приглашённых:</b>\n"

        for ref_id, username, first_name, joined in referrals:
            admin_msg += f"• ID: <code>{ref_id}</code> @{username or 'нет'} {first_name or ''} (с {joined[:10]})\n"

        admin_msg += f"\nВыберите действие:"

        markup = types.InlineKeyboardMarkup(row_width=2)
        markup.add(
            types.InlineKeyboardButton("✅ Одобрить", callback_data=f"approve_gift_{uid}"),
            types.InlineKeyboardButton("❌ Отклонить", callback_data=f"reject_gift_{uid}")
        )

        sent_to_admins = False
        for admin_id in ADMIN_IDS:
            try:
                bot.send_message(admin_id, admin_msg, parse_mode="HTML", reply_markup=markup)
                sent_to_admins = True
            except:
                pass

        if sent_to_admins:
            bot.answer_callback_query(call.id, "✅ Заявка отправлена администраторам.")
        else:
            bot.answer_callback_query(call.id, "❌ Не удалось отправить заявку (нет доступных админов).", show_alert=True)

    elif call.data.startswith("approve_gift_"):
        if not is_admin(uid):
            bot.answer_callback_query(call.id)
            return
        target_uid = int(call.data.replace("approve_gift_", ""))

        # Проверяем, не отправляли ли уже подарок
        with db_lock:
            already_sent = cursor.execute("SELECT COUNT(*) FROM referrals WHERE referrer_id=? AND gift_sent=1", (target_uid,)).fetchone()[0]
            if already_sent > 0:
                bot.send_message(cid, "⚠️ Подарок этому пользователю уже был отправлен ранее.")
                bot.answer_callback_query(call.id)
                return

            gift_id = int(get_setting('gift_id', 0))
            if not gift_id:
                bot.send_message(cid, "❌ ID подарка не настроен. Используйте /setgiftid <число>")
                bot.answer_callback_query(call.id)
                return

        # Отправляем подарок (вне блокировки, чтобы не держать БД)
        bot.send_message(cid, f"⏳ Отправляю подарок пользователю {target_uid}...")
        try:
            success = send_gift_direct(target_uid, gift_id)
            if success:
                with db_lock:
                    cursor.execute("UPDATE referrals SET status='approved', gift_sent=1 WHERE referrer_id=? AND status='pending'", (target_uid,))
                    conn.commit()
                bot.send_message(cid, f"✅ Подарок успешно отправлен пользователю {target_uid}.")
                try:
                    bot.send_message(target_uid, "🎁 Поздравляем! Вы получили подарок за приглашение друзей!")
                except:
                    pass
            else:
                bot.send_message(cid, "❌ Не удалось отправить подарок. Проверьте баланс администратора и ID подарка.")
        except Exception as e:
            log_error(f"Ошибка при отправке подарка: {e}")
            bot.send_message(cid, f"❌ Ошибка при отправке подарка: {e}")

        bot.answer_callback_query(call.id)

    elif call.data.startswith("reject_gift_"):
        if not is_admin(uid):
            bot.answer_callback_query(call.id)
            return
        target_uid = int(call.data.replace("reject_gift_", ""))
        with db_lock:
            cursor.execute("UPDATE referrals SET status='rejected' WHERE referrer_id=? AND status='pending'", (target_uid,))
            conn.commit()
        bot.send_message(cid, f"❌ Заявка на подарок от пользователя {target_uid} отклонена.")
        try:
            bot.send_message(target_uid, "❌ Ваша заявка на подарок отклонена.")
        except:
            pass
        bot.answer_callback_query(call.id)

    # ---- Тикеты ----

    elif call.data.startswith("reply_ticket:"):
        parts = call.data.split(":")
        if len(parts) == 3:
            target_user_id = int(parts[1])
            ticket_msg_id = int(parts[2])
            text = "🧑‍💻 Введите ваш ответ для пользователя (можно с фото):"
            markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Назад", callback_data=f"back_to_ticket:{ticket_msg_id}"))
            try:
                bot.edit_message_text(text, cid, mid, reply_markup=markup)
            except:
                bot.send_message(cid, text, reply_markup=markup)
            user_states[uid] = {"state": "admin_reply_to_ticket", "target_user_id": target_user_id, "ticket_msg_id": ticket_msg_id}
        bot.answer_callback_query(call.id)
    elif call.data.startswith("close_ticket:"):
        ticket_msg_id = int(call.data.split(":",1)[1])
        try:
            bot.delete_message(cid, ticket_msg_id)
        except: pass
        if ticket_msg_id in tickets:
            del tickets[ticket_msg_id]
        bot.answer_callback_query(call.id, "Тикет скрыт")
    elif call.data == "user_reply":
        text = get_text(uid, 'ticket_user_reply_prompt')
        markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="cancel_user_reply"))
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup)
        except:
            bot.send_message(cid, text, reply_markup=markup)
        user_states[uid] = {"state": "user_reply_to_admin"}
        bot.answer_callback_query(call.id)
    elif call.data == "cancel_user_reply":
        del user_states[uid]
        bot.delete_message(cid, mid)
    elif call.data == "helpful":
        text = get_text(uid, 'ticket_reply_helpful')
        try:
            bot.edit_message_text(text, cid, mid)
        except:
            bot.send_message(cid, text)
        time.sleep(1)
        bot.delete_message(cid, mid)
        start_private(call.message)
        for admin_id in ADMIN_IDS:
            try: bot.send_message(admin_id, f"Пользователю @{call.from_user.username} помог ответ! 👍")
            except: pass
        bot.answer_callback_query(call.id)
    elif call.data.startswith("back_to_ticket:"):
        ticket_msg_id = int(call.data.split(":",1)[1])
        if ticket_msg_id in tickets:
            data = tickets[ticket_msg_id]
            user_id = data['user_id']
            text = data['text']
            photo_id = data.get('photo_id')
            username = bot.get_chat(user_id).username or "None"
            if photo_id:
                msg = bot.send_photo(cid, photo_id, caption=f"<b>📨 Запрос в поддержку!</b>\n\n👤 Пользователь: @{username}\n🆔 ID: {user_id}\n📝 Проблема: {text}", parse_mode="HTML")
            else:
                msg = bot.send_message(cid, f"<b>📨 Запрос в поддержку!</b>\n\n👤 Пользователь: @{username}\n🆔 ID: {user_id}\n📝 Проблема: {text}", parse_mode="HTML")
            tickets[msg.message_id] = data
            markup = types.InlineKeyboardMarkup().add(
                types.InlineKeyboardButton("💬 Ответить", callback_data=f"reply_ticket:{user_id}:{msg.message_id}"),
                types.InlineKeyboardButton("❌ Скрыть", callback_data=f"close_ticket:{msg.message_id}")
            )
            bot.edit_message_reply_markup(cid, msg.message_id, reply_markup=markup)
            bot.delete_message(cid, mid)
        else:
            bot.answer_callback_query(call.id, "Тикет не найден")

    # ---- Заявки на карту ----

    elif call.data == "admin_card_requests":
        if not is_admin(uid): return
        payments = get_pending_card_requests()
        if not payments:
            bot.send_message(cid, "Нет pending заявок.")
            bot.answer_callback_query(call.id)
            return
        text = "📋 Выберите заявку:"
        markup = card_requests_keyboard(uid, payments)
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup)
        except:
            bot.send_message(cid, text, reply_markup=markup)
        bot.answer_callback_query(call.id)
    elif call.data.startswith("card_request_"):
        if not is_admin(uid): return
        pid = int(call.data.replace("card_request_",""))
        text = f"Заявка #{pid}"
        markup = card_request_action_keyboard(pid)
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup)
        except:
            bot.send_message(cid, text, reply_markup=markup)
        bot.answer_callback_query(call.id)

    elif call.data.startswith("card_approve_"):
        if not is_admin(uid):
            bot.answer_callback_query(call.id)
            return
        try:
            log_error(f"card_approve_ callback: entered try block")
            pid = int(call.data.replace("card_approve_",""))
            log_error(f"card_approve_ callback: pid={pid}")
            log_error(f"Admin {uid} approving card payment {pid} - start")
            success, result = approve_card_payment(pid, uid)
            log_error(f"Admin {uid} - approve_card_payment returned {success}, {result}")
            if success:
                user_id = result
                amount_row = cursor.execute("SELECT amount_rub FROM card_payments WHERE id=?", (pid,)).fetchone()
                amount = amount_row[0] if amount_row else "?"
                bot.answer_callback_query(call.id, "✅ Заявка подтверждена")
                bot.send_message(cid, get_text(uid, 'card_approved', id=pid, amount=amount))
                try:
                    bot.send_message(user_id, "✅ Ваша заявка на пополнение картой подтверждена. Баланс начислен.")
                except Exception as e:
                    log_error(f"Error notifying user {user_id}: {e}")
            else:
                bot.answer_callback_query(call.id, "❌ Ошибка", show_alert=True)
                bot.send_message(cid, result)
            log_error(f"Admin {uid} - after sending messages, now deleting message {mid}")
            try:
                bot.delete_message(cid, mid)
            except Exception as e:
                log_error(f"Error deleting message {mid}: {e}")
            log_error(f"Admin {uid} - after delete, now showing pending requests")
            payments = get_pending_card_requests()
            if payments:
                bot.send_message(cid, "📋 Выберите заявку:", reply_markup=card_requests_keyboard(uid, payments))
            else:
                bot.send_message(cid, "📋 Нет pending заявок.", reply_markup=admin_panel_keyboard(uid))
            log_error(f"Admin {uid} - finished processing card_approve_")
        except Exception as e:
            log_error(f"Error in card_approve_ callback: {traceback.format_exc()}")
            bot.answer_callback_query(call.id, "Произошла ошибка", show_alert=True)

    elif call.data.startswith("card_reject_"):
        if not is_admin(uid):
            bot.answer_callback_query(call.id)
            return
        try:
            pid = int(call.data.replace("card_reject_",""))
            reject_card_payment(pid)
            bot.answer_callback_query(call.id, "❌ Заявка отклонена")
            bot.send_message(cid, get_text(uid, 'card_rejected', id=pid))
            user_row = cursor.execute("SELECT user_id FROM card_payments WHERE id=?", (pid,)).fetchone()
            if user_row:
                try:
                    bot.send_message(user_row[0], "❌ Ваша заявка на пополнение картой отклонена.")
                except Exception as e:
                    log_error(f"Error notifying user {user_row[0]}: {e}")
            try:
                bot.delete_message(cid, mid)
            except Exception as e:
                log_error(f"Error deleting message {mid}: {e}")
            payments = get_pending_card_requests()
            if payments:
                bot.send_message(cid, "📋 Выберите заявку:", reply_markup=card_requests_keyboard(uid, payments))
            else:
                bot.send_message(cid, "📋 Нет pending заявок.", reply_markup=admin_panel_keyboard(uid))
        except Exception as e:
            log_error(f"Error in card_reject_ callback: {traceback.format_exc()}")
            bot.answer_callback_query(call.id, "Произошла ошибка", show_alert=True)

    # ========== ЗАЯВКИ НА ПОДАРКИ ==========

    elif call.data == "admin_gift_requests":
        if not is_admin(uid):
            bot.answer_callback_query(call.id)
            return
        with db_lock:
            rows = cursor.execute("SELECT DISTINCT referrer_id FROM referrals WHERE status='pending'").fetchall()
        if not rows:
            bot.send_message(cid, "Нет ожидающих заявок на подарки.")
            bot.answer_callback_query(call.id)
            return

        text = "📨 <b>Заявки на подарки</b>\n\n"
        markup = types.InlineKeyboardMarkup(row_width=1)
        for (ref_id,) in rows:
            count = get_referral_count(ref_id, 'pending')
            user_info = cursor.execute("SELECT username, first_name FROM users WHERE user_id=?", (ref_id,)).fetchone()
            uname = user_info[0] if user_info else ""
            fname = user_info[1] if user_info else ""
            display = f"@{uname}" if uname else (fname or str(ref_id))
            text += f"👤 {display} (ID: {ref_id}) — приглашено: {count}\n"
            markup.add(types.InlineKeyboardButton(
                f"👤 {display} ({count})",
                callback_data=f"show_gift_request_{ref_id}"
            ))
        markup.add(types.InlineKeyboardButton("◀️ Назад", callback_data="admin_panel"))
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)

    elif call.data.startswith("show_gift_request_"):
        if not is_admin(uid):
            bot.answer_callback_query(call.id)
            return
        ref_id = int(call.data.replace("show_gift_request_", ""))
        referrals = get_referrals_list(ref_id, 'pending')
        if not referrals:
            bot.send_message(cid, "У этого пользователя нет ожидающих рефералов.")
            bot.answer_callback_query(call.id)
            return

        text = f"📨 <b>Заявка от пользователя {ref_id}</b>\n\n"
        for refd_id, username, first_name, joined in referrals:
            text += f"• ID: <code>{refd_id}</code> @{username or 'нет'} {first_name or ''} (с {joined[:10]})\n"
        markup = types.InlineKeyboardMarkup(row_width=2)
        markup.add(
            types.InlineKeyboardButton("✅ Одобрить", callback_data=f"approve_gift_{ref_id}"),
            types.InlineKeyboardButton("❌ Отклонить", callback_data=f"reject_gift_{ref_id}"),
            types.InlineKeyboardButton("◀️ Назад", callback_data="admin_gift_requests")
        )
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)

    # ---- Расширенное управление ценами ----

    elif call.data == "admin_prices_management":
        if not is_admin(uid): return
        text = "⚙️ <b>Управление ценами</b>\n\nВыберите, что хотите изменить:"
        markup = admin_prices_management_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)

    elif call.data == "admin_edit_profit_multiplier":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_profit_multiplier"
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="admin_prices_management"))
        bot.send_message(cid, f"Текущий коэффициент наценки: <b>{PROFIT_MULTIPLIER}</b>\n\nВведите новое значение (например, 2.5):", parse_mode="HTML", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)

    elif call.data == "admin_edit_aitunnel_prices":
        if not is_admin(uid): return
        text = "🤖 <b>Цены моделей AITUNNEL</b>\n\nВыберите модель для редактирования:"
        markup = aitunnel_price_edit_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)

    elif call.data.startswith("aem_"):
        model_id = call.data.replace("aem_", "")
        prices = AITUNNEL_PRICES.get(model_id, {})
        text = f"🤖 <b>Модель: {model_id}</b>\n\n"
        text += f"🔹 Input price: <code>{prices.get('input', 'N/A')}</code> коп./токен\n"
        text += f"🔸 Output price: <code>{prices.get('output', 'N/A')}</code> коп./токен\n\n"
        text += "Выберите, что изменить:"
        markup = aitunnel_model_price_edit_keyboard(uid, model_id)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)

    elif call.data.startswith("asp_"):
        parts = call.data.replace("asp_", "").split("_")
        if len(parts) == 2:
            model_id, price_type = parts
            user_states[uid] = f"admin_waiting_aitunnel_price:{model_id}:{price_type}"
            price_name = "input" if price_type == "input" else "output"
            cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data=f"aem_{model_id}"))
            bot.send_message(cid, f"Введите новую цену <b>{price_name}</b> для модели <b>{model_id}</b> в копейках за токен (например, 0.00192):", parse_mode="HTML", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)

    elif call.data == "admin_add_aitunnel_model":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_aitunnel_new_model_id"
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="admin_edit_aitunnel_prices"))
        bot.send_message(cid, "Введите <b>идентификатор</b> новой модели (например, gpt-5):", parse_mode="HTML", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)

    elif call.data == "admin_edit_media_generate":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_media_price:generate"
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="admin_prices_management"))
        bot.send_message(cid, f"Текущая базовая цена генерации изображения: <b>{MEDIA_BASE_PRICES['generate']}</b> ₽\n\nВведите новую цену в рублях (например, 9.75):", parse_mode="HTML", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)

    elif call.data == "admin_edit_media_video":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_media_price:video_per_second"
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="admin_prices_management"))
        bot.send_message(cid, f"Текущая базовая цена видео (за секунду): <b>{MEDIA_BASE_PRICES['video_per_second']}</b> ₽\n\nВведите новую цену в рублях (например, 10.0):", parse_mode="HTML", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)

    elif call.data == "admin_edit_media_tts":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_media_price:tts_per_1000_chars"
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="admin_prices_management"))
        bot.send_message(cid, f"Текущая базовая цена TTS (за 1000 символов): <b>{MEDIA_BASE_PRICES['tts_per_1000_chars']}</b> ₽\n\nВведите новую цену в рублях (например, 10.0):", parse_mode="HTML", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)

    elif call.data == "admin_edit_media_voice":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_media_price:voice"
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="admin_prices_management"))
        bot.send_message(cid, f"Текущая цена за распознавание голоса: <b>{PRICES.get('voice', 500)/100:.2f}</b> ₽\n\nВведите новую цену в рублях (например, 5.0):", parse_mode="HTML", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)

    # ---- Новая рассылка как бот ----

    elif call.data == "admin_broadcast_as_bot":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_broadcast_as_bot"
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="admin_panel"))
        bot.send_message(cid, "Введите текст (или отправьте фото с подписью) для рассылки (как бот, без подписи администратора):", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)

    # ---- Управление ценами (старая, оставлена для совместимости) ----

    elif call.data == "admin_prices":
        if not is_admin(uid): return
        text = "💰 <b>Текущие цены (в рублях):</b>\n\n"
        for s, p in PRICES.items():
            text += f"• {s}: {kopecks_to_rubles(p)} ₽\n"
        markup = price_edit_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)
    elif call.data.startswith("admin_setprice_"):
        if not is_admin(uid): return
        service = call.data.replace("admin_setprice_","")
        user_states[uid] = f"admin_waiting_price:{service}"
        bot.send_message(cid, f"Введите новую цену для <b>{service}</b> в рублях (можно дробную, например 0.1):")
        bot.answer_callback_query(call.id)

    # ---- Управление курсами ----

    elif call.data == "admin_rates":
        if not is_admin(uid): return
        text = "Выберите курс для изменения:"
        markup = types.InlineKeyboardMarkup(row_width=2)
        markup.add(
            types.InlineKeyboardButton("💵 Курс USD", callback_data="admin_set_usd_rate"),
            types.InlineKeyboardButton("⭐ Курс звезды", callback_data="admin_set_star_rate"),
            types.InlineKeyboardButton("◀️ Назад", callback_data="admin_panel")
        )
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup)
        except:
            bot.send_message(cid, text, reply_markup=markup)
        bot.answer_callback_query(call.id)

    elif call.data == "admin_set_usd_rate":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_usd_rate"
        markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Отмена", callback_data="admin_rates"))
        bot.send_message(cid, "Введите новый курс USD к рублю (например 80.5):", reply_markup=markup)
        bot.answer_callback_query(call.id)

    elif call.data == "admin_set_star_rate":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_star_rate"
        markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Отмена", callback_data="admin_rates"))
        bot.send_message(cid, "Введите новый курс звезды к рублю (например 1.5):", reply_markup=markup)
        bot.answer_callback_query(call.id)

    # ---- TTS модели ----

    elif call.data == "admin_tts_models":
        if not is_admin(uid): return
        text = "🎤 <b>Управление TTS моделями</b>"
        markup = tts_models_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)
    elif call.data == "admin_tts_list":
        if not is_admin(uid): return
        text = "📋 <b>TTS модели:</b>\n\n"
        for mid, name in AVAILABLE_TTS_MODELS.items():
            text += f"• <b>{name}</b>\n  ID: <code>{mid}</code>\n\n"
        bot.send_message(cid, text, parse_mode="HTML")
        back_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Назад в админ-панель", callback_data="admin_panel"))
        bot.send_message(cid, "Выберите действие:", reply_markup=back_markup)
        bot.answer_callback_query(call.id)
    elif call.data == "admin_tts_add":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_tts_add_name"
        bot.send_message(cid, "➕ Введите <b>название</b> новой TTS модели:", parse_mode="HTML")
        bot.answer_callback_query(call.id)
    elif call.data == "admin_tts_delete":
        if not is_admin(uid): return
        if len(AVAILABLE_TTS_MODELS) <= 1:
            bot.answer_callback_query(call.id, "❌ Нельзя удалить последнюю TTS модель!")
            return
        markup = types.InlineKeyboardMarkup(row_width=1)
        for mid, name in AVAILABLE_TTS_MODELS.items():
            markup.add(types.InlineKeyboardButton(name, callback_data=f"admin_tts_del_{mid}"))
        markup.add(types.InlineKeyboardButton("◀️ Назад", callback_data="admin_tts_models"))
        text = "Выберите TTS модель для удаления:"
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup)
        except:
            bot.send_message(cid, text, reply_markup=markup)
        bot.answer_callback_query(call.id)
    elif call.data.startswith("admin_tts_del_"):
        if not is_admin(uid): return
        model_id = call.data.replace("admin_tts_del_","")
        try:
            remove_tts_model(model_id)
            bot.send_message(cid, f"✅ TTS модель {model_id} удалена.", parse_mode="HTML")
            back_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Назад в админ-панель", callback_data="admin_panel"))
            bot.send_message(cid, "Выберите действие:", reply_markup=back_markup)
        except Exception as e:
            bot.send_message(cid, f"❌ Ошибка: {e}")
        bot.answer_callback_query(call.id)

    # ---- Видео модели ----

    elif call.data == "admin_video_models":
        if not is_admin(uid): return
        text = "🎥 <b>Управление видео моделями</b>"
        markup = video_models_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)
    elif call.data == "admin_video_list":
        if not is_admin(uid): return
        text = "📋 <b>Видео модели:</b>\n\n"
        for mid, name in AVAILABLE_VIDEO_MODELS.items():
            text += f"• <b>{name}</b>\n  ID: <code>{mid}</code>\n\n"
        bot.send_message(cid, text, parse_mode="HTML")
        back_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Назад в админ-панель", callback_data="admin_panel"))
        bot.send_message(cid, "Выберите действие:", reply_markup=back_markup)
        bot.answer_callback_query(call.id)
    elif call.data == "admin_video_add":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_video_add_name"
        bot.send_message(cid, "➕ Введите <b>название</b> новой видео модели:", parse_mode="HTML")
        bot.answer_callback_query(call.id)
    elif call.data == "admin_video_delete":
        if not is_admin(uid): return
        if len(AVAILABLE_VIDEO_MODELS) <= 1:
            bot.answer_callback_query(call.id, "❌ Нельзя удалить последнюю видео модель!")
            return
        markup = types.InlineKeyboardMarkup(row_width=1)
        for mid, name in AVAILABLE_VIDEO_MODELS.items():
            markup.add(types.InlineKeyboardButton(name, callback_data=f"admin_video_del_{mid}"))
        markup.add(types.InlineKeyboardButton("◀️ Назад", callback_data="admin_video_models"))
        text = "Выберите видео модель для удаления:"
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup)
        except:
            bot.send_message(cid, text, reply_markup=markup)
        bot.answer_callback_query(call.id)
    elif call.data.startswith("admin_video_del_"):
        if not is_admin(uid): return
        model_id = call.data.replace("admin_video_del_","")
        try:
            remove_video_model(model_id)
            bot.send_message(cid, f"✅ Видео модель {model_id} удалена.", parse_mode="HTML")
            back_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Назад в админ-панель", callback_data="admin_panel"))
            bot.send_message(cid, "Выберите действие:", reply_markup=back_markup)
        except Exception as e:
            bot.send_message(cid, f"❌ Ошибка: {e}")
        bot.answer_callback_query(call.id)

    # ---- Модели изображений ----

    elif call.data == "admin_image_models":
        if not is_admin(uid): return
        text = "🖼 <b>Управление моделями изображений</b>"
        markup = image_models_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)
    elif call.data == "admin_image_list":
        if not is_admin(uid): return
        text = "📋 <b>Модели изображений:</b>\n\n"
        for mid, name in AVAILABLE_IMAGE_MODELS.items():
            text += f"• <b>{name}</b>\n  ID: <code>{mid}</code>\n\n"
        bot.send_message(cid, text, parse_mode="HTML")
        back_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Назад в админ-панель", callback_data="admin_panel"))
        bot.send_message(cid, "Выберите действие:", reply_markup=back_markup)
        bot.answer_callback_query(call.id)
    elif call.data == "admin_image_add":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_image_add_name"
        bot.send_message(cid, "➕ Введите <b>название</b> новой модели изображений:", parse_mode="HTML")
        bot.answer_callback_query(call.id)
    elif call.data == "admin_image_delete":
        if not is_admin(uid): return
        if len(AVAILABLE_IMAGE_MODELS) <= 1:
            bot.answer_callback_query(call.id, "❌ Нельзя удалить последнюю модель изображений!")
            return
        markup = types.InlineKeyboardMarkup(row_width=1)
        for mid, name in AVAILABLE_IMAGE_MODELS.items():
            markup.add(types.InlineKeyboardButton(name, callback_data=f"admin_image_del_{mid}"))
        markup.add(types.InlineKeyboardButton("◀️ Назад", callback_data="admin_image_models"))
        text = "Выберите модель изображений для удаления:"
        try:
            bot.edit_message_text(text, cid, mid, reply_markup=markup)
        except:
            bot.send_message(cid, text, reply_markup=markup)
        bot.answer_callback_query(call.id)
    elif call.data.startswith("admin_image_del_"):
        if not is_admin(uid): return
        model_id = call.data.replace("admin_image_del_","")
        try:
            remove_image_model(model_id)
            bot.send_message(cid, f"✅ Модель изображений {model_id} удалена.", parse_mode="HTML")
            back_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Назад в админ-панель", callback_data="admin_panel"))
            bot.send_message(cid, "Выберите действие:", reply_markup=back_markup)
        except Exception as e:
            bot.send_message(cid, f"❌ Ошибка: {e}")
        bot.answer_callback_query(call.id)

    # ---- Управление админами ----

    elif call.data == "admin_admins":
        if not is_admin(uid): return
        text = "👑 <b>Управление администраторами</b>"
        markup = admin_admins_keyboard(uid)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML", reply_markup=markup)
        except:
            bot.send_message(cid, text, parse_mode="HTML", reply_markup=markup)
        bot.answer_callback_query(call.id)
    elif call.data == "admin_add_admin":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_add_admin"
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="admin_admins"))
        bot.send_message(cid, "Введите ID пользователя, которого нужно сделать администратором:", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)
    elif call.data == "admin_remove_admin":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_remove_admin"
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="admin_admins"))
        bot.send_message(cid, "Введите ID пользователя, которого нужно лишить прав администратора:", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)
    elif call.data == "admin_list_admins":
        if not is_admin(uid): return
        admins = cursor.execute("SELECT user_id, username, first_name FROM users WHERE is_admin=1").fetchall()
        text = "👑 <b>Список администраторов:</b>\n\n"
        for aid, uname, fname in admins:
            text += f"• ID: <code>{aid}</code> @{escape_html(uname)} {escape_html(fname)}\n"
        bot.send_message(cid, text, parse_mode="HTML")
        back_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Назад в админ-панель", callback_data="admin_panel"))
        bot.send_message(cid, "Выберите действие:", reply_markup=back_markup)
        bot.answer_callback_query(call.id)

    # ---- История платежей ----

    elif call.data == "admin_transactions":
        if not is_admin(uid): return
        transactions = get_all_transactions(50)
        text = "📜 <b>Последние 50 транзакций:</b>\n\n"
        for user_id, amount, typ, desc, ts in transactions:
            sign = "+" if amount > 0 else ""
            text += f"{ts[:19]} | User {user_id} | {sign}{kopecks_to_rubles(amount)} ₽ | {typ} | {desc}\n"
        if not transactions:
            text += "Нет транзакций."
        send_long_message(cid, text)
        back_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Назад в админ-панель", callback_data="admin_panel"))
        bot.send_message(cid, "Выберите действие:", reply_markup=back_markup)
        bot.answer_callback_query(call.id)

    # ---- Логи пользователя ----

    elif call.data == "admin_user_logs":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_userid_for_logs"
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="admin_panel"))
        bot.send_message(cid, "Введите ID пользователя для просмотра его последних запросов:", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)

    # ---- Список пользователей ----

    elif call.data == "admin_list":
        if not is_admin(uid): return
        user_states[uid] = "admin_list_page:0"
        show_user_page(uid, cid, mid, 0)
        bot.answer_callback_query(call.id)
    elif call.data.startswith("admin_list_page_"):
        if not is_admin(uid): return
        try:
            page = int(call.data.replace("admin_list_page_",""))
        except:
            bot.answer_callback_query(call.id, "Ошибка")
            return
        user_states[uid] = f"admin_list_page:{page}"
        show_user_page(uid, cid, mid, page)
        bot.answer_callback_query(call.id)

    # ---- Изменение баланса ----

    elif call.data == "admin_balance":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_balance_input"
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="admin_panel"))
        bot.send_message(cid, "Введите ID и сумму в рублях (положительную для пополнения, отрицательную для снятия). Пример: `12345 100` или `12345 -50`:", parse_mode="HTML", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)

    # ---- Рассылка (с подписью) ----

    elif call.data == "admin_broadcast":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_broadcast"
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="admin_panel"))
        bot.send_message(cid, "Введите текст рассылки (будет добавлена подпись 'Сообщение от администратора'):", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)

    # ---- Блокировка ----

    elif call.data == "admin_block":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_block"
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="admin_panel"))
        bot.send_message(cid, "Введите ID для блокировки:", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)

    # ---- Разблокировка ----

    elif call.data == "admin_unblock":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_unblock"
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="admin_panel"))
        bot.send_message(cid, "Введите ID для разблокировки:", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)

    # ---- Завершение чата ----

    elif call.data == "admin_end_chat":
        if not is_admin(uid): return
        user_states[uid] = "admin_waiting_end_chat"
        cancel_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("❌ Отмена", callback_data="admin_panel"))
        bot.send_message(cid, "Введите ID пользователя для завершения чата:", reply_markup=cancel_markup)
        bot.answer_callback_query(call.id)

    # ---- Статистика ----

    elif call.data == "admin_stats":
        if not is_admin(uid): return
        total_users = cursor.execute("SELECT COUNT(*) FROM users").fetchone()[0]
        total_blocked = cursor.execute("SELECT COUNT(*) FROM users WHERE blocked=1").fetchone()[0]
        total_queries = cursor.execute("SELECT COUNT(*) FROM queries_log").fetchone()[0]
        total_balance = cursor.execute("SELECT SUM(balance) FROM users").fetchone()[0] or 0
        stats = f"📊 <b>Статистика</b>\n\n"
        stats += f"👥 Всего пользователей: {total_users}\n"
        stats += f"🚫 Заблокировано: {total_blocked}\n"
        stats += f"📝 Всего запросов: {total_queries}\n"
        stats += f"💰 Общий баланс пользователей: {total_balance/100:.2f} ₽"
        try:
            bot.edit_message_text(stats, cid, mid, parse_mode="HTML")
        except:
            bot.send_message(cid, stats, parse_mode="HTML")
        back_markup = types.InlineKeyboardMarkup().add(types.InlineKeyboardButton("◀️ Назад в админ-панель", callback_data="admin_panel"))
        bot.send_message(cid, "Выберите действие:", reply_markup=back_markup)
        bot.answer_callback_query(call.id)

    # ---- Ответы и завершение чата ----

    elif call.data.startswith("reply_to_"):
        if not is_admin(uid): return
        target = int(call.data.split("_")[2])
        user_states[uid] = f"replying_to_{target}"
        bot.send_message(uid, f"✏️ Теперь вы отвечаете пользователю {target}. Все следующие сообщения будут отправляться ему. Чтобы завершить, используйте кнопку «Завершить чат» в админ-панели.")
        bot.answer_callback_query(call.id)
    elif call.data.startswith("end_chat_"):
        if not is_admin(uid): return
        target = int(call.data.split("_")[2])
        if uid in user_states and isinstance(user_states[uid], str) and user_states[uid].startswith("replying_to_") and int(user_states[uid].split("_")[2]) == target:
            del user_states[uid]
        set_in_chat(target, 0)
        markup = rating_keyboard(target)
        try:
            bot.send_message(target, get_text(target, 'chat_ended'), reply_markup=markup, parse_mode="HTML")
            bot.send_message(uid, f"✅ Чат с {target} завершён, запрос оценки отправлен.")
        except Exception as e:
            bot.send_message(uid, f"❌ Не удалось отправить сообщение пользователю {target}: {e}")
        bot.answer_callback_query(call.id)
    elif call.data.startswith("rate_"):
        rating = int(call.data.split("_")[1])
        user_states[uid] = f"awaiting_review_comment:{rating}"
        stars = "⭐" * rating
        text = get_text(uid, 'rate_prompt', rating=rating, stars=stars)
        try:
            bot.edit_message_text(text, cid, mid, parse_mode="HTML")
        except:
            bot.send_message(cid, text, parse_mode="HTML")
        bot.answer_callback_query(call.id)

# ========== ФУНКЦИЯ ДЛЯ ОТОБРАЖЕНИЯ СТРАНИЦЫ ПОЛЬЗОВАТЕЛЕЙ ==========

def show_user_page(admin_uid, chat_id, message_id, page):
    users_per_page = 5
    with db_lock:
        total = cursor.execute("SELECT COUNT(*) FROM users").fetchone()[0]
        offset = page * users_per_page
        users = cursor.execute("""
            SELECT user_id, username, first_name, balance, blocked, reg_date, text_model, voice_enabled,
                   in_chat, last_admin_message, language, theme, mode, internet_search_enabled
            FROM users ORDER BY reg_date DESC LIMIT ? OFFSET ?
        """, (users_per_page, offset)).fetchall()
    total_pages = (total + users_per_page - 1) // users_per_page
    if page < 0: page = 0
    if page >= total_pages and total_pages > 0: page = total_pages - 1

    msg = f"📋 <b>Список пользователей</b> (страница {page+1} из {max(total_pages,1)}):\n\n"
    for user in users:
        uid, uname, fname, bal, blocked, reg, model, voice, in_chat, last_msg, lang, theme, mode, search = user
        status = "🚫" if blocked else "✅"
        voice_emoji = "🎤" if voice else ""
        chat_emoji = "💬" if in_chat else ""
        search_status = "🔍" if search else ""
        safe_uname = escape_html(uname) if uname else ""
        safe_fname = escape_html(fname) if fname else ""
        display_model = next((d for d,m in AVAILABLE_MODELS.items() if m==model), model)
        lang_name = get_language_name(lang)
        theme_name = get_theme_name(theme)
        mode_name = "Обычный" if mode=='normal' else "Без цензуры"
        msg += (f"{status} ID: <code>{uid}</code> @{safe_uname} {safe_fname} {chat_emoji} {search_status}\n"
                f"Баланс: {bal/100:.2f} Модель: {display_model} {voice_emoji}\n"
                f"Язык: {lang_name}, Тема: {theme_name}, Режим: {mode_name}\n"
                f"Рег: {reg[:10]}\n\n")

    markup = types.InlineKeyboardMarkup(row_width=3)
    buttons = []
    if page > 0:
        buttons.append(types.InlineKeyboardButton("◀️ Назад", callback_data=f"admin_list_page_{page-1}"))
    else:
        buttons.append(types.InlineKeyboardButton("•", callback_data="noop"))
    buttons.append(types.InlineKeyboardButton(f"{page+1}/{total_pages}", callback_data="noop"))
    if page < total_pages - 1:
        buttons.append(types.InlineKeyboardButton("Вперёд ▶️", callback_data=f"admin_list_page_{page+1}"))
    else:
        buttons.append(types.InlineKeyboardButton("•", callback_data="noop"))
    markup.add(*buttons)
    markup.add(types.InlineKeyboardButton("🔙 Назад в админку", callback_data="admin_panel"))

    try:
        bot.edit_message_text(msg, chat_id, message_id, parse_mode="HTML", reply_markup=markup)
    except:
        bot.send_message(chat_id, msg, parse_mode="HTML", reply_markup=markup)

# ========== ПЛАТЕЖИ ==========

@bot.pre_checkout_query_handler(func=lambda q: True)
def pre_checkout_query(q):
    bot.answer_pre_checkout_query(q.id, ok=True)

@bot.message_handler(content_types=['successful_payment'])
def successful_payment(message):
    payload = message.successful_payment.invoice_payload
    if payload.startswith("stars_"):
        parts = payload.split("_")
        if len(parts) == 3:
            try:
                uid = int(parts[1])
                rub = int(parts[2])
                update_balance(uid, rub*100, reason="Пополнение звёздами")
                bot.send_message(message.chat.id, f"✅ Баланс пополнен на {rub} ₽.")
            except:
                bot.send_message(message.chat.id, "❌ Ошибка начисления.")
    else:
        parts = payload.split('_')
        if len(parts) == 4 and parts[0] == 'topup':
            try:
                uid = int(parts[1])
                stars = int(parts[2])
                amount = int(parts[3])
                STAR_PACKAGES = {100:10000,250:25000,500:50000,1000:100000}
                if stars in STAR_PACKAGES and STAR_PACKAGES[stars] == amount:
                    update_balance(uid, amount, reason="Звёзды")
                    bot.send_message(message.chat.id, f"✅ Баланс пополнен на {amount/100:.2f} ₽.")
                else:
                    bot.send_message(message.chat.id, "❌ Ошибка начисления.")
            except:
                bot.send_message(message.chat.id, "❌ Ошибка.")
        else:
            bot.send_message(message.chat.id, "❌ Неизвестный платёж.")

# ========== ФУНКЦИИ УДАЛЕНИЯ МОДЕЛЕЙ ==========

def remove_text_model(model_id):
    if len(AVAILABLE_MODELS) <= 1:
        raise ValueError("Нельзя удалить последнюю модель")
    with db_lock:
        cursor.execute("DELETE FROM text_models WHERE model_id=?", (model_id,))
        conn.commit()
        name_to_remove = None
        for name, mid in list(AVAILABLE_MODELS.items()):
            if mid == model_id:
                name_to_remove = name
                break
        if name_to_remove:
            del AVAILABLE_MODELS[name_to_remove]
            del MODEL_PRICES[model_id]
        cursor.execute("UPDATE users SET text_model=? WHERE text_model=?", (DEFAULT_MODEL, model_id))
        conn.commit()

def remove_tts_model(model_id):
    if len(AVAILABLE_TTS_MODELS) <= 1:
        raise ValueError("Нельзя удалить последнюю TTS модель")
    with db_lock:
        cursor.execute("DELETE FROM tts_models WHERE model_id=?", (model_id,))
        conn.commit()
        del AVAILABLE_TTS_MODELS[model_id]
        if get_tts_model() == model_id:
            new_default = next(iter(AVAILABLE_TTS_MODELS))
            set_tts_model(new_default)

def remove_video_model(model_id):
    if len(AVAILABLE_VIDEO_MODELS) <= 1:
        raise ValueError("Нельзя удалить последнюю видео модель")
    with db_lock:
        cursor.execute("DELETE FROM video_models WHERE model_id=?", (model_id,))
        conn.commit()
        del AVAILABLE_VIDEO_MODELS[model_id]
        if get_video_model() == model_id:
            new_default = next(iter(AVAILABLE_VIDEO_MODELS))
            set_video_model(new_default)

def remove_image_model(model_id):
    if len(AVAILABLE_IMAGE_MODELS) <= 1:
        raise ValueError("Нельзя удалить последнюю модель изображений")
    with db_lock:
        cursor.execute("DELETE FROM image_models WHERE model_id=?", (model_id,))
        conn.commit()
        del AVAILABLE_IMAGE_MODELS[model_id]
        if get_image_model() == model_id:
            new_default = next(iter(AVAILABLE_IMAGE_MODELS))
            set_image_model(new_default)

# ========== ЗАПУСК ==========

if __name__ == "__main__":
    print(f"✅ Бот {BOT_NAME} (версия {BOT_VERSION}) запущен!")
    os.makedirs("sessions", exist_ok=True)
    while True:
        try:
            bot.infinity_polling(timeout=120, long_polling_timeout=180)
        except Exception as e:
            log_error(f"Polling error: {e}")
            time.sleep(5)
